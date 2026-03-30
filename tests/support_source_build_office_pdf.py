"""Office, spreadsheet, and PDF source-build tests."""

from __future__ import annotations

import io
import json
import os
import tempfile
import threading
import time
import unittest
import warnings
from pathlib import Path
from unittest import mock

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE

from docmason.commands import (
    ACTION_REQUIRED,
    DEGRADED,
    status_workspace,
    sync_workspace,
    validate_knowledge_base,
)
from docmason.coordination import lease_dir, workspace_lease
from docmason.control_plane import load_shared_job
from docmason.evidence_artifacts import compile_pptx_visual_artifacts
from docmason.hybrid import (
    required_overlay_slots,
    select_lane_b_batch,
)
from docmason.knowledge import (
    _ThirdPartyDiagnosticCapture,
    build_single_source_artifacts,
    build_docx_source,
    build_pptx_source,
    build_xlsx_source,
    locate_previous_source_dir,
    render_pdf_document,
    sanitize_text,
    source_artifact_contract_complete,
)
from docmason.project import WorkspacePaths, read_json, write_json
from docmason.semantic_overlays import semantic_overlay_candidates, write_semantic_overlay
from docmason.versioning import apply_snapshot_retention
from tests.support_ready_workspace import seed_self_contained_bootstrap_state


class SourceBuildOfficePdfTests(unittest.TestCase):
    """Cover office, spreadsheet, and PDF build behavior."""

    def make_workspace(self) -> WorkspacePaths:
        tempdir = tempfile.TemporaryDirectory()
        self.addCleanup(tempdir.cleanup)
        root = Path(tempdir.name)

        (root / "src" / "docmason").mkdir(parents=True)
        (root / "skills" / "canonical" / "workspace-bootstrap").mkdir(parents=True)
        (root / "original_doc").mkdir()
        (root / "knowledge_base").mkdir()
        (root / "runtime").mkdir()
        (root / "pyproject.toml").write_text(
            "[project]\nname = 'docmason'\nversion = '0.0.0'\n",
            encoding="utf-8",
        )
        (root / "docmason.yaml").write_text(
            "workspace:\n  source_dir: original_doc\n",
            encoding="utf-8",
        )
        (root / "src" / "docmason" / "__init__.py").write_text(
            "__version__ = '0.0.0'\n",
            encoding="utf-8",
        )
        (root / "AGENTS.md").write_text("# Agents\n", encoding="utf-8")
        (root / "skills" / "canonical" / "workspace-bootstrap" / "SKILL.md").write_text(
            "# Workspace Bootstrap\n",
            encoding="utf-8",
        )
        return WorkspacePaths(root=root)

    def ready_probe(self, _workspace: WorkspacePaths) -> tuple[bool, str]:
        return True, "Editable install resolves to the workspace source tree."

    def mark_environment_ready(self, workspace: WorkspacePaths) -> None:
        seed_self_contained_bootstrap_state(
            workspace,
            prepared_at="2026-03-16T00:00:00Z",
        )

    def create_pdf(self, path: Path) -> None:
        from pypdf import PdfWriter

        writer = PdfWriter()
        writer.add_blank_page(width=144, height=144)
        with path.open("wb") as handle:
            writer.write(handle)

    def create_pdf_with_chart(self, path: Path) -> None:
        try:
            import pymupdf  # type: ignore[import-not-found]
        except ImportError:  # pragma: no cover - compatibility import
            import fitz as pymupdf  # type: ignore[import-not-found]

        document = pymupdf.open()
        page = document.new_page(width=420, height=320)
        page.insert_text((48, 36), "Quarterly Revenue Chart")
        page.insert_text((48, 58), "Q1 Q2 Q3 Q4 Actual Budget")
        page.draw_line((48, 260), (360, 260), color=(0, 0, 0), width=1.2)
        page.draw_line((48, 84), (48, 260), color=(0, 0, 0), width=1.2)
        for index, height in enumerate((44, 86, 62, 118), start=0):
            left = 84 + (index * 56)
            page.draw_rect(
                (left, 260 - height, left + 28, 260),
                color=(0.1, 0.2, 0.8),
                fill=(0.1, 0.2, 0.8),
            )
            page.insert_text((left - 2, 278), f"Q{index + 1}")
        page.insert_text((300, 92), "Revenue")
        document.save(path)
        document.close()

    def create_pdf_with_full_page_image(self, path: Path) -> None:
        try:
            import pymupdf  # type: ignore[import-not-found]
        except ImportError:  # pragma: no cover - compatibility import
            import fitz as pymupdf  # type: ignore[import-not-found]
        from PIL import Image, ImageDraw

        with tempfile.TemporaryDirectory() as tempdir_name:
            image_path = Path(tempdir_name) / "page.png"
            image = Image.new("RGB", (1200, 1600), color=(245, 245, 245))
            draw = ImageDraw.Draw(image)
            draw.rectangle((80, 120, 1120, 1480), outline=(32, 64, 128), width=14)
            draw.rectangle((170, 300, 1030, 540), fill=(210, 225, 245))
            draw.rectangle((170, 660, 1030, 1180), fill=(225, 235, 250))
            draw.rectangle((170, 1230, 760, 1380), fill=(235, 240, 250))
            image.save(image_path)

            document = pymupdf.open()
            page = document.new_page(width=595, height=842)
            page.insert_image(page.rect, filename=str(image_path))
            document.save(path)
            document.close()

    def create_pdf_with_sections_and_tables(self, path: Path) -> None:
        try:
            import pymupdf  # type: ignore[import-not-found]
        except ImportError:  # pragma: no cover - compatibility import
            import fitz as pymupdf  # type: ignore[import-not-found]

        document = pymupdf.open()
        for page_index in range(2):
            page = document.new_page(width=520, height=720)
            page.insert_text((48, 48), "1. Executive Summary", fontsize=18)
            page.insert_text(
                (48, 78),
                "The document summarises KPI trends and table evidence for follow-up analysis.",
                fontsize=11,
            )
            page.insert_text((48, 118), "Table 1. KPI Summary", fontsize=12)
            top = 150
            left = 48
            width = 320
            row_height = 30
            col_width = 100
            for row in range(4):
                y = top + (row * row_height)
                page.draw_line((left, y), (left + width, y), color=(0, 0, 0), width=1)
            for col in range(4):
                x = left + (col * col_width)
                page.draw_line((x, top), (x, top + (3 * row_height)), color=(0, 0, 0), width=1)
            headers = ["Quarter", "Actual", "Budget"]
            values = [
                ["Q1", "10", "12"],
                [
                    "Q2" if page_index == 0 else "Q3",
                    "15" if page_index == 0 else "18",
                    "14" if page_index == 0 else "17",
                ],
            ]
            for col, header in enumerate(headers):
                page.insert_text((left + 8 + (col * col_width), top + 18), header, fontsize=10)
            for row, row_values in enumerate(values, start=1):
                for col, value in enumerate(row_values):
                    page.insert_text(
                        (left + 8 + (col * col_width), top + 18 + (row * row_height)),
                        value,
                        fontsize=10,
                    )
            if page_index == 1:
                page.insert_text(
                    (48, 320),
                    (
                        "1. Continue validating the metric deltas.\n"
                        "2. Escalate anomalies to the operating owner.\n"
                        "3. Review the supporting appendix."
                    ),
                    fontsize=11,
                )
        document.save(path)
        document.close()

    def create_pptx(self, path: Path) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Test Deck"
        slide.placeholders[1].text = "Hello from Office source build."
        presentation.save(path)

    def create_pptx_with_hidden_slide(self, path: Path) -> None:
        presentation = Presentation()
        for index in range(3):
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = f"Slide {index + 1}"
            slide.placeholders[1].text = f"Body {index + 1}"
            if index == 1:
                slide._element.set("show", "0")
        presentation.save(path)

    def create_pptx_with_connector_flow(self, path: Path) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        first = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 900000, 1200000, 1800000, 700000
        )
        first.text = "Collect"
        second = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 3600000, 1200000, 1800000, 700000
        )
        second.text = "Review"
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            2700000,
            1550000,
            3600000,
            1550000,
        )
        connector.name = "Flow Connector"
        label = slide.shapes.add_textbox(2100000, 700000, 1400000, 400000)
        label.text_frame.text = "Approval Flow"
        presentation.save(path)

    def create_docx(self, path: Path) -> None:
        from docx import Document

        document = Document()
        document.add_heading("Legacy Doc", level=1)
        document.add_paragraph("Legacy Word content.")
        document.save(path)

    def create_docx_with_structure(self, path: Path) -> None:
        from docx import Document

        document = Document()
        document.add_heading("Implementation Overview", level=1)
        document.add_paragraph("Figure 1. Process sketch")
        document.add_paragraph("The document explains the sequence in a structured way.")
        document.add_paragraph("1. Capture the source inputs.")
        document.add_paragraph("2. Review the semantic evidence.")
        document.add_paragraph("3. Publish the validated outputs.")
        document.add_heading("Supporting Notes", level=2)
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "Coverage"
        table.cell(1, 1).text = "High"
        document.save(path)

    def create_docx_with_embedded_picture(self, path: Path) -> None:
        from docx import Document
        from PIL import Image

        with tempfile.TemporaryDirectory() as tempdir_name:
            image_path = Path(tempdir_name) / "figure.png"
            Image.new("RGB", (120, 80), color=(32, 96, 160)).save(image_path, format="PNG")
            document = Document()
            document.add_heading("Architecture Sketch", level=1)
            document.add_paragraph("Figure 1. Logical view")
            document.add_picture(str(image_path))
            document.add_paragraph("The figure captures the main logical components.")
            document.save(path)

    def create_xlsx(self, path: Path) -> None:
        from openpyxl import Workbook

        workbook = Workbook()
        worksheet = workbook.active
        worksheet.title = "Sheet 1"
        worksheet["A1"] = "Metric"
        worksheet["B1"] = "Value"
        worksheet["A2"] = "Budget"
        worksheet["B2"] = 42
        workbook.save(path)

    def create_xlsx_with_table(self, path: Path) -> None:
        from openpyxl import Workbook
        from openpyxl.worksheet.table import Table, TableStyleInfo

        workbook = Workbook()
        worksheet = workbook.active
        worksheet.title = "Sheet 1"
        worksheet["A1"] = "Metric"
        worksheet["B1"] = "Value"
        worksheet["A2"] = "Budget"
        worksheet["B2"] = 42
        table = Table(displayName="MetricsTable", ref="A1:B2")
        table.tableStyleInfo = TableStyleInfo(
            name="TableStyleMedium2",
            showFirstColumn=False,
            showLastColumn=False,
            showRowStripes=True,
            showColumnStripes=False,
        )
        worksheet.add_table(table)
        workbook.save(path)

    def create_xlsx_with_chart(self, path: Path) -> None:
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, Reference
        from openpyxl.worksheet.table import Table, TableStyleInfo

        workbook = Workbook()
        worksheet = workbook.active
        worksheet.title = "Dashboard"
        for row in [
            ("Quarter", "Actual", "Budget"),
            ("Q1", 10, 12),
            ("Q2", 15, 14),
            ("Q3", 13, 16),
            ("Q4", 18, 17),
        ]:
            worksheet.append(row)
        table = Table(displayName="MetricsTable", ref="A1:C5")
        table.tableStyleInfo = TableStyleInfo(
            name="TableStyleMedium2",
            showFirstColumn=False,
            showLastColumn=False,
            showRowStripes=True,
            showColumnStripes=False,
        )
        worksheet.add_table(table)
        chart = BarChart()
        chart.title = "Quarterly Revenue"
        chart.x_axis.title = "Quarter"
        chart.y_axis.title = "Revenue"
        data = Reference(worksheet, min_col=2, max_col=3, min_row=1, max_row=5)
        cats = Reference(worksheet, min_col=1, min_row=2, max_row=5)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        worksheet.add_chart(chart, "E2")
        workbook.save(path)

    def create_xlsx_with_embedded_image(self, path: Path) -> None:
        from openpyxl import Workbook
        from openpyxl.drawing.image import Image as OpenpyxlImage
        from PIL import Image

        with tempfile.TemporaryDirectory() as tempdir_name:
            image_path = Path(tempdir_name) / "logo.png"
            Image.new("RGB", (80, 40), color=(32, 96, 160)).save(image_path, format="PNG")
            workbook = Workbook()
            worksheet = workbook.active
            worksheet.title = "Quote"
            worksheet["A1"] = "Item"
            worksheet["B1"] = "Cost"
            worksheet["A2"] = "Design"
            worksheet["B2"] = 120
            worksheet.add_image(OpenpyxlImage(str(image_path)), "D2")
            workbook.save(path)

    def seed_agent_outputs(self, source_dir: Path) -> None:
        source_manifest = read_json(source_dir / "source_manifest.json")
        evidence_manifest = read_json(source_dir / "evidence_manifest.json")
        first_unit_id = evidence_manifest["units"][0]["unit_id"]
        knowledge = {
            "source_id": source_manifest["source_id"],
            "source_fingerprint": source_manifest["source_fingerprint"],
            "title": "Seeded knowledge object",
            "source_language": "en",
            "summary_en": "A short English summary for the staged source.",
            "summary_source": "A short English summary for the staged source.",
            "document_type": source_manifest["document_type"],
            "key_points": [
                {
                    "text_en": "The document contains one staged evidence unit.",
                    "text_source": "The document contains one staged evidence unit.",
                    "citations": [{"unit_id": first_unit_id, "support": "primary"}],
                }
            ],
            "entities": [],
            "claims": [
                {
                    "statement_en": "The sync pipeline produced evidence for the source.",
                    "statement_source": "The sync pipeline produced evidence for the source.",
                    "citations": [{"unit_id": first_unit_id, "support": "primary"}],
                }
            ],
            "known_gaps": [],
            "ambiguities": [],
            "confidence": {
                "level": "medium",
                "notes_en": "Based on a minimal test fixture.",
                "notes_source": "Based on a minimal test fixture.",
            },
            "citations": [{"unit_id": first_unit_id, "support": "primary"}],
            "related_sources": [],
        }
        write_json(source_dir / "knowledge.json", knowledge)
        summary = "\n".join(
            [
                "# Seeded knowledge object",
                "",
                f"Source ID: {source_manifest['source_id']}",
                "",
                "## English Summary",
                "A short English summary for the staged source.",
                "",
                "## Source-Language Summary",
                "A short English summary for the staged source.",
                "",
            ]
        )
        (source_dir / "summary.md").write_text(summary, encoding="utf-8")

    def seed_snapshot_version(
        self,
        workspace: WorkspacePaths,
        *,
        snapshot_id: str,
        published_at: str,
        source_signature: str,
    ) -> Path:
        snapshot_dir = workspace.knowledge_version_dir(snapshot_id)
        snapshot_dir.mkdir(parents=True, exist_ok=True)
        write_json(
            snapshot_dir / "publish_manifest.json",
            {
                "snapshot_id": snapshot_id,
                "published_at": published_at,
                "validation_status": "valid",
                "published_source_signature": source_signature,
            },
        )
        write_json(
            snapshot_dir / "validation_report.json",
            {
                "status": "valid",
                "source_signature": source_signature,
            },
        )
        return snapshot_dir

    def test_sync_pdf_only_creates_staging_and_pending_synthesis(self) -> None:
        workspace = self.make_workspace()
        self.create_pdf_with_chart(workspace.source_dir / "example.pdf")

        report = sync_workspace(workspace, autonomous=False)

        self.assertEqual(report.exit_code, 2)
        self.assertEqual(report.payload["status"], DEGRADED)
        self.assertEqual(report.payload["sync_status"], "pending-synthesis")
        self.assertTrue(workspace.staging_catalog_path.exists())
        self.assertTrue(workspace.source_index_path.exists())
        self.assertEqual(len(report.payload["pending_sources"]), 1)
        evidence_manifest = read_json(
            workspace.knowledge_base_staging_dir
            / "sources"
            / report.payload["pending_sources"][0]["source_id"]
            / "evidence_manifest.json"
        )
        affordances = read_json(
            workspace.knowledge_base_staging_dir
            / "sources"
            / report.payload["pending_sources"][0]["source_id"]
            / "derived_affordances.json"
        )
        self.assertEqual(evidence_manifest["document_type"], "pdf")
        self.assertTrue(evidence_manifest["units"])
        self.assertEqual(evidence_manifest["artifact_index_asset"], "artifact_index.json")
        self.assertTrue(evidence_manifest["visual_layout_assets"])
        self.assertEqual(affordances["artifact_type"], "derived-affordances")
        self.assertTrue(affordances["source_affordances"]["available_channels"])
        artifact_index = read_json(
            workspace.knowledge_base_staging_dir
            / "sources"
            / report.payload["pending_sources"][0]["source_id"]
            / "artifact_index.json"
        )
        self.assertTrue(artifact_index["artifacts"])
        self.assertTrue(
            any(
                item["artifact_type"] in {"chart", "major-region", "text-region"}
                for item in artifact_index["artifacts"]
                if isinstance(item, dict)
            )
        )

    def test_sync_pdf_builds_pdf_document_and_hybrid_candidate_summary(self) -> None:
        workspace = self.make_workspace()
        self.create_pdf_with_sections_and_tables(workspace.source_dir / "example.pdf")

        report = sync_workspace(workspace, autonomous=False)

        self.assertEqual(report.payload["sync_status"], "pending-synthesis")
        pending_source = report.payload["pending_sources"][0]["source_id"]
        source_dir = workspace.knowledge_base_staging_dir / "sources" / pending_source
        evidence_manifest = read_json(source_dir / "evidence_manifest.json")
        self.assertEqual(evidence_manifest["pdf_document_asset"], "pdf_document.json")
        pdf_document = read_json(source_dir / "pdf_document.json")
        self.assertTrue(pdf_document["page_contexts"])
        self.assertTrue(pdf_document["caption_links"])
        self.assertTrue(pdf_document["continuation_links"])
        self.assertTrue(pdf_document["procedure_spans"])
        hybrid = report.payload["hybrid_enrichment"]
        self.assertGreaterEqual(hybrid["eligible_unit_count"], 1)
        self.assertEqual(hybrid["mode"], "candidate-prepared")

    def test_sync_builds_focus_render_assets_and_extended_hybrid_contract(self) -> None:
        workspace = self.make_workspace()
        self.create_pdf_with_chart(workspace.source_dir / "example.pdf")

        report = sync_workspace(workspace, autonomous=False)

        self.assertEqual(report.payload["sync_status"], "pending-synthesis")
        pending_source = report.payload["pending_sources"][0]["source_id"]
        source_dir = workspace.knowledge_base_staging_dir / "sources" / pending_source
        artifact_index = read_json(source_dir / "artifact_index.json")
        focus_artifacts = [
            artifact
            for artifact in artifact_index["artifacts"]
            if artifact.get("artifact_type") in {"chart", "major-region", "page-image"}
        ]
        self.assertTrue(focus_artifacts)
        self.assertTrue(any(artifact.get("focus_render_assets") for artifact in focus_artifacts))
        for artifact in focus_artifacts:
            for asset in artifact.get("focus_render_assets", []):
                self.assertTrue((source_dir / asset).exists())

        layout_path = source_dir / "visual_layout" / "page-001.json"
        visual_layout = read_json(layout_path)
        chart_region = next(
            region
            for region in visual_layout["regions"]
            if region.get("artifact_id") == focus_artifacts[0]["artifact_id"]
        )
        self.assertEqual(
            chart_region["focus_render_assets"],
            focus_artifacts[0]["focus_render_assets"],
        )

        hybrid_work = read_json(workspace.knowledge_base_staging_dir / "hybrid_work.json")
        source_packet = hybrid_work["sources"][0]
        self.assertTrue(source_packet["source_fingerprint"])
        self.assertEqual(source_packet["source_hybrid_status"], "candidate-prepared")
        self.assertEqual(
            source_packet["remaining_candidate_count"],
            source_packet["candidate_unit_count"],
        )
        candidate = source_packet["units"][0]
        self.assertTrue(candidate["unit_title"])
        self.assertTrue(candidate["unit_evidence_fingerprint"])
        self.assertTrue(candidate["required_overlay_slots"])
        self.assertTrue(candidate["suggested_overlay_kinds"])
        self.assertTrue(candidate["target_focus_render_assets"])
        self.assertEqual(candidate["coverage_status"], "candidate-prepared")
        lane_b_batch = select_lane_b_batch(hybrid_work)
        self.assertLessEqual(len(lane_b_batch), 4)
        self.assertLessEqual(sum(len(item["units"]) for item in lane_b_batch), 12)

    def test_third_party_diagnostic_capture_swallows_stderr_and_filters_benign_messages(
        self,
    ) -> None:
        outer_stderr = io.StringIO()
        with mock.patch("sys.stderr", outer_stderr):
            capture = _ThirdPartyDiagnosticCapture()
            with capture:
                warnings.warn(
                    "Conditional Formatting extension is not supported and will be removed",
                    UserWarning,
                    stacklevel=2,
                )
                warnings.warn("Meaningful parser warning", UserWarning, stacklevel=2)
                os.write(2, b"Ignoring wrong pointing object 12 0 R\n")
                os.write(2, b"low-level parser failure\n")

        self.assertEqual(outer_stderr.getvalue(), "")
        self.assertCountEqual(
            capture.messages,
            ["Meaningful parser warning", "low-level parser failure"],
        )

    def test_render_pdf_document_uses_fast_png_settings(self) -> None:
        renders_dir = Path(tempfile.mkdtemp())
        self.addCleanup(lambda: renders_dir.exists() and renders_dir.rmdir())

        document = mock.Mock()
        document.__len__ = mock.Mock(return_value=1)
        page = mock.Mock()
        bitmap = mock.Mock()
        image = mock.Mock()
        document.__getitem__ = mock.Mock(return_value=page)
        page.render.return_value = bitmap
        bitmap.to_pil.return_value = image
        pdfium = mock.Mock()
        pdfium.PdfDocument.return_value = document

        with mock.patch("docmason.knowledge.import_pdf_modules", return_value=(pdfium, object())):
            rendered_assets, failures = render_pdf_document(Path("example.pdf"), renders_dir)

        self.assertEqual(rendered_assets, ["renders/page-001.png"])
        self.assertEqual(failures, [])
        image.save.assert_called_once_with(
            renders_dir / "page-001.png",
            format="PNG",
            compress_level=1,
            optimize=False,
        )

    def test_sync_requires_office_renderer_for_office_sources(self) -> None:
        workspace = self.make_workspace()
        self.create_pptx(workspace.source_dir / "deck.pptx")

        with mock.patch(
            "docmason.knowledge.validate_soffice_binary",
            return_value={
                "ready": False,
                "binary": None,
                "version": None,
                "detail": "No LibreOffice command candidate was detected.",
            },
        ):
            report = sync_workspace(workspace)

        self.assertEqual(report.exit_code, 1)
        self.assertEqual(report.payload["status"], ACTION_REQUIRED)
        self.assertEqual(report.payload["sync_status"], "action-required")

    def test_build_single_source_artifacts_routes_legacy_office_by_document_type(self) -> None:
        workspace = self.make_workspace()
        source_dir = workspace.knowledge_base_staging_dir / "sources" / "source-legacy"
        source_dir.mkdir(parents=True, exist_ok=True)
        cases = [
            ("legacy-deck.ppt", "pptx", "build_pptx_source"),
            ("legacy-doc.doc", "docx", "build_docx_source"),
            ("legacy-sheet.xls", "xlsx", "build_xlsx_source"),
        ]

        for filename, document_type, builder_name in cases:
            with self.subTest(filename=filename, document_type=document_type):
                source_path = workspace.source_dir / filename
                source_path.write_bytes(b"legacy")
                source_entry = {
                    "source_id": f"source-{document_type}",
                    "source_fingerprint": f"fingerprint-{document_type}",
                    "current_path": f"original_doc/{filename}",
                    "prior_paths": [],
                    "document_type": document_type,
                    "source_extension": source_path.suffix.lstrip("."),
                    "first_seen_at": "2026-03-26T00:00:00Z",
                    "last_seen_at": "2026-03-26T00:00:00Z",
                    "identity_confidence": "new",
                }
                sentinel_manifest = {"document_type": document_type}
                sentinel_evidence = {"document_type": document_type, "units": []}
                builder_path = f"docmason.knowledge.{builder_name}"

                with mock.patch(
                    builder_path,
                    return_value=(sentinel_manifest, sentinel_evidence),
                ) as builder:
                    manifest, evidence = build_single_source_artifacts(
                        workspace,
                        source_entry,
                        source_dir,
                        office_binary="soffice",
                    )

                builder.assert_called_once()
                self.assertEqual(manifest, sentinel_manifest)
                self.assertEqual(evidence, sentinel_evidence)

    def test_hidden_pptx_slides_do_not_consume_visible_render_slots(self) -> None:
        workspace = self.make_workspace()
        source_path = workspace.source_dir / "hidden-deck.pptx"
        source_dir = workspace.knowledge_base_staging_dir / "sources" / "source-1"
        source_dir.mkdir(parents=True, exist_ok=True)
        self.create_pptx_with_hidden_slide(source_path)

        with (
            mock.patch(
                "docmason.knowledge.convert_office_to_pdf",
                return_value=(Path("dummy.pdf"), []),
            ),
            mock.patch(
                "docmason.knowledge.render_pdf_document",
                return_value=(["renders/page-001.png", "renders/page-002.png"], []),
            ),
        ):
            _source_manifest, evidence_manifest = build_pptx_source(
                workspace,
                source_path,
                {
                    "source_id": "source-1",
                    "source_fingerprint": "fingerprint-1",
                    "prior_paths": [],
                    "document_type": "pptx",
                    "first_seen_at": "2026-03-16T00:00:00Z",
                    "last_seen_at": "2026-03-16T00:00:00Z",
                    "identity_confidence": "new",
                },
                source_dir,
                soffice_binary="soffice",
            )

        self.assertEqual(evidence_manifest["units"][0]["rendered_asset"], "renders/page-001.png")
        self.assertTrue(evidence_manifest["units"][1]["hidden"])
        self.assertIsNone(evidence_manifest["units"][1]["rendered_asset"])
        self.assertEqual(evidence_manifest["units"][2]["rendered_asset"], "renders/page-002.png")
        artifact_index = read_json(source_dir / "artifact_index.json")
        first_slide_artifact = next(
            item for item in artifact_index["artifacts"] if item.get("unit_id") == "slide-001"
        )
        self.assertEqual(first_slide_artifact["render_page_span"], {"start": 1, "end": 1})

    def test_legacy_ppt_reuses_pptx_builder_via_libreoffice_normalization(self) -> None:
        workspace = self.make_workspace()
        source_path = workspace.source_dir / "legacy-deck.ppt"
        normalized_path = workspace.source_dir / "normalized-deck.pptx"
        source_dir = workspace.knowledge_base_staging_dir / "sources" / "source-ppt"
        source_dir.mkdir(parents=True, exist_ok=True)
        self.create_pptx(normalized_path)
        source_path.write_bytes(normalized_path.read_bytes())

        with (
            mock.patch(
                "docmason.knowledge.normalize_legacy_office_source",
                return_value=(normalized_path, []),
            ),
            mock.patch(
                "docmason.knowledge.convert_office_to_pdf",
                return_value=(Path("dummy.pdf"), []),
            ),
            mock.patch(
                "docmason.knowledge.render_pdf_document",
                return_value=(["renders/page-001.png"], []),
            ),
        ):
            _source_manifest, evidence_manifest = build_pptx_source(
                workspace,
                source_path,
                {
                    "source_id": "source-ppt",
                    "source_fingerprint": "fingerprint-ppt",
                    "prior_paths": [],
                    "document_type": "pptx",
                    "source_extension": "ppt",
                    "first_seen_at": "2026-03-19T00:00:00Z",
                    "last_seen_at": "2026-03-19T00:00:00Z",
                    "identity_confidence": "new",
                },
                source_dir,
                soffice_binary="soffice",
            )

        self.assertEqual(evidence_manifest["document_type"], "pptx")
        self.assertTrue(evidence_manifest["units"])

    def test_legacy_doc_reuses_docx_builder_via_libreoffice_normalization(self) -> None:
        workspace = self.make_workspace()
        source_path = workspace.source_dir / "legacy-doc.doc"
        normalized_path = workspace.source_dir / "normalized-doc.docx"
        source_dir = workspace.knowledge_base_staging_dir / "sources" / "source-doc"
        source_dir.mkdir(parents=True, exist_ok=True)
        self.create_docx(normalized_path)
        source_path.write_bytes(normalized_path.read_bytes())

        with (
            mock.patch(
                "docmason.knowledge.normalize_legacy_office_source",
                return_value=(normalized_path, []),
            ),
            mock.patch(
                "docmason.knowledge.convert_office_to_pdf",
                return_value=(Path("dummy.pdf"), []),
            ),
            mock.patch(
                "docmason.knowledge.render_pdf_document",
                return_value=(["renders/page-001.png"], []),
            ),
        ):
            _source_manifest, evidence_manifest = build_docx_source(
                workspace,
                source_path,
                {
                    "source_id": "source-doc",
                    "source_fingerprint": "fingerprint-doc",
                    "prior_paths": [],
                    "document_type": "docx",
                    "source_extension": "doc",
                    "first_seen_at": "2026-03-19T00:00:00Z",
                    "last_seen_at": "2026-03-19T00:00:00Z",
                    "identity_confidence": "new",
                },
                source_dir,
                soffice_binary="soffice",
            )

        self.assertEqual(evidence_manifest["document_type"], "docx")
        self.assertTrue(evidence_manifest["units"])

    def test_legacy_xls_reuses_xlsx_builder_via_libreoffice_normalization(self) -> None:
        workspace = self.make_workspace()
        source_path = workspace.source_dir / "legacy-sheet.xls"
        normalized_path = workspace.source_dir / "normalized-sheet.xlsx"
        source_dir = workspace.knowledge_base_staging_dir / "sources" / "source-xls"
        source_dir.mkdir(parents=True, exist_ok=True)
        self.create_xlsx(normalized_path)
        source_path.write_bytes(normalized_path.read_bytes())

        with (
            mock.patch(
                "docmason.knowledge.normalize_legacy_office_source",
                return_value=(normalized_path, []),
            ),
            mock.patch(
                "docmason.knowledge.convert_office_to_pdf",
                return_value=(Path("dummy.pdf"), []),
            ),
            mock.patch(
                "docmason.knowledge.render_pdf_document",
                return_value=(["renders/page-001.png"], []),
            ),
        ):
            _source_manifest, evidence_manifest = build_xlsx_source(
                workspace,
                source_path,
                {
                    "source_id": "source-xls",
                    "source_fingerprint": "fingerprint-xls",
                    "prior_paths": [],
                    "document_type": "xlsx",
                    "source_extension": "xls",
                    "first_seen_at": "2026-03-19T00:00:00Z",
                    "last_seen_at": "2026-03-19T00:00:00Z",
                    "identity_confidence": "new",
                },
                source_dir,
                soffice_binary="soffice",
            )

        self.assertEqual(evidence_manifest["document_type"], "xlsx")
        self.assertTrue(evidence_manifest["units"])

    def test_build_xlsx_source_reads_table_refs_across_openpyxl_variants(self) -> None:
        workspace = self.make_workspace()
        source_path = workspace.source_dir / "table-sheet.xlsx"
        source_dir = workspace.knowledge_base_staging_dir / "sources" / "source-table"
        source_dir.mkdir(parents=True, exist_ok=True)
        self.create_xlsx_with_table(source_path)

        with (
            mock.patch(
                "docmason.knowledge.convert_office_to_pdf",
                return_value=(Path("dummy.pdf"), []),
            ),
            mock.patch(
                "docmason.knowledge.render_pdf_document",
                return_value=(["renders/page-001.png"], []),
            ),
        ):
            _source_manifest, evidence_manifest = build_xlsx_source(
                workspace,
                source_path,
                {
                    "source_id": "source-table",
                    "source_fingerprint": "fingerprint-table",
                    "prior_paths": [],
                    "document_type": "xlsx",
                    "first_seen_at": "2026-03-20T00:00:00Z",
                    "last_seen_at": "2026-03-20T00:00:00Z",
                    "identity_confidence": "new",
                },
                source_dir,
                soffice_binary="soffice",
            )

        structure_path = source_dir / evidence_manifest["units"][0]["structure_asset"]
        structure = read_json(structure_path)
        self.assertEqual(structure["tables"], [{"name": "MetricsTable", "ref": "A1:B2"}])

    def test_build_xlsx_source_recovers_richer_chart_registry(self) -> None:
        workspace = self.make_workspace()
        source_path = workspace.source_dir / "chart-sheet.xlsx"
        source_dir = workspace.knowledge_base_staging_dir / "sources" / "source-chart-sheet"
        source_dir.mkdir(parents=True, exist_ok=True)
        self.create_xlsx_with_chart(source_path)

        with (
            mock.patch(
                "docmason.knowledge.convert_office_to_pdf",
                return_value=(Path("dummy.pdf"), []),
            ),
            mock.patch(
                "docmason.knowledge.render_xlsx_sheet_documents",
                return_value=(
                    {"Dashboard": ["renders/sheet-001-page-001.png"]},
                    ["renders/sheet-001-page-001.png"],
                    [],
                ),
            ),
        ):
            _source_manifest, evidence_manifest = build_xlsx_source(
                workspace,
                source_path,
                {
                    "source_id": "source-chart-sheet",
                    "source_fingerprint": "fingerprint-chart-sheet",
                    "prior_paths": [],
                    "document_type": "xlsx",
                    "first_seen_at": "2026-03-22T00:00:00Z",
                    "last_seen_at": "2026-03-22T00:00:00Z",
                    "identity_confidence": "new",
                },
                source_dir,
                soffice_binary="soffice",
            )

        workbook_payload = read_json(source_dir / evidence_manifest["spreadsheet_workbook_asset"])
        self.assertTrue(workbook_payload["chart_registry"])
        chart = workbook_payload["chart_registry"][0]
        self.assertEqual(chart["title"], "Quarterly Revenue")
        self.assertEqual(chart["x_axis_title"], "Quarter")
        self.assertEqual(chart["y_axis_title"], "Revenue")
        self.assertEqual(chart["linked_table_names"], ["MetricsTable"])
        sheet_payload = read_json(source_dir / evidence_manifest["spreadsheet_sheet_assets"][0])
        self.assertIn("dashboard-like", sheet_payload["sheet_role_hints"])

    def test_xlsx_picture_only_sheet_does_not_force_dashboard_chart_intent(self) -> None:
        workspace = self.make_workspace()
        source_path = workspace.source_dir / "image-sheet.xlsx"
        source_dir = workspace.knowledge_base_staging_dir / "sources" / "source-image-sheet"
        source_dir.mkdir(parents=True, exist_ok=True)
        self.create_xlsx_with_embedded_image(source_path)

        with (
            mock.patch(
                "docmason.knowledge.convert_office_to_pdf",
                return_value=(Path("dummy.pdf"), []),
            ),
            mock.patch(
                "docmason.knowledge.render_xlsx_sheet_documents",
                return_value=(
                    {"Quote": ["renders/sheet-001-page-001.png"]},
                    ["renders/sheet-001-page-001.png"],
                    [],
                ),
            ),
        ):
            _source_manifest, evidence_manifest = build_xlsx_source(
                workspace,
                source_path,
                {
                    "source_id": "source-image-sheet",
                    "source_fingerprint": "fingerprint-image-sheet",
                    "prior_paths": [],
                    "document_type": "xlsx",
                    "first_seen_at": "2026-03-22T00:00:00Z",
                    "last_seen_at": "2026-03-22T00:00:00Z",
                    "identity_confidence": "new",
                },
                source_dir,
                soffice_binary="soffice",
            )

        sheet_payload = read_json(source_dir / evidence_manifest["spreadsheet_sheet_assets"][0])
        self.assertEqual(sheet_payload["sheet_role_hints"], ["image-heavy"])
        artifact_index = read_json(source_dir / "artifact_index.json")
        picture_artifact = next(
            artifact
            for artifact in artifact_index["artifacts"]
            if artifact.get("artifact_type") == "picture"
        )
        self.assertTrue(picture_artifact["focus_render_assets"])
        self.assertTrue(picture_artifact["focus_render_assets"][0].startswith("media/"))
        self.assertTrue((source_dir / picture_artifact["focus_render_assets"][0]).exists())
        self.assertEqual(picture_artifact["image_ref"], "image-001")
        visual_layout = read_json(source_dir / evidence_manifest["visual_layout_assets"][0])
        picture_region = next(
            region
            for region in visual_layout["regions"]
            if region.get("artifact_type") == "picture"
        )
        self.assertEqual(
            picture_region["focus_render_assets"][0],
            picture_artifact["focus_render_assets"][0],
        )
        self.assertEqual(picture_region["image_ref"], "image-001")
        candidates = semantic_overlay_candidates(
            source_dir,
            evidence_manifest=evidence_manifest,
        )
        self.assertEqual(len(candidates), 1)
        self.assertNotIn("dashboard-like-sheet", candidates[0]["all_reasons"])
        self.assertNotIn("chart-intent", required_overlay_slots(candidates[0]))

    def test_xlsx_picture_focus_contract_requires_image_ref_for_reuse(self) -> None:
        from PIL import Image

        workspace = self.make_workspace()
        source_dir = workspace.knowledge_base_staging_dir / "sources" / "source-image-contract"
        (source_dir / "visual_layout").mkdir(parents=True, exist_ok=True)
        (source_dir / "spreadsheet_sheet").mkdir(parents=True, exist_ok=True)
        (source_dir / "media").mkdir(parents=True, exist_ok=True)
        (source_dir / "renders").mkdir(parents=True, exist_ok=True)

        media_asset = source_dir / "media" / "sheet-001-image-001.png"
        render_asset = source_dir / "renders" / "sheet-001-page-001.png"
        Image.new("RGB", (32, 32), color=(32, 96, 160)).save(media_asset, format="PNG")
        Image.new("RGB", (64, 64), color=(255, 255, 255)).save(render_asset, format="PNG")

        write_json(
            source_dir / "source_manifest.json",
            {
                "source_id": "source-image-contract",
                "document_type": "xlsx",
                "current_path": "original_doc/image-sheet.xlsx",
                "source_fingerprint": "fingerprint-image-contract",
            },
        )
        write_json(
            source_dir / "evidence_manifest.json",
            {
                "source_id": "source-image-contract",
                "document_type": "xlsx",
                "spreadsheet_workbook_asset": "spreadsheet_workbook.json",
                "spreadsheet_sheet_assets": ["spreadsheet_sheet/sheet-001.json"],
                "visual_layout_assets": ["visual_layout/sheet-001.json"],
                "units": [
                    {
                        "unit_id": "sheet-001",
                        "unit_type": "sheet",
                        "render_assets": ["renders/sheet-001-page-001.png"],
                        "render_page_span": {"start": 1, "end": 1},
                    }
                ],
            },
        )
        write_json(
            source_dir / "spreadsheet_workbook.json",
            {"artifact_type": "spreadsheet-workbook"},
        )
        write_json(
            source_dir / "spreadsheet_sheet" / "sheet-001.json",
            {"artifact_type": "spreadsheet-sheet", "unit_id": "sheet-001"},
        )
        write_json(
            source_dir / "artifact_index.json",
            {
                "artifacts": [
                    {
                        "artifact_id": "sheet-001:picture-001",
                        "artifact_type": "picture",
                        "unit_id": "sheet-001",
                        "title": "Quote logo",
                        "focus_render_assets": [
                            "media/sheet-001-image-001.png",
                            "renders/sheet-001-page-001.png",
                        ],
                        "render_assets": ["renders/sheet-001-page-001.png"],
                        "render_page_span": {"start": 1, "end": 1},
                        "image_ref": "image-001",
                    }
                ]
            },
        )
        write_json(
            source_dir / "visual_layout" / "sheet-001.json",
            {
                "artifact_type": "visual-layout",
                "unit_id": "sheet-001",
                "focus_render_assets": [
                    "media/sheet-001-image-001.png",
                    "renders/sheet-001-page-001.png",
                ],
                "regions": [
                    {
                        "artifact_id": "sheet-001:picture-001",
                        "artifact_type": "picture",
                        "focus_render_assets": [
                            "media/sheet-001-image-001.png",
                            "renders/sheet-001-page-001.png",
                        ],
                        "image_ref": "image-001",
                    }
                ],
            },
        )

        self.assertTrue(source_artifact_contract_complete(source_dir, document_type="xlsx"))

        artifact_index = read_json(source_dir / "artifact_index.json")
        artifact_index["artifacts"][0].pop("image_ref")
        write_json(source_dir / "artifact_index.json", artifact_index)

        self.assertFalse(source_artifact_contract_complete(source_dir, document_type="xlsx"))

    def test_pptx_visual_compiler_tolerates_unsupported_chart_properties(self) -> None:
        workspace = self.make_workspace()
        source_dir = workspace.knowledge_base_staging_dir / "sources" / "source-chart"
        source_dir.mkdir(parents=True, exist_ok=True)

        class ExplodingSeries:
            def __iter__(self):
                raise ValueError("unsupported series access")

        class PlotArea:
            def iterchildren(self):
                class Child:
                    tag = "{http://schemas.openxmlformats.org/drawingml/2006/chart}ofPieChart"

                return [Child()]

        class ChartSpace:
            plotArea = PlotArea()

        class ExplodingChart:
            has_title = False
            chart_title = None
            _chartSpace = ChartSpace()

            @property
            def chart_type(self):
                raise ValueError("unsupported plot type")

            @property
            def series(self):
                return ExplodingSeries()

        class Shape:
            has_chart = True
            has_table = False
            has_text_frame = False
            left = 10
            top = 20
            width = 200
            height = 120
            name = "Unsupported chart"
            chart = ExplodingChart()

        class Slide:
            shapes = [Shape()]

        class Presentation:
            slide_width = 720
            slide_height = 540
            slides = [Slide()]

        result = compile_pptx_visual_artifacts(
            source_dir,
            source_id="source-chart",
            presentation=Presentation(),
            units=[
                {
                    "unit_id": "slide-001",
                    "ordinal": 1,
                    "title": "Slide 1",
                    "rendered_asset": "renders/page-001.png",
                    "render_ordinal": 1,
                }
            ],
        )

        self.assertTrue(result["artifact_index"]["artifacts"])
        artifact = result["artifact_index"]["artifacts"][0]
        self.assertEqual(artifact["artifact_type"], "chart")
        self.assertEqual(artifact["visual_hints"], ["ofPieChart"])

    def test_pptx_visual_compiler_links_connectors_and_labels(self) -> None:
        workspace = self.make_workspace()
        source_path = workspace.source_dir / "connector-deck.pptx"
        self.create_pptx_with_connector_flow(source_path)
        source_dir = workspace.knowledge_base_staging_dir / "sources" / "source-connector"
        source_dir.mkdir(parents=True, exist_ok=True)
        presentation = Presentation(str(source_path))
        result = compile_pptx_visual_artifacts(
            source_dir,
            source_id="source-connector",
            presentation=presentation,
            units=[
                {
                    "unit_id": "slide-001",
                    "ordinal": 1,
                    "title": "Approval Flow",
                    "rendered_asset": "renders/page-001.png",
                    "render_ordinal": 1,
                }
            ],
        )

        connector = next(
            item
            for item in result["artifact_index"]["artifacts"]
            if item["artifact_type"] == "connector"
        )
        self.assertTrue(connector["linked_shape_artifact_ids"])
        self.assertTrue(
            any(
                isinstance(item.get("caption_text"), str) and item.get("caption_text")
                for item in result["artifact_index"]["artifacts"]
                if item["artifact_type"] in {"auto-shape", "group"}
            )
        )
        self.assertTrue(
            all(
                item.get("render_page_span") == {"start": 1, "end": 1}
                for item in result["artifact_index"]["artifacts"]
            )
        )

    def test_docx_builder_uses_heading_titles_and_aliases(self) -> None:
        workspace = self.make_workspace()
        source_path = workspace.source_dir / "structured.docx"
        self.create_docx_with_structure(source_path)
        source_dir = workspace.knowledge_base_staging_dir / "sources" / "source-structured-docx"
        source_dir.mkdir(parents=True, exist_ok=True)

        with (
            mock.patch(
                "docmason.knowledge.convert_office_to_pdf",
                return_value=(Path("dummy.pdf"), []),
            ),
            mock.patch(
                "docmason.knowledge.render_pdf_document",
                return_value=(["renders/page-001.png"], []),
            ),
        ):
            _source_manifest, evidence_manifest = build_docx_source(
                workspace,
                source_path,
                {
                    "source_id": "source-structured-docx",
                    "source_fingerprint": "fingerprint-structured-docx",
                    "prior_paths": [],
                    "document_type": "docx",
                    "first_seen_at": "2026-03-22T00:00:00Z",
                    "last_seen_at": "2026-03-22T00:00:00Z",
                    "identity_confidence": "new",
                },
                source_dir,
                soffice_binary="soffice",
            )

        first_unit = evidence_manifest["units"][0]
        self.assertEqual(first_unit["title"], "Implementation Overview")
        self.assertIn("Implementation Overview", first_unit["heading_aliases"])
        structure = read_json(source_dir / first_unit["structure_asset"])
        self.assertTrue(structure["procedure_spans"])
        self.assertTrue(structure["captions"])

    def test_docx_builder_extracts_embedded_picture_focus_assets(self) -> None:
        from PIL import Image

        workspace = self.make_workspace()
        source_path = workspace.source_dir / "figure.docx"
        self.create_docx_with_embedded_picture(source_path)
        source_dir = workspace.knowledge_base_staging_dir / "sources" / "source-figure-docx"
        source_dir.mkdir(parents=True, exist_ok=True)

        def fake_render_pdf_document(
            _pdf_path: Path,
            output_dir: Path,
        ) -> tuple[list[str], list[dict[str, str]]]:
            output_dir.mkdir(parents=True, exist_ok=True)
            Image.new("RGB", (200, 120), color=(255, 255, 255)).save(
                output_dir / "page-001.png",
                format="PNG",
            )
            return (["renders/page-001.png"], [])

        with (
            mock.patch(
                "docmason.knowledge.convert_office_to_pdf",
                return_value=(Path("dummy.pdf"), []),
            ),
            mock.patch(
                "docmason.knowledge.render_pdf_document",
                side_effect=fake_render_pdf_document,
            ),
        ):
            source_manifest, evidence_manifest = build_docx_source(
                workspace,
                source_path,
                {
                    "source_id": "source-figure-docx",
                    "source_fingerprint": "fingerprint-figure-docx",
                    "prior_paths": [],
                    "document_type": "docx",
                    "first_seen_at": "2026-03-22T00:00:00Z",
                    "last_seen_at": "2026-03-22T00:00:00Z",
                    "identity_confidence": "new",
                },
                source_dir,
                soffice_binary="soffice",
            )
        write_json(source_dir / "source_manifest.json", source_manifest)
        write_json(source_dir / "evidence_manifest.json", evidence_manifest)

        artifact_index = read_json(source_dir / "artifact_index.json")
        picture_artifact = next(
            item
            for item in artifact_index["artifacts"]
            if item.get("artifact_type") == "picture"
        )
        self.assertTrue(picture_artifact["focus_render_assets"])
        self.assertTrue(picture_artifact["focus_render_assets"][0].startswith("media/rId"))
        self.assertTrue((source_dir / picture_artifact["focus_render_assets"][0]).exists())

        layout = read_json(source_dir / "visual_layout" / "section-001.json")
        picture_region = next(
            item
            for item in layout["regions"]
            if item.get("artifact_type") == "picture"
        )
        self.assertEqual(
            picture_region["focus_render_assets"][0],
            picture_artifact["focus_render_assets"][0],
        )
        self.assertTrue(
            source_artifact_contract_complete(source_dir, document_type="docx")
        )

    def test_locate_previous_source_dir_prefers_richer_current_semantic_snapshot(self) -> None:
        workspace = self.make_workspace()
        source_id = "source-rich-current"
        staging_dir = workspace.knowledge_base_staging_dir / "sources" / source_id
        current_dir = workspace.knowledge_base_current_dir / "sources" / source_id
        staging_dir.mkdir(parents=True, exist_ok=True)
        current_dir.mkdir(parents=True, exist_ok=True)
        for source_dir in (staging_dir, current_dir):
            write_json(source_dir / "source_manifest.json", {"source_id": source_id})
            write_json(
                source_dir / "evidence_manifest.json",
                {
                    "source_id": source_id,
                    "document_type": "pdf",
                    "units": [],
                },
            )
            write_json(
                source_dir / "artifact_index.json",
                {"source_id": source_id, "artifacts": []},
            )
        write_json(current_dir / "knowledge.json", {"source_id": source_id})
        (current_dir / "summary.md").write_text("summary", encoding="utf-8")
        overlay_dir = current_dir / "semantic_overlay"
        overlay_dir.mkdir(parents=True, exist_ok=True)
        write_json(
            overlay_dir / "page-001.json",
            {
                "artifact_type": "semantic-overlay",
                "source_id": source_id,
                "unit_id": "page-001",
            },
        )

        chosen = locate_previous_source_dir(workspace, source_id)
        self.assertEqual(chosen, current_dir)

    def test_sanitize_text_replaces_unpaired_surrogates(self) -> None:
        self.assertEqual(sanitize_text("ok\ud835bad"), "ok?bad")

    def test_status_reports_invalid_stage_after_pending_sync(self) -> None:
        workspace = self.make_workspace()
        self.mark_environment_ready(workspace)
        self.create_pdf(workspace.source_dir / "example.pdf")

        report = sync_workspace(workspace, autonomous=False)
        self.assertEqual(report.payload["sync_status"], "pending-synthesis")

        status = status_workspace(workspace, editable_install_probe=self.ready_probe)
        self.assertEqual(status.payload["stage"], "knowledge-base-invalid")
        self.assertEqual(status.payload["knowledge_base"]["validation_status"], "pending-synthesis")

    def test_validate_and_publish_after_seeded_agent_outputs(self) -> None:
        workspace = self.make_workspace()
        self.mark_environment_ready(workspace)
        self.create_pdf(workspace.source_dir / "example.pdf")

        pending = sync_workspace(workspace, autonomous=False)
        source_id = pending.payload["pending_sources"][0]["source_id"]
        source_dir = workspace.knowledge_base_staging_dir / "sources" / source_id
        self.seed_agent_outputs(source_dir)

        validation = validate_knowledge_base(workspace, target="staging")
        self.assertEqual(validation.exit_code, 0)
        self.assertEqual(validation.payload["validation_status"], "valid")

        published = sync_workspace(workspace)
        self.assertEqual(published.exit_code, 2)
        self.assertEqual(published.payload["sync_status"], "valid")
        self.assertIn(
            published.payload["hybrid_enrichment"]["mode"],
            {"candidate-prepared", "partially-covered"},
        )
        self.assertTrue(published.payload["lane_b_follow_up"]["work_path"])
        self.assertTrue(workspace.current_publish_manifest_path.exists())
        self.assertTrue((workspace.knowledge_base_dir / "current-pointer.json").exists())
        self.assertTrue(workspace.knowledge_base_current_dir.is_symlink())
        self.assertTrue(
            (
                workspace.knowledge_base_current_dir / "sources" / source_id / "knowledge.json"
            ).exists()
        )
        self.assertTrue(
            (
                workspace.knowledge_base_current_dir
                / "sources"
                / source_id
                / "derived_affordances.json"
            ).exists()
        )
        current_publish_manifest = read_json(workspace.current_publish_manifest_path)
        current_validation = read_json(workspace.current_validation_report_path)
        self.assertIsInstance(current_publish_manifest.get("snapshot_id"), str)
        self.assertFalse(str(current_publish_manifest.get("snapshot_id")).startswith("unknown-"))
        self.assertEqual(
            current_validation.get("source_signature"),
            current_publish_manifest.get("published_source_signature"),
        )

        status = status_workspace(workspace, editable_install_probe=self.ready_probe)
        self.assertEqual(status.payload["stage"], "knowledge-base-present")

    def test_sync_stages_governed_lane_b_follow_up_batch(self) -> None:
        workspace = self.make_workspace()
        self.mark_environment_ready(workspace)
        self.create_pdf_with_full_page_image(workspace.source_dir / "scan.pdf")

        pending = sync_workspace(workspace, autonomous=False)
        source_id = pending.payload["pending_sources"][0]["source_id"]
        source_dir = workspace.knowledge_base_staging_dir / "sources" / source_id
        self.seed_agent_outputs(source_dir)

        published = sync_workspace(workspace)

        self.assertEqual(published.payload["sync_status"], "valid")
        self.assertEqual(
            published.payload["hybrid_enrichment"]["mode"],
            "candidate-prepared",
        )
        follow_up = published.payload["lane_b_follow_up"]
        follow_up_summary = published.payload["lane_b_follow_up_summary"]
        self.assertTrue(follow_up["triggered"])
        self.assertEqual(follow_up["state"], "running")
        self.assertEqual(follow_up_summary["state"], "running")
        self.assertGreaterEqual(follow_up_summary["selected_source_count"], 1)
        self.assertTrue((workspace.root / follow_up["work_path"]).exists())
        self.assertLessEqual(follow_up["selected_unit_count"], 12)
        self.assertLessEqual(len(follow_up["selected_source_ids"]), 4)
        manifest = load_shared_job(workspace, str(follow_up["job_id"]))
        self.assertEqual(manifest["job_family"], "lane-b")
        self.assertEqual(manifest["scope"]["target"], "staging")
        self.assertIn("Lane B follow-up: state=running", "\n".join(published.lines))

    def test_sync_settles_governed_lane_b_follow_up_after_overlay_write(self) -> None:
        workspace = self.make_workspace()
        self.mark_environment_ready(workspace)
        self.create_pdf_with_full_page_image(workspace.source_dir / "scan.pdf")

        pending = sync_workspace(workspace, autonomous=False)
        source_id = pending.payload["pending_sources"][0]["source_id"]
        source_dir = workspace.knowledge_base_staging_dir / "sources" / source_id
        self.seed_agent_outputs(source_dir)

        first_sync = sync_workspace(workspace)
        follow_up = first_sync.payload["lane_b_follow_up"]
        work_packet = read_json(workspace.root / follow_up["work_path"])
        selected_source = work_packet["sources"][0]

        for unit in selected_source["units"]:
            focus_assets = [
                asset
                for asset in unit.get("target_focus_render_assets", [])
                if isinstance(asset, str) and asset
            ]
            render_assets = [
                asset
                for asset in unit.get("target_render_assets", [])
                if isinstance(asset, str) and asset
            ]
            consumed_inputs = {}
            if focus_assets:
                consumed_inputs["focus_render_assets"] = focus_assets
            if render_assets:
                consumed_inputs["render_assets"] = render_assets
            write_semantic_overlay(
                source_dir,
                {
                    "source_id": source_id,
                    "unit_id": unit["unit_id"],
                    "eligible_reason": unit.get("eligible_reason") or "image-only-page",
                    "consumed_inputs": consumed_inputs,
                    "semantic_labels": [
                        {
                            "label": slot,
                            "text": f"Covered {slot} for the staged multimodal unit.",
                            "confidence": "high",
                        }
                        for slot in unit.get("required_overlay_slots", [])
                    ],
                    "artifact_annotations": [],
                    "cross_region_relations": [],
                    "uncertainty_notes": [],
                },
            )

        second_sync = sync_workspace(workspace)

        self.assertEqual(second_sync.payload["hybrid_enrichment"]["mode"], "covered")
        settled = second_sync.payload["lane_b_follow_up"]["settled_job"]
        self.assertEqual(settled["job_id"], follow_up["job_id"])
        self.assertEqual(load_shared_job(workspace, str(follow_up["job_id"]))["status"], "completed")

    def test_snapshot_retention_keeps_current_run_baseline_and_review_pins(self) -> None:
        workspace = self.make_workspace()
        self.seed_snapshot_version(
            workspace,
            snapshot_id="snapshot-current",
            published_at="2026-03-28T00:00:00Z",
            source_signature="sig-current",
        )
        self.seed_snapshot_version(
            workspace,
            snapshot_id="snapshot-run",
            published_at="2026-03-27T00:00:00Z",
            source_signature="sig-run",
        )
        self.seed_snapshot_version(
            workspace,
            snapshot_id="snapshot-baseline",
            published_at="2026-03-26T00:00:00Z",
            source_signature="sig-baseline",
        )
        self.seed_snapshot_version(
            workspace,
            snapshot_id="snapshot-review",
            published_at="2026-03-25T00:00:00Z",
            source_signature="sig-review",
        )
        self.seed_snapshot_version(
            workspace,
            snapshot_id="snapshot-recent-a",
            published_at="2026-03-24T00:00:00Z",
            source_signature="sig-recent-a",
        )
        self.seed_snapshot_version(
            workspace,
            snapshot_id="snapshot-recent-b",
            published_at="2026-03-23T00:00:00Z",
            source_signature="sig-recent-b",
        )
        self.seed_snapshot_version(
            workspace,
            snapshot_id="snapshot-expired",
            published_at="2026-03-10T00:00:00Z",
            source_signature="sig-expired",
        )
        write_json(
            workspace.knowledge_base_dir / "current-pointer.json",
            {"snapshot_id": "snapshot-current"},
        )
        write_json(
            workspace.runs_dir / "run-1" / "state.json",
            {
                "status": "active",
                "version_context": {"published_snapshot_id": "snapshot-run"},
            },
        )
        write_json(
            workspace.eval_baseline_path("broad"),
            {"version_context": {"published_snapshot_id": "snapshot-baseline"}},
        )
        write_json(
            workspace.snapshot_pins_path,
            {
                "pins": [
                    {
                        "snapshot_id": "snapshot-review",
                        "pin_kind": "review-case",
                        "pin_id": "case-1",
                    }
                ]
            },
        )

        retention = apply_snapshot_retention(workspace)

        self.assertTrue(workspace.knowledge_version_dir("snapshot-current").exists())
        self.assertTrue(workspace.knowledge_version_dir("snapshot-run").exists())
        self.assertTrue(workspace.knowledge_version_dir("snapshot-baseline").exists())
        self.assertTrue(workspace.knowledge_version_dir("snapshot-review").exists())
        self.assertTrue(workspace.knowledge_version_dir("snapshot-recent-a").exists())
        self.assertTrue(workspace.knowledge_version_dir("snapshot-recent-b").exists())
        self.assertFalse(workspace.knowledge_version_dir("snapshot-expired").exists())
        self.assertIn("snapshot-expired", retention["deleted_snapshot_ids"])
        review_entry = next(
            item for item in retention["snapshots"] if item["snapshot_id"] == "snapshot-review"
        )
        self.assertIn("review-case:case-1", review_entry["pin_reasons"])

    def test_validate_rejects_placeholder_knowledge(self) -> None:
        workspace = self.make_workspace()
        self.create_pdf(workspace.source_dir / "example.pdf")

        pending = sync_workspace(workspace, autonomous=False)
        source_id = pending.payload["pending_sources"][0]["source_id"]
        source_dir = workspace.knowledge_base_staging_dir / "sources" / source_id
        source_manifest = read_json(source_dir / "source_manifest.json")
        evidence_manifest = read_json(source_dir / "evidence_manifest.json")
        first_unit_id = evidence_manifest["units"][0]["unit_id"]
        write_json(
            source_dir / "knowledge.json",
            {
                "source_id": source_manifest["source_id"],
                "source_fingerprint": source_manifest["source_fingerprint"],
                "title": "TODO",
                "source_language": "en",
                "summary_en": "TODO",
                "summary_source": "TODO",
                "document_type": source_manifest["document_type"],
                "key_points": [],
                "entities": [],
                "claims": [{"citations": [{"unit_id": first_unit_id, "support": "primary"}]}],
                "known_gaps": [],
                "ambiguities": [],
                "confidence": {"level": "low"},
                "citations": [{"unit_id": first_unit_id, "support": "primary"}],
                "related_sources": [],
            },
        )
        (source_dir / "summary.md").write_text(
            "\n".join(
                [
                    "# TODO",
                    "",
                    f"Source ID: {source_manifest['source_id']}",
                    "",
                    "## English Summary",
                    "TODO",
                    "",
                    "## Source-Language Summary",
                    "TODO",
                    "",
                ]
            ),
            encoding="utf-8",
        )

        validation = validate_knowledge_base(workspace, target="staging")

        self.assertEqual(validation.exit_code, 1)
        self.assertEqual(validation.payload["status"], ACTION_REQUIRED)
        self.assertEqual(validation.payload["validation_status"], "blocking-errors")
        report = json.dumps(validation.payload["validation"], sort_keys=True)
        self.assertIn("placeholder", report.lower())

    def test_sync_waits_for_sync_lease_before_mutating_shared_state(self) -> None:
        workspace = self.make_workspace()
        result: dict[str, object] = {}
        finished = threading.Event()

        def run_sync() -> None:
            result["payload"] = sync_workspace(workspace, autonomous=False)
            finished.set()

        with workspace_lease(workspace, "sync", timeout_seconds=1.0):
            thread = threading.Thread(target=run_sync)
            thread.start()
            time.sleep(0.2)
            self.assertFalse(
                finished.is_set(),
                "Sync should wait while another sync lease is active.",
            )
        thread.join(timeout=5.0)
        self.assertFalse(thread.is_alive())
        payload = result["payload"]
        self.assertEqual(payload.payload["sync_status"], "valid")

    def test_workspace_lease_recovers_from_conflicting_file_path(self) -> None:
        workspace = self.make_workspace()
        conflicting_path = lease_dir(workspace, "conversation:test-thread")
        conflicting_path.parent.mkdir(parents=True, exist_ok=True)
        conflicting_path.write_text("stale", encoding="utf-8")

        with workspace_lease(workspace, "conversation:test-thread") as payload:
            self.assertTrue(conflicting_path.is_dir())
            lease_payload = read_json(conflicting_path / "lease.json")
            self.assertEqual(lease_payload["owner"], payload["owner"])


if __name__ == "__main__":
    unittest.main()
