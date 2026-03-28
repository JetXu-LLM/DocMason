"""Phase 4 incremental sync, evidence preparation, validation, retrieval, and trace."""

from __future__ import annotations

import hashlib
import io
import os
import re
import shutil
import subprocess
import sys
import tempfile
import uuid
import warnings
import zipfile
from collections import Counter
from collections.abc import Iterable
from contextlib import redirect_stderr
from datetime import UTC, datetime
from difflib import SequenceMatcher
from pathlib import Path
from time import perf_counter
from typing import Any, Literal
from urllib.parse import urlparse

from .affordances import (
    DEFAULT_AFFORDANCE_FILENAME,
    derive_source_affordances,
    merge_derived_affordances,
    validate_derived_affordances,
)
from .artifacts import validate_artifact_index, validate_pdf_document
from .control_plane import repair_stale_shared_jobs
from .coordination import workspace_lease
from .email_sources import parse_email_source
from .evidence_artifacts import (
    compile_docx_visual_compatibility,
    compile_pdf_visual_artifacts,
    compile_pptx_visual_artifacts,
    compile_xlsx_artifacts,
    write_empty_artifact_index,
)
from .hybrid import (
    build_source_hybrid_packet,
    focus_render_contract_complete,
    materialize_focus_render_assets,
    summarize_hybrid_work,
)
from .interaction import (
    build_promoted_interaction_memories,
    interaction_ingest_snapshot,
    load_promoted_interaction_contexts,
    mark_promoted_interaction_entries,
    repair_interaction_memory_related_sources,
)
from .project import (
    WorkspacePaths,
    isoformat_timestamp,
    read_json,
    relative_paths,
    source_index,
    source_inventory_signature,
    source_type_definition,
    source_type_definition_for_path,
    supported_source_documents,
    sync_state,
    write_json,
)
from .projections import projection_state_summary, queue_projection_refresh
from .retrieval import (
    build_retrieval_artifacts,
    build_trace_artifacts,
    normalize_filename_stem,
)
from .semantic_overlays import (
    collect_semantic_overlay_assets,
    load_semantic_overlays,
    overlay_confidence,
    overlay_search_strings,
    validate_hybrid_work,
    validate_semantic_overlay,
)
from .source_references import (
    enrich_evidence_manifest_reference_fields,
    enrich_source_manifest_reference_fields,
)
from .text_sources import ParsedUnit, parse_text_source
from .versioning import publish_staging_snapshot

PLACEHOLDER_TERMS = ("todo", "tbd", "placeholder", "lorem ipsum", "fill in")
DOCX_ORDERED_STEP_PATTERN = re.compile(r"^\s*(\d+)(?:[.)-])\s+")
DOCX_CAPTION_PREFIX_PATTERN = re.compile(
    r"^(figure|fig\.|table|chart|diagram|exhibit)\s+[A-Za-z0-9.-]+\s*[:.-]?\s+",
    re.IGNORECASE,
)
REQUIRED_KNOWLEDGE_KEYS = (
    "source_id",
    "source_fingerprint",
    "title",
    "source_language",
    "summary_en",
    "summary_source",
    "document_type",
    "key_points",
    "entities",
    "claims",
    "known_gaps",
    "ambiguities",
    "confidence",
    "citations",
    "related_sources",
)
KNOWLEDGE_STATUSES = {"valid", "warnings", "blocking-errors", "pending-synthesis", "not-run"}
TEXT_DOCUMENT_TYPES = {"markdown", "plaintext", "mdx", "yaml", "tex", "csv", "tsv"}
NON_RENDERED_DOCUMENT_TYPES = TEXT_DOCUMENT_TYPES | {"email"}
EMAIL_MAX_ATTACHMENT_DEPTH = 2
DERIVED_SOURCE_ORIGIN = "derived-attachment"
EMAIL_ATTACHMENT_RELATION_TYPE = "email-attachment"
LEGACY_OFFICE_EXTENSION_MAP = {"ppt": "pptx", "doc": "docx", "xls": "xlsx"}
BENIGN_THIRD_PARTY_DIAGNOSTIC_SUBSTRINGS = (
    "conditional formatting extension is not supported and will be removed",
    "data validation extension is not supported and will be removed",
    "cannot parse header or footer so it will be ignored",
    "ignoring wrong pointing object",
)


def utc_now() -> str:
    """Return the current UTC timestamp in ISO 8601 form."""
    return datetime.now(tz=UTC).isoformat().replace("+00:00", "Z")


def file_sha256(path: Path) -> str:
    """Return the SHA-256 hex digest for a file."""
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        while True:
            chunk = handle.read(1024 * 1024)
            if not chunk:
                break
            digest.update(chunk)
    return digest.hexdigest()


def sanitize_text(value: Any) -> str:
    """Normalize structured content into compact human-readable text."""
    if value is None:
        return ""
    text = str(value).replace("\r\n", "\n").replace("\r", "\n")
    # Some Office/PDF extractors surface lone surrogate code points that cannot be
    # persisted with strict UTF-8 encoding. Replace them conservatively.
    text = text.encode("utf-8", errors="replace").decode("utf-8")
    return "\n".join(line.rstrip() for line in text.splitlines()).strip()


def detect_source_language(texts: Iterable[str]) -> str:
    """Make a conservative language guess from extracted text."""
    joined = " ".join(texts)
    if not joined.strip():
        return "unknown"
    ascii_ratio = sum(1 for character in joined if ord(character) < 128) / max(len(joined), 1)
    if ascii_ratio > 0.95:
        return "en"
    return "mixed-or-non-en"


def module_available(module_name: str) -> bool:
    """Return whether a Python module can be imported."""
    try:
        __import__(module_name)
    except ImportError:
        return False
    return True


def _is_benign_third_party_diagnostic(message: str) -> bool:
    normalized = " ".join(message.lower().split())
    return any(token in normalized for token in BENIGN_THIRD_PARTY_DIAGNOSTIC_SUBSTRINGS)


def _normalize_third_party_diagnostics(lines: list[str]) -> list[str]:
    messages: list[str] = []
    for value in lines:
        text = " ".join(sanitize_text(value).split()).strip()
        if not text or _is_benign_third_party_diagnostic(text):
            continue
        messages.append(text)
    return list(dict.fromkeys(messages))


class _ThirdPartyDiagnosticCapture:
    """Capture third-party warnings and stderr without polluting CLI output."""

    def __init__(self) -> None:
        self.messages: list[str] = []
        self._warning_context: Any = None
        self._warning_records: list[Any] = []
        self._stderr_redirect: Any = None
        self._stderr_proxy = io.StringIO()
        self._stderr_file: Any = None
        self._stderr_fd: int | None = None
        self._fd_redirected = False

    def __enter__(self) -> _ThirdPartyDiagnosticCapture:
        self._warning_context = warnings.catch_warnings(record=True)
        self._warning_records = self._warning_context.__enter__()
        warnings.simplefilter("always")
        self._stderr_file = tempfile.TemporaryFile(mode="w+b")
        try:
            self._stderr_fd = os.dup(2)
            try:
                sys.stderr.flush()
            except Exception:
                pass
            os.dup2(self._stderr_file.fileno(), 2)
            self._fd_redirected = True
        except OSError:
            self._stderr_fd = None
            self._fd_redirected = False
        self._stderr_redirect = redirect_stderr(self._stderr_proxy)
        self._stderr_redirect.__enter__()
        return self

    def __exit__(self, exc_type: Any, exc: Any, tb: Any) -> Literal[False]:
        if self._stderr_redirect is not None:
            self._stderr_redirect.__exit__(exc_type, exc, tb)
        if self._fd_redirected and self._stderr_fd is not None:
            try:
                sys.stderr.flush()
            except Exception:
                pass
            os.dup2(self._stderr_fd, 2)
            os.close(self._stderr_fd)
        raw_lines: list[str] = []
        if self._stderr_file is not None:
            self._stderr_file.seek(0)
            raw_lines.extend(
                self._stderr_file.read().decode("utf-8", errors="replace").splitlines()
            )
            self._stderr_file.close()
        raw_lines.extend(self._stderr_proxy.getvalue().splitlines())
        raw_lines.extend(
            str(getattr(record, "message", record))
            for record in self._warning_records
            if str(getattr(record, "message", record)).strip()
        )
        if self._warning_context is not None:
            self._warning_context.__exit__(exc_type, exc, tb)
        self.messages = _normalize_third_party_diagnostics(raw_lines)
        return False


def _append_captured_diagnostic_failures(
    failures: list[dict[str, str]],
    capture: _ThirdPartyDiagnosticCapture,
    *,
    stage: str,
) -> None:
    existing = {(item.get("stage"), item.get("detail")) for item in failures}
    for message in capture.messages:
        key = (stage, message)
        if key in existing:
            continue
        failures.append({"stage": stage, "detail": message})
        existing.add(key)


def pdf_renderer_snapshot() -> dict[str, Any]:
    """Describe PDF rendering dependency readiness."""
    missing: list[str] = []
    for module_name in ("pypdfium2", "pypdf", "PIL"):
        if not module_available(module_name):
            missing.append(module_name)
    if not (module_available("pymupdf") or module_available("fitz")):
        missing.append("PyMuPDF")
    ready = not missing
    detail = "PDF rendering and extraction dependencies are available."
    if missing:
        detail = "Missing Python packages for PDF rendering/extraction: " + ", ".join(missing)
    return {"ready": ready, "detail": detail, "missing": missing}


def find_soffice_binary() -> str | None:
    """Resolve a validated LibreOffice soffice executable from common locations."""
    for candidate in (
        shutil.which("soffice"),
        shutil.which("libreoffice"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ):
        validation = validate_soffice_binary(candidate)
        if validation["ready"]:
            return str(validation["binary"])
    return None


def office_source_documents(paths: WorkspacePaths) -> list[Path]:
    """Return Office documents that require the LibreOffice rendering path."""
    return [
        path
        for path in supported_source_documents(paths)
        if (definition := source_type_definition_for_path(path)) is not None
        and definition.requires_office_renderer
    ]


def office_renderer_snapshot(paths: WorkspacePaths) -> dict[str, Any]:
    """Describe Office renderer readiness for the current workspace."""
    discovery = discover_soffice_binary()
    validation = validate_soffice_binary(discovery)
    soffice_binary = str(validation["binary"]) if validation["ready"] else None
    office_sources = office_source_documents(paths)
    required = bool(office_sources)
    ready = soffice_binary is not None
    if not required:
        detail = "LibreOffice is optional until PowerPoint, Word, or Excel sources are present."
    elif ready:
        version_suffix = (
            f" ({validation['version']})"
            if isinstance(validation.get("version"), str) and validation["version"]
            else ""
        )
        detail = f"LibreOffice rendering is available at {soffice_binary}{version_suffix}."
    else:
        candidate_path = validation.get("binary") or discovery
        if candidate_path:
            detail = (
                "LibreOffice `soffice` is required to render PowerPoint, Word, and Excel "
                f"sources, but the detected candidate `{candidate_path}` is not a validated "
                f"LibreOffice install. {validation['detail']}"
            )
        else:
            detail = (
                "LibreOffice `soffice` is required to render PowerPoint, Word, and Excel "
                "sources, but it is not available."
            )
    return {
        "ready": ready,
        "required": required,
        "binary": soffice_binary,
        "candidate_binary": validation.get("binary") or discovery,
        "validation_detail": validation["detail"],
        "validated": bool(validation["ready"]),
        "version": validation.get("version"),
        "detail": detail,
        "office_sources": relative_paths(paths, office_sources),
    }


def discover_soffice_binary() -> str | None:
    """Locate the most likely LibreOffice command path without trusting it yet."""
    for candidate in (
        shutil.which("soffice"),
        shutil.which("libreoffice"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ):
        if candidate and Path(candidate).exists():
            return str(Path(candidate))
    return None


def validate_soffice_binary(candidate: str | None) -> dict[str, Any]:
    """Validate that a detected command is a usable LibreOffice renderer."""
    if not candidate:
        return {
            "ready": False,
            "binary": None,
            "version": None,
            "detail": "No LibreOffice command candidate was detected.",
        }
    binary = Path(candidate)
    if not binary.exists():
        return {
            "ready": False,
            "binary": str(binary),
            "version": None,
            "detail": "The detected LibreOffice command path does not exist.",
        }
    if not os.access(binary, os.X_OK):
        return {
            "ready": False,
            "binary": str(binary),
            "version": None,
            "detail": "The detected LibreOffice command path is not executable.",
        }
    try:
        completed = subprocess.run(
            [str(binary), "--version"],
            capture_output=True,
            text=True,
            check=False,
        )
    except OSError as exc:
        return {
            "ready": False,
            "binary": str(binary),
            "version": None,
            "detail": f"The detected LibreOffice command failed to execute: {exc}.",
        }
    output = (completed.stdout or completed.stderr or "").strip()
    if completed.returncode != 0:
        return {
            "ready": False,
            "binary": str(binary),
            "version": None,
            "detail": (
                "The detected LibreOffice command failed the version probe: "
                f"{output or f'exit code {completed.returncode}'}."
            ),
        }
    normalized_output = output.lower()
    if "libreoffice" not in normalized_output:
        return {
            "ready": False,
            "binary": str(binary),
            "version": output or None,
            "detail": "The detected command did not identify itself as LibreOffice.",
        }
    return {
        "ready": True,
        "binary": str(binary),
        "version": output or None,
        "detail": "Validated LibreOffice renderer capability.",
    }


def ensure_directory(path: Path) -> None:
    """Create a directory tree when it does not already exist."""
    path.mkdir(parents=True, exist_ok=True)


def write_text(path: Path, content: str) -> None:
    """Persist a UTF-8 text file and ensure its parent directory exists."""
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")


def build_relative_path_lineage(source_relative_path: Path) -> list[str]:
    """Build the path lineage beneath `original_doc/` for trust modeling."""
    lineage: list[str] = []
    current = Path()
    for part in source_relative_path.parts[:-1]:
        current /= part
        lineage.append(str(current))
    return lineage


def build_trust_prior_payload(
    source_relative_path: Path,
    *,
    modified_at: str | None,
) -> dict[str, Any]:
    """Build the trust-prior payload from a relative path plus timestamp."""
    return {
        "first_level_subtree": source_relative_path.parts[0]
        if len(source_relative_path.parts) > 1
        else "",
        "local_branch_depth": max(len(source_relative_path.parts) - 1, 0),
        "relative_path_lineage": build_relative_path_lineage(source_relative_path),
        "modified_at": modified_at,
        "graph_centrality": None,
        "corroboration": None,
    }


def build_trust_prior(path: Path, source_relative_path: Path) -> dict[str, Any]:
    """Build the trust-prior payload for a source or evidence unit."""
    return build_trust_prior_payload(
        source_relative_path,
        modified_at=isoformat_timestamp(path.stat().st_mtime),
    )


def source_relative_path_from_current_path(current_path: str) -> Path:
    """Return the corpus-relative path lineage for a source-like current path."""
    normalized = str(current_path or "")
    if normalized.startswith("original_doc/"):
        normalized = normalized.removeprefix("original_doc/")
    normalized = normalized.split("#attachment/", 1)[0]
    return Path(normalized)


def build_trust_prior_from_source_entry(source_entry: dict[str, Any]) -> dict[str, Any]:
    """Build trust-prior inputs for original or derived corpus sources."""
    return build_trust_prior_payload(
        source_relative_path_from_current_path(str(source_entry.get("current_path") or "")),
        modified_at=str(source_entry.get("modified_at") or "") or None,
    )


def source_content_path(paths: WorkspacePaths, source_entry: dict[str, Any]) -> Path:
    """Return the readable filesystem path for a source entry."""
    content_path = source_entry.get("content_path")
    if isinstance(content_path, str) and content_path:
        return Path(content_path)
    return paths.root / str(source_entry["current_path"])


def normalize_source_entry_for_build(
    paths: WorkspacePaths,
    source_path: Path,
    source_entry: dict[str, Any],
) -> dict[str, Any]:
    """Backfill missing build-time fields for direct builder callers and tests."""
    if isinstance(source_entry.get("current_path"), str) and source_entry.get("current_path"):
        return source_entry
    current_path = (
        str(source_path.relative_to(paths.root))
        if source_path.is_absolute() and source_path.is_relative_to(paths.root)
        else str(source_path)
    )
    stat = source_path.stat() if source_path.exists() else None
    return {
        **source_entry,
        "current_path": current_path,
        "source_extension": source_entry.get("source_extension")
        or source_path.suffix.lower().lstrip("."),
        "file_size": source_entry.get("file_size")
        if source_entry.get("file_size") is not None
        else (stat.st_size if stat is not None else 0),
        "modified_at": source_entry.get("modified_at")
        or (isoformat_timestamp(stat.st_mtime) if stat is not None else None),
        "prior_paths": source_entry.get("prior_paths", []),
        "path_history": source_entry.get("path_history", [current_path]),
        "first_seen_at": source_entry.get("first_seen_at", utc_now()),
        "last_seen_at": source_entry.get("last_seen_at", utc_now()),
        "identity_confidence": source_entry.get("identity_confidence", "unknown"),
    }


def source_definition_for_entry(source_entry: dict[str, Any], source_path: Path) -> Any:
    """Resolve the source-type definition for original or derived source entries."""
    extension = source_entry.get("source_extension")
    definition = source_type_definition(str(extension)) if isinstance(extension, str) else None
    return definition or source_type_definition_for_path(source_path)


def write_bytes(path: Path, payload: bytes) -> None:
    """Persist raw bytes and ensure the parent directory exists."""
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(payload)


def read_zip_entries(path: Path) -> list[str]:
    """Return the ordered ZIP entry list for an OOXML file."""
    with zipfile.ZipFile(path) as archive:
        return sorted(archive.namelist())


def relocation_candidate_score(
    current_path: str,
    current_size: int,
    candidate_entry: dict[str, Any],
) -> float:
    """Return a deterministic score for relocation-heuristic identity reuse."""
    candidate_path = candidate_entry.get("current_path")
    if not isinstance(candidate_path, str):
        return 0.0
    current_name = normalize_filename_stem(Path(current_path).stem)
    candidate_name = normalize_filename_stem(Path(candidate_path).stem)
    if not current_name or not candidate_name:
        return 0.0
    name_score = SequenceMatcher(None, current_name, candidate_name).ratio()
    candidate_size = candidate_entry.get("file_size")
    if isinstance(candidate_size, int) and candidate_size > 0 and current_size > 0:
        size_score = min(candidate_size, current_size) / max(candidate_size, current_size)
    else:
        size_score = 0.9
    if name_score < 0.65 or size_score < 0.9:
        return 0.0
    return (0.7 * name_score) + (0.3 * size_score)


def append_unique_strings(values: Iterable[str]) -> list[str]:
    """Deduplicate strings while preserving order."""
    seen: set[str] = set()
    ordered: list[str] = []
    for value in values:
        if value and value not in seen:
            seen.add(value)
            ordered.append(value)
    return ordered


def build_change_record(
    entry: dict[str, Any],
    *,
    previous_path: str | None,
) -> dict[str, Any]:
    """Build a compact change record for the latest sync state."""
    return {
        "source_id": entry.get("source_id"),
        "current_path": entry.get("current_path"),
        "previous_path": previous_path,
        "document_type": entry.get("document_type"),
        "change_classification": entry.get("change_classification"),
        "change_traits": entry.get("change_traits", []),
        "change_reason": entry.get("change_reason"),
        "identity_basis": entry.get("identity_basis", entry.get("identity_confidence")),
        "matched_source_ids": entry.get("matched_source_ids", []),
    }


def _initial_change_traits(
    existing_entry: dict[str, Any] | None,
    *,
    current_path: str,
    fingerprint: str,
) -> list[str]:
    """Return the initial additive change traits available before staging begins."""
    if existing_entry is None:
        return []
    traits: list[str] = []
    previous_path = existing_entry.get("current_path")
    previous_fingerprint = existing_entry.get("source_fingerprint")
    if isinstance(previous_path, str) and previous_path != current_path:
        traits.append("path_changed")
    if isinstance(previous_fingerprint, str) and previous_fingerprint != fingerprint:
        traits.append("binary_changed")
    return traits


def _initial_change_reason(
    *,
    change_classification: str,
    identity_basis: str,
    change_traits: list[str],
) -> str:
    """Return a concise operator-facing explanation for the detected source change."""
    if change_classification == "unchanged":
        return "Path and binary fingerprint matched the previously indexed source."
    if change_classification == "added":
        return "This source is new to the workspace corpus."
    if change_classification == "deleted":
        return "This previously indexed source is no longer present under `original_doc/`."
    if change_classification == "ambiguous":
        return (
            "Multiple historical sources could plausibly match this path; "
            "operator review is required."
        )
    if "path_changed" in change_traits and "binary_changed" in change_traits:
        return (
            "The source path changed and the binary fingerprint also changed; staged evidence must "
            "confirm whether semantic outputs can be reused."
        )
    if "path_changed" in change_traits:
        return (
            f"The source path changed and identity was preserved via `{identity_basis}` matching."
        )
    if "binary_changed" in change_traits:
        return "The source path is stable but the binary fingerprint changed."
    return f"The source was classified as `{change_classification}` via `{identity_basis}`."


def _compute_source_index(
    paths: WorkspacePaths,
    *,
    persist: bool,
) -> tuple[dict[str, Any], list[dict[str, Any]], bool, dict[str, Any]]:
    """Compute the current source identity view, optionally persisting it."""
    now = utc_now()
    existing = source_index(paths)
    existing_sources = existing.get("sources", [])
    if not isinstance(existing_sources, list):
        existing_sources = []

    path_lookup: dict[str, dict[str, Any]] = {}
    all_fingerprint_lookup: dict[str, list[dict[str, Any]]] = {}
    active_existing: list[dict[str, Any]] = []
    inactive_existing: list[dict[str, Any]] = []
    for raw_entry in existing_sources:
        if not isinstance(raw_entry, dict):
            continue
        fingerprint = raw_entry.get("source_fingerprint")
        if isinstance(fingerprint, str):
            all_fingerprint_lookup.setdefault(fingerprint, []).append(raw_entry)
        if raw_entry.get("active", True):
            active_existing.append(raw_entry)
            current_path = raw_entry.get("current_path")
            if isinstance(current_path, str):
                path_lookup[current_path] = raw_entry
        else:
            inactive_existing.append(raw_entry)

    active_entries: list[dict[str, Any]] = []
    ambiguous_match = False
    seen_source_ids: set[str] = set()
    matched_existing_ids: set[str] = set()
    changes: list[dict[str, Any]] = []
    for path in supported_source_documents(paths):
        current_path = str(path.relative_to(paths.root))
        fingerprint = file_sha256(path)
        stat = path.stat()
        file_size = stat.st_size
        definition = source_type_definition_for_path(path)
        if definition is None:
            continue
        document_type = definition.document_type
        existing_entry = path_lookup.get(current_path)
        identity_basis = "path"
        matched_source_ids: list[str] = []
        change_classification = "unchanged"
        previous_path: str | None = None

        if existing_entry is None:
            matches = [
                entry
                for entry in all_fingerprint_lookup.get(fingerprint, [])
                if isinstance(entry.get("source_id"), str)
                and str(entry["source_id"]) not in matched_existing_ids
            ]
            matched_source_ids = [
                str(entry["source_id"])
                for entry in matches
                if isinstance(entry.get("source_id"), str)
            ]
            if len(matches) == 1:
                existing_entry = matches[0]
                identity_basis = "fingerprint"
                change_classification = "moved-or-renamed"
            elif len(matches) > 1:
                ambiguous_match = True
                identity_basis = "ambiguous"
                change_classification = "ambiguous"
            else:
                relocation_candidates = [
                    (candidate, relocation_candidate_score(current_path, file_size, candidate))
                    for candidate in active_existing
                    if candidate.get("active", True)
                    and isinstance(candidate.get("source_id"), str)
                    and str(candidate["source_id"]) not in matched_existing_ids
                    and candidate.get("document_type") == document_type
                ]
                relocation_candidates = [
                    (candidate, score) for candidate, score in relocation_candidates if score > 0
                ]
                relocation_candidates.sort(
                    key=lambda item: (-float(item[1]), str(item[0].get("source_id", "")))
                )
                if len(relocation_candidates) == 1:
                    existing_entry = relocation_candidates[0][0]
                    identity_basis = "relocation-heuristic"
                    change_classification = "moved-or-renamed"
                    matched_source_ids = [str(existing_entry["source_id"])]
                elif len(relocation_candidates) > 1:
                    ambiguous_match = True
                    identity_basis = "ambiguous"
                    change_classification = "ambiguous"
                    matched_source_ids = [
                        str(candidate["source_id"])
                        for candidate, _score in relocation_candidates
                        if isinstance(candidate.get("source_id"), str)
                    ]

        prior_paths: list[str] = []
        first_seen_at = now
        source_id = str(uuid.uuid4())
        change_traits: list[str] = []
        if existing_entry is not None:
            source_id = str(existing_entry.get("source_id", source_id))
            first_seen_at = str(existing_entry.get("first_seen_at", now))
            existing_prior = existing_entry.get("prior_paths", [])
            if isinstance(existing_prior, list):
                prior_paths.extend(str(value) for value in existing_prior)
            old_path = existing_entry.get("current_path")
            if isinstance(old_path, str) and old_path != current_path:
                prior_paths.append(old_path)
                previous_path = old_path
            elif isinstance(old_path, str):
                previous_path = old_path
            if identity_basis == "path":
                change_classification = (
                    "unchanged"
                    if existing_entry.get("source_fingerprint") == fingerprint
                    else "modified"
                )
            change_traits = _initial_change_traits(
                existing_entry,
                current_path=current_path,
                fingerprint=fingerprint,
            )
            matched_existing_ids.add(source_id)
        else:
            identity_basis = "new" if identity_basis != "ambiguous" else identity_basis
            change_classification = (
                "added" if change_classification != "ambiguous" else change_classification
            )

        path_history = (
            append_unique_strings(
                [
                    *(
                        value
                        for value in existing_entry.get("path_history", [])
                        if isinstance(existing_entry, dict)
                        and isinstance(existing_entry.get("path_history"), list)
                        and isinstance(value, str)
                    ),
                    *prior_paths,
                    current_path,
                ]
            )
            if existing_entry is not None
            else [current_path]
        )

        entry = {
            "source_id": source_id,
            "current_path": current_path,
            "prior_paths": append_unique_strings(prior_paths),
            "path_history": path_history,
            "source_fingerprint": fingerprint,
            "file_size": file_size,
            "first_seen_at": first_seen_at,
            "last_seen_at": now,
            "identity_confidence": identity_basis,
            "identity_basis": identity_basis,
            "change_classification": change_classification,
            "change_traits": change_traits,
            "change_reason": _initial_change_reason(
                change_classification=change_classification,
                identity_basis=identity_basis,
                change_traits=change_traits,
            ),
            "ambiguous_match": identity_basis == "ambiguous",
            "matched_source_ids": matched_source_ids,
            "document_type": document_type,
            "source_extension": definition.extension,
            "support_tier": definition.support_tier,
            "archived_at": None,
            "deleted_at": None,
            "active": True,
        }
        active_entries.append(entry)
        seen_source_ids.add(source_id)
        changes.append(build_change_record(entry, previous_path=previous_path))

    archived_entries: list[dict[str, Any]] = []
    for raw_entry in active_existing:
        if not isinstance(raw_entry, dict):
            continue
        source_id_value = raw_entry.get("source_id")
        if not isinstance(source_id_value, str) or source_id_value in seen_source_ids:
            continue
        archived_entry = dict(raw_entry)
        archived_entry["active"] = False
        archived_entry["change_classification"] = "deleted"
        archived_entry["change_traits"] = []
        archived_entry["change_reason"] = _initial_change_reason(
            change_classification="deleted",
            identity_basis=str(archived_entry.get("identity_basis") or "path"),
            change_traits=[],
        )
        archived_entry["archived_at"] = archived_entry.get("archived_at") or now
        archived_entry["deleted_at"] = archived_entry.get("deleted_at") or now
        archived_entries.append(archived_entry)
        changes.append(
            build_change_record(
                archived_entry,
                previous_path=raw_entry.get("current_path"),
            )
        )

    archived_entries.extend(inactive_existing)
    classification_counts = Counter(
        str(change["change_classification"])
        for change in changes
        if change.get("change_classification")
    )
    change_set = {
        "generated_at": now,
        "source_signature": source_inventory_signature(paths),
        "stats": {
            "unchanged": classification_counts.get("unchanged", 0),
            "added": classification_counts.get("added", 0),
            "modified": classification_counts.get("modified", 0),
            "moved_or_renamed": classification_counts.get("moved-or-renamed", 0),
            "deleted": classification_counts.get("deleted", 0),
            "ambiguous": classification_counts.get("ambiguous", 0),
        },
        "changes": changes,
    }

    payload = {
        "generated_at": now,
        "sources": sorted(active_entries, key=lambda item: str(item["current_path"]))
        + sorted(
            archived_entries,
            key=lambda item: (str(item.get("current_path", "")), str(item.get("source_id", ""))),
        ),
    }
    if persist:
        write_json(paths.source_index_path, payload)
    return payload, active_entries, ambiguous_match, change_set


def update_source_index(
    paths: WorkspacePaths,
) -> tuple[dict[str, Any], list[dict[str, Any]], bool, dict[str, Any]]:
    """Update and persist the stable source identity index plus the current change set."""
    return _compute_source_index(paths, persist=True)


def preview_source_changes(
    paths: WorkspacePaths,
) -> tuple[dict[str, Any], list[dict[str, Any]], bool, dict[str, Any]]:
    """Compute the current change set without mutating persisted source-index state."""
    return _compute_source_index(paths, persist=False)


def import_pdf_modules() -> tuple[Any, Any]:
    """Import the PDF dependencies lazily."""
    import pypdfium2 as pdfium  # type: ignore[import-untyped]
    from pypdf import PdfReader

    return pdfium, PdfReader


def render_pdf_document(
    pdf_path: Path,
    renders_dir: Path,
    *,
    prefix: str = "page",
) -> tuple[list[str], list[dict[str, str]]]:
    """Render a PDF into PNG images and return relative asset names."""
    failures: list[dict[str, str]] = []
    rendered_assets: list[str] = []
    pdfium, _ = import_pdf_modules()
    diagnostics = _ThirdPartyDiagnosticCapture()

    try:
        with diagnostics:
            document = pdfium.PdfDocument(str(pdf_path))
            try:
                for index in range(len(document)):
                    page = document[index]
                    bitmap = page.render(scale=2)
                    image = bitmap.to_pil()
                    filename = f"{prefix}-{index + 1:03d}.png"
                    output_path = renders_dir / filename
                    image.save(
                        output_path,
                        format="PNG",
                        compress_level=1,
                        optimize=False,
                    )
                    rendered_assets.append(str(Path("renders") / filename))
            finally:
                document.close()
    except Exception as exc:  # pragma: no cover - defensive against third-party failures
        _append_captured_diagnostic_failures(
            failures,
            diagnostics,
            stage="render-pdf-diagnostic",
        )
        failures.append({"stage": "render-pdf", "detail": str(exc)})
    else:
        _append_captured_diagnostic_failures(
            failures,
            diagnostics,
            stage="render-pdf-diagnostic",
        )
    return rendered_assets, failures


def extract_pdf_text(pdf_path: Path) -> tuple[list[str], list[dict[str, str]]]:
    """Extract page text from a PDF and return page-aligned strings."""
    failures: list[dict[str, str]] = []
    texts: list[str] = []
    _, PdfReader = import_pdf_modules()
    diagnostics = _ThirdPartyDiagnosticCapture()
    try:
        with diagnostics:
            reader = PdfReader(str(pdf_path))
            for page in reader.pages:
                texts.append(sanitize_text(page.extract_text() or ""))
    except Exception as exc:  # pragma: no cover - defensive against third-party failures
        _append_captured_diagnostic_failures(
            failures,
            diagnostics,
            stage="extract-pdf-text-diagnostic",
        )
        failures.append({"stage": "extract-pdf-text", "detail": str(exc)})
    else:
        _append_captured_diagnostic_failures(
            failures,
            diagnostics,
            stage="extract-pdf-text-diagnostic",
        )
    return texts, failures


def convert_office_to_pdf(
    source_path: Path,
    output_dir: Path,
    soffice_binary: str,
) -> tuple[Path | None, list[dict[str, str]]]:
    """Convert an Office document to PDF via LibreOffice."""
    return convert_office_to_format(
        source_path,
        output_dir,
        soffice_binary,
        target_format="pdf",
    )


def _load_workbook_quietly(
    workbook_path: Path,
    *,
    data_only: bool,
    failures: list[dict[str, str]],
    stage: str,
) -> Any | None:
    from openpyxl import load_workbook  # type: ignore[import-untyped]

    diagnostics = _ThirdPartyDiagnosticCapture()
    try:
        with diagnostics:
            workbook = load_workbook(filename=str(workbook_path), data_only=data_only)
    except Exception as exc:  # pragma: no cover - defensive against third-party failures
        _append_captured_diagnostic_failures(failures, diagnostics, stage=f"{stage}-diagnostic")
        failures.append(
            {
                "stage": stage,
                "detail": f"Could not load workbook `{workbook_path.name}`: {exc}",
            }
        )
        return None
    _append_captured_diagnostic_failures(failures, diagnostics, stage=f"{stage}-diagnostic")
    return workbook


def _save_workbook_quietly(
    workbook: Any,
    output_path: Path,
    *,
    failures: list[dict[str, str]],
    stage: str,
) -> bool:
    diagnostics = _ThirdPartyDiagnosticCapture()
    try:
        with diagnostics:
            workbook.save(output_path)
    except Exception as exc:  # pragma: no cover - defensive against third-party failures
        _append_captured_diagnostic_failures(failures, diagnostics, stage=f"{stage}-diagnostic")
        failures.append(
            {
                "stage": stage,
                "detail": f"Could not save workbook `{output_path.name}`: {exc}",
            }
        )
        return False
    _append_captured_diagnostic_failures(failures, diagnostics, stage=f"{stage}-diagnostic")
    return True


def convert_office_to_format(
    source_path: Path,
    output_dir: Path,
    soffice_binary: str,
    *,
    target_format: str,
) -> tuple[Path | None, list[dict[str, str]]]:
    """Convert an Office document to another format via LibreOffice."""
    failures: list[dict[str, str]] = []
    command = [
        soffice_binary,
        "--headless",
        "--norestore",
        "--nolockcheck",
        "--nodefault",
        "--convert-to",
        target_format,
        "--outdir",
        str(output_dir),
        str(source_path),
    ]
    completed = subprocess.run(command, capture_output=True, text=True, check=False)
    if completed.returncode != 0:
        detail = completed.stderr.strip() or completed.stdout.strip() or "no output"
        failures.append({"stage": "convert-office-to-pdf", "detail": detail})
        return None, failures

    extension = target_format.split(":", 1)[0].lower()
    converted = sorted(output_dir.glob(f"*.{extension}"))
    if not converted:
        failures.append(
            {
                "stage": f"convert-office-to-{extension}",
                "detail": f"LibreOffice completed without producing a .{extension} output.",
            }
        )
        return None, failures
    return converted[0], failures


def normalize_legacy_office_source(
    source_path: Path,
    *,
    source_entry: dict[str, Any],
    output_dir: Path,
    soffice_binary: str,
) -> tuple[Path | None, list[dict[str, str]]]:
    """Convert legacy binary Office formats to OOXML so existing builders can reuse them."""
    source_extension = str(
        source_entry.get("source_extension") or source_path.suffix.lower().lstrip(".")
    )
    target_extension = LEGACY_OFFICE_EXTENSION_MAP.get(source_extension)
    if target_extension is None:
        return source_path, []
    return convert_office_to_format(
        source_path,
        output_dir,
        soffice_binary,
        target_format=target_extension,
    )


def render_xlsx_sheet_documents(
    workbook_path: Path,
    *,
    sheet_names: list[str],
    renders_dir: Path,
    tempdir: Path,
    soffice_binary: str,
) -> tuple[dict[str, list[str]], list[str], list[dict[str, str]]]:
    """Render one workbook into per-sheet PNG assets via temporary isolated exports."""
    sheet_renders: dict[str, list[str]] = {}
    document_renders: list[str] = []
    failures: list[dict[str, str]] = []
    if not sheet_names:
        return sheet_renders, document_renders, failures

    for index, sheet_name in enumerate(sheet_names, start=1):
        isolated_workbook_path = tempdir / f"sheet-{index:03d}.xlsx"
        try:
            workbook = _load_workbook_quietly(
                workbook_path,
                data_only=False,
                failures=failures,
                stage="xlsx-sheet-export-load",
            )
            if workbook is None:
                continue
            target_sheet = workbook[sheet_name]
            for worksheet in workbook.worksheets:
                worksheet.sheet_state = "visible" if worksheet.title == sheet_name else "hidden"
            workbook.active = workbook.worksheets.index(target_sheet)
            if not _save_workbook_quietly(
                workbook,
                isolated_workbook_path,
                failures=failures,
                stage="xlsx-sheet-export-save",
            ):
                continue
        except Exception as exc:  # pragma: no cover - defensive against third-party failures
            failures.append(
                {
                    "stage": "xlsx-sheet-export",
                    "detail": f"Could not prepare isolated workbook for `{sheet_name}`: {exc}",
                }
            )
            continue

        converted_pdf, conversion_failures = convert_office_to_pdf(
            isolated_workbook_path,
            tempdir,
            soffice_binary,
        )
        failures.extend(conversion_failures)
        if converted_pdf is None:
            continue
        rendered_assets, render_failures = render_pdf_document(
            converted_pdf,
            renders_dir,
            prefix=f"sheet-{index:03d}-page",
        )
        failures.extend(render_failures)
        sheet_renders[sheet_name] = rendered_assets
        document_renders.extend(rendered_assets)
    return sheet_renders, document_renders, failures


def pptx_texts_from_slide(slide: Any) -> list[str]:
    """Collect visible text content from a PPTX slide."""
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = sanitize_text(shape.text)
            if text:
                texts.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    text = sanitize_text(cell.text)
                    if text:
                        texts.append(text)
    return texts


def pptx_notes_text(slide: Any) -> str:
    """Collect slide notes text when it is available."""
    notes_slide = getattr(slide, "notes_slide", None)
    if notes_slide is None:
        return ""
    texts: list[str] = []
    for shape in notes_slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = sanitize_text(shape.text)
            if text:
                texts.append(text)
    return "\n".join(texts).strip()


def extract_pptx_media_refs(slide: Any) -> list[str]:
    """Collect embedded media references from a slide part."""
    refs: list[str] = []
    for relationship in slide.part.rels.values():
        target = getattr(relationship, "target_ref", "")
        if "/media/" in target:
            refs.append(str(target))
    return sorted(set(refs))


def pptx_slide_hidden(slide: Any) -> bool:
    """Return whether a PPTX slide is marked hidden."""
    return bool(slide._element.get("show") == "0")


def build_pdf_source(
    paths: WorkspacePaths,
    source_path: Path,
    source_entry: dict[str, Any],
    source_dir: Path,
) -> tuple[dict[str, Any], dict[str, Any]]:
    """Build manifests and extracted artifacts for a PDF source."""
    source_entry = normalize_source_entry_for_build(paths, source_path, source_entry)
    extracted_dir = source_dir / "extracted"
    renders_dir = source_dir / "renders"
    ensure_directory(extracted_dir)
    ensure_directory(renders_dir)

    rendered_assets, render_failures = render_pdf_document(source_path, renders_dir)
    page_texts, text_failures = extract_pdf_text(source_path)
    page_count = max(len(rendered_assets), len(page_texts))
    units: list[dict[str, Any]] = []
    for index in range(page_count):
        unit_id = f"page-{index + 1:03d}"
        text = page_texts[index] if index < len(page_texts) else ""
        text_asset = Path("extracted") / f"{unit_id}.txt"
        structure_asset = Path("extracted") / f"{unit_id}.json"
        rendered_asset = rendered_assets[index] if index < len(rendered_assets) else None
        write_text(source_dir / text_asset, text + ("\n" if text else ""))
        write_json(
            source_dir / structure_asset,
            {
                "unit_id": unit_id,
                "ordinal": index + 1,
                "text_excerpt": text,
            },
        )
        units.append(
            {
                "unit_id": unit_id,
                "unit_type": "page",
                "ordinal": index + 1,
                "title": f"Page {index + 1}",
                "rendered_asset": rendered_asset,
                "render_assets": [rendered_asset] if rendered_asset else [],
                "render_page_span": (
                    {"start": index + 1, "end": index + 1} if rendered_asset else None
                ),
                "text_asset": str(text_asset),
                "structure_asset": str(structure_asset),
                "embedded_media": [],
                "extraction_confidence": "high" if text else "low",
                "trust_prior_inputs": build_trust_prior_from_source_entry(source_entry),
            }
        )

    extracted_texts = [text for text in page_texts if text]
    source_manifest = build_source_manifest(paths, source_entry, "python-pdf")
    phase_three = compile_pdf_visual_artifacts(
        source_dir,
        source_id=str(source_entry["source_id"]),
        pdf_path=source_path,
        units=units,
        page_texts=page_texts,
    )
    for unit in units:
        unit_update = phase_three["unit_updates"].get(str(unit.get("unit_id")), {})
        if isinstance(unit_update, dict):
            unit.update(unit_update)
    materialize_focus_render_assets(source_dir, evidence_manifest={"units": units})
    evidence_manifest = {
        "source_id": source_entry["source_id"],
        "document_type": "pdf",
        "source_fingerprint": source_entry["source_fingerprint"],
        "generated_at": utc_now(),
        "rendering": {"renderer": "pypdfium2", "status": "ready" if rendered_assets else "failed"},
        "document_renders": rendered_assets,
        "units": units,
        "failures": render_failures + text_failures,
        "warnings": list(phase_three.get("warnings", [])),
        "structure_assets": [
            str(Path("extracted") / f"page-{index + 1:03d}.json") for index in range(page_count)
        ],
        "artifact_index_asset": "artifact_index.json",
        "pdf_document_asset": phase_three.get("pdf_document_asset"),
        "visual_layout_assets": phase_three["visual_layout_assets"],
        "embedded_media": [],
        "language_candidates": [detect_source_language(extracted_texts)],
    }
    return (
        source_manifest,
        enrich_evidence_manifest_reference_fields(
            source_manifest,
            evidence_manifest,
            source_dir=source_dir,
        ),
    )


def build_pptx_source(
    paths: WorkspacePaths,
    source_path: Path,
    source_entry: dict[str, Any],
    source_dir: Path,
    soffice_binary: str,
) -> tuple[dict[str, Any], dict[str, Any]]:
    """Build manifests and extracted artifacts for a PPTX source."""
    from pptx import Presentation

    source_entry = normalize_source_entry_for_build(paths, source_path, source_entry)
    extracted_dir = source_dir / "extracted"
    renders_dir = source_dir / "renders"
    ensure_directory(extracted_dir)
    ensure_directory(renders_dir)

    source_manifest = build_source_manifest(paths, source_entry, "libreoffice-pdf")
    with tempfile.TemporaryDirectory() as tempdir_name:
        tempdir = Path(tempdir_name)
        normalized_source_path, normalization_failures = normalize_legacy_office_source(
            source_path,
            source_entry=source_entry,
            output_dir=tempdir,
            soffice_binary=soffice_binary,
        )
        presentation = (
            Presentation(str(normalized_source_path))
            if normalized_source_path is not None
            else None
        )
        converted_pdf, conversion_failures = convert_office_to_pdf(
            source_path, tempdir, soffice_binary
        )
        render_failures: list[dict[str, str]] = []
        rendered_assets: list[str] = []
        if converted_pdf is not None:
            rendered_assets, render_failures = render_pdf_document(converted_pdf, renders_dir)

    units: list[dict[str, Any]] = []
    all_texts: list[str] = []
    embedded_media: list[str] = []
    visible_render_index = 0
    if presentation is not None:
        for index, slide in enumerate(presentation.slides, start=1):
            unit_id = f"slide-{index:03d}"
            hidden = pptx_slide_hidden(slide)
            visible_texts = pptx_texts_from_slide(slide)
            slide_title = ""
            try:
                slide_title = sanitize_text(slide.shapes.title.text)
            except Exception:
                slide_title = visible_texts[0] if visible_texts else ""
            notes = pptx_notes_text(slide)
            media_refs = extract_pptx_media_refs(slide)
            all_texts.extend(visible_texts)
            if notes:
                all_texts.append(notes)
            embedded_media.extend(media_refs)
            text_asset = Path("extracted") / f"{unit_id}.txt"
            structure_asset = Path("extracted") / f"{unit_id}.json"
            write_text(
                source_dir / text_asset,
                "\n".join(visible_texts + ([notes] if notes else [])) + "\n",
            )
            write_json(
                source_dir / structure_asset,
                {
                    "unit_id": unit_id,
                    "ordinal": index,
                    "visible_text": visible_texts,
                    "notes_text": notes,
                    "embedded_media": media_refs,
                    "hidden": hidden,
                },
            )
            rendered_asset = None
            if not hidden and visible_render_index < len(rendered_assets):
                rendered_asset = rendered_assets[visible_render_index]
                visible_render_index += 1
            render_ordinal = visible_render_index if rendered_asset else None
            units.append(
                {
                    "unit_id": unit_id,
                    "unit_type": "slide",
                    "ordinal": index,
                    "title": slide_title or f"Slide {index}",
                    "rendered_asset": rendered_asset,
                    "render_assets": [rendered_asset] if rendered_asset else [],
                    "render_ordinal": render_ordinal,
                    "render_page_span": (
                        {"start": render_ordinal, "end": render_ordinal}
                        if isinstance(render_ordinal, int)
                        else None
                    ),
                    "text_asset": str(text_asset),
                    "structure_asset": str(structure_asset),
                    "embedded_media": media_refs,
                    "extraction_confidence": "high" if visible_texts or notes else "low",
                    "hidden": hidden,
                    "heading_aliases": [slide_title] if slide_title else [],
                    "locator_aliases": [slide_title] if slide_title else [f"Slide {index}"],
                    "trust_prior_inputs": build_trust_prior_from_source_entry(source_entry),
                }
            )

    phase_three = (
        compile_pptx_visual_artifacts(
            source_dir,
            source_id=str(source_entry["source_id"]),
            presentation=presentation,
            units=units,
        )
        if presentation is not None
        else {
            "artifact_index": write_empty_artifact_index(
                source_dir, source_id=str(source_entry["source_id"])
            ),
            "visual_layout_assets": [],
            "unit_updates": {},
        }
    )
    for unit in units:
        unit_update = phase_three["unit_updates"].get(str(unit.get("unit_id")), {})
        if isinstance(unit_update, dict):
            unit.update(unit_update)
    materialize_focus_render_assets(source_dir, evidence_manifest={"units": units})
    evidence_manifest = {
        "source_id": source_entry["source_id"],
        "document_type": "pptx",
        "source_fingerprint": source_entry["source_fingerprint"],
        "generated_at": utc_now(),
        "rendering": {
            "renderer": "libreoffice->pdf->png",
            "status": "ready" if rendered_assets else "failed",
        },
        "document_renders": rendered_assets,
        "units": units,
        "failures": normalization_failures + conversion_failures + render_failures,
        "structure_assets": [
            str(Path("extracted") / f"slide-{index:03d}.json") for index in range(1, len(units) + 1)
        ],
        "artifact_index_asset": "artifact_index.json",
        "visual_layout_assets": phase_three["visual_layout_assets"],
        "embedded_media": sorted(set(embedded_media)),
        "language_candidates": [detect_source_language(all_texts)],
    }
    return (
        source_manifest,
        enrich_evidence_manifest_reference_fields(
            source_manifest,
            evidence_manifest,
            source_dir=source_dir,
        ),
    )


def extract_docx_blocks(source_path: Path) -> tuple[list[dict[str, Any]], list[str]]:
    """Extract section-order block structure and embedded image references from a DOCX."""
    from docx import Document

    document = Document(str(source_path))
    sections: list[dict[str, Any]] = []
    embedded_images: list[str] = []
    current_blocks: list[dict[str, Any]] = []
    current_section = 1

    def finalize_section() -> None:
        nonlocal current_blocks, current_section
        if current_blocks or not sections:
            sections.append(
                {
                    "unit_id": f"section-{current_section:03d}",
                    "ordinal": current_section,
                    "blocks": current_blocks,
                }
            )
            current_blocks = []
            current_section += 1

    for child in document.element.body.iterchildren():
        tag = child.tag.rsplit("}", 1)[-1]
        if tag == "sectPr":
            finalize_section()
            continue

        texts = [
            sanitize_text(node.text)
            for node in child.iter()
            if node.tag.rsplit("}", 1)[-1] == "t" and sanitize_text(node.text)
        ]
        style_name = next(
            (
                sanitize_text(
                    node.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val")
                )
                for node in child.iter()
                if node.tag.rsplit("}", 1)[-1] == "pStyle"
                and sanitize_text(
                    node.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val")
                )
            ),
            "",
        )
        has_numbering = any(node.tag.rsplit("}", 1)[-1] == "numPr" for node in child.iter())
        image_refs = [
            str(
                node.get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                )
            )
            for node in child.iter()
            if node.tag.rsplit("}", 1)[-1] == "blip"
            and node.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
        ]
        embedded_images.extend(image_refs)
        text_value = "\n".join(texts).strip()
        is_heading = style_name.lower().startswith("heading")
        if is_heading and current_blocks:
            finalize_section()
        block = {
            "kind": "table" if tag == "tbl" else "paragraph",
            "text": text_value,
            "image_refs": image_refs,
            "style_name": style_name or None,
            "is_heading": is_heading,
            "list_kind": (
                "ordered"
                if DOCX_ORDERED_STEP_PATTERN.match(text_value)
                else "bullet"
                if has_numbering or text_value.lstrip().startswith(("-", "*", "•"))
                else None
            ),
            "caption_kind": ("figure" if DOCX_CAPTION_PREFIX_PATTERN.match(text_value) else None),
        }
        if block["text"] or block["image_refs"]:
            current_blocks.append(block)

        has_section_break = any(node.tag.rsplit("}", 1)[-1] == "sectPr" for node in child.iter())
        if has_section_break:
            finalize_section()

    if current_blocks:
        finalize_section()
    for section in sections:
        blocks = [block for block in section.get("blocks", []) if isinstance(block, dict)]
        headings = [
            str(block.get("text"))
            for block in blocks
            if block.get("is_heading") and isinstance(block.get("text"), str) and block.get("text")
        ]
        section["headings"] = headings
        section["procedure_spans"] = [
            {
                "kind": str(block.get("list_kind")),
                "text_excerpt": sanitize_text(block.get("text")),
            }
            for block in blocks
            if isinstance(block.get("list_kind"), str) and sanitize_text(block.get("text"))
        ]
        captions = [
            sanitize_text(block.get("text"))
            for block in blocks
            if block.get("caption_kind") and sanitize_text(block.get("text"))
        ]
        section["captions"] = captions
        image_count = sum(
            len(block.get("image_refs", []))
            for block in blocks
            if isinstance(block.get("image_refs"), list)
        )
        role_hints: list[str] = []
        if image_count:
            role_hints.append("image-heavy")
        if any(block.get("kind") == "table" for block in blocks):
            role_hints.append("table-heavy")
        if section["procedure_spans"]:
            role_hints.append("procedure-like")
        section["role_hints"] = role_hints
    return sections, sorted(set(embedded_images))


def extract_docx_embedded_media(
    source_path: Path,
    *,
    image_refs: list[str],
    source_dir: Path,
) -> dict[str, str]:
    """Extract embedded DOCX images into stable per-source media assets."""
    from docx import Document

    document = Document(str(source_path))
    media_dir = source_dir / "media"
    media_dir.mkdir(parents=True, exist_ok=True)
    assets_by_ref: dict[str, str] = {}
    for image_ref in sorted(
        {
            value
            for value in image_refs
            if isinstance(value, str) and value.strip()
        }
    ):
        relationship = document.part.rels.get(image_ref)
        if relationship is None:
            continue
        if not str(getattr(relationship, "reltype", "")).endswith("/image"):
            continue
        target_part = getattr(relationship, "target_part", None)
        if target_part is None:
            continue
        blob = getattr(target_part, "blob", None)
        partname = Path(str(getattr(target_part, "partname", "") or ""))
        suffix = partname.suffix.lower() or ".bin"
        if not isinstance(blob, (bytes, bytearray)) or not blob:
            continue
        asset_path = media_dir / f"{image_ref}{suffix}"
        asset_path.write_bytes(bytes(blob))
        assets_by_ref[image_ref] = str(asset_path.relative_to(source_dir))
    return assets_by_ref


def extract_xlsx_embedded_media(
    workbook_value: Any,
    *,
    units: list[dict[str, Any]],
    source_dir: Path,
) -> dict[str, dict[str, str]]:
    """Extract embedded XLSX images into stable per-source media assets."""
    media_dir = source_dir / "media"
    media_dir.mkdir(parents=True, exist_ok=True)
    unit_lookup = {
        str(unit.get("title") or ""): str(unit.get("unit_id") or "")
        for unit in units
        if isinstance(unit.get("title"), str) and isinstance(unit.get("unit_id"), str)
    }
    assets_by_unit: dict[str, dict[str, str]] = {}
    for worksheet in getattr(workbook_value, "worksheets", []) or []:
        worksheet_title = str(getattr(worksheet, "title", "") or "")
        unit_id = unit_lookup.get(worksheet_title)
        if not unit_id:
            continue
        unit_assets: dict[str, str] = {}
        for index, image in enumerate(getattr(worksheet, "_images", []), start=1):
            image_ref = f"image-{index:03d}"
            try:
                blob = image._data()
            except Exception:
                continue
            if not isinstance(blob, (bytes, bytearray)) or not blob:
                continue
            extension = str(getattr(image, "format", "") or "").lower().strip(".")
            suffix = f".{extension}" if extension else ".bin"
            asset_path = media_dir / f"{unit_id}-{image_ref}{suffix}"
            asset_path.write_bytes(bytes(blob))
            unit_assets[image_ref] = str(asset_path.relative_to(source_dir))
        if unit_assets:
            assets_by_unit[unit_id] = unit_assets
    return assets_by_unit


def build_docx_source(
    paths: WorkspacePaths,
    source_path: Path,
    source_entry: dict[str, Any],
    source_dir: Path,
    soffice_binary: str,
) -> tuple[dict[str, Any], dict[str, Any]]:
    """Build manifests and extracted artifacts for a DOCX source."""
    source_entry = normalize_source_entry_for_build(paths, source_path, source_entry)
    extracted_dir = source_dir / "extracted"
    renders_dir = source_dir / "renders"
    ensure_directory(extracted_dir)
    ensure_directory(renders_dir)

    source_manifest = build_source_manifest(paths, source_entry, "libreoffice-pdf")
    with tempfile.TemporaryDirectory() as tempdir_name:
        tempdir = Path(tempdir_name)
        normalized_source_path, normalization_failures = normalize_legacy_office_source(
            source_path,
            source_entry=source_entry,
            output_dir=tempdir,
            soffice_binary=soffice_binary,
        )
        sections, embedded_images = (
            extract_docx_blocks(normalized_source_path)
            if normalized_source_path is not None
            else ([], [])
        )
        embedded_media_assets = (
            extract_docx_embedded_media(
                normalized_source_path,
                image_refs=embedded_images,
                source_dir=source_dir,
            )
            if normalized_source_path is not None
            else {}
        )
        converted_pdf, conversion_failures = convert_office_to_pdf(
            source_path, tempdir, soffice_binary
        )
        render_failures: list[dict[str, str]] = []
        rendered_assets: list[str] = []
        if converted_pdf is not None:
            rendered_assets, render_failures = render_pdf_document(converted_pdf, renders_dir)

    structure_assets: list[str] = []
    all_texts: list[str] = []
    units: list[dict[str, Any]] = []
    for section in sections:
        unit_id = str(section["unit_id"])
        structure_asset = Path("extracted") / f"{unit_id}.json"
        text_asset = Path("extracted") / f"{unit_id}.txt"
        section_text = "\n".join(
            block["text"]
            for block in section["blocks"]
            if isinstance(block.get("text"), str) and block["text"]
        ).strip()
        all_texts.append(section_text)
        write_json(source_dir / structure_asset, section)
        write_text(source_dir / text_asset, section_text + ("\n" if section_text else ""))
        structure_assets.append(str(structure_asset))
        headings = [
            value
            for value in section.get("headings", [])
            if isinstance(value, str) and value.strip()
        ]
        locator_aliases = headings[:3] or [f"Section {section['ordinal']}"]
        units.append(
            {
                "unit_id": unit_id,
                "unit_type": "section",
                "ordinal": section["ordinal"],
                "title": headings[0] if headings else f"Section {section['ordinal']}",
                "rendered_asset": None,
                "render_reference_ids": [Path(asset).stem for asset in rendered_assets],
                "text_asset": str(text_asset),
                "structure_asset": str(structure_asset),
                "embedded_media": sorted(
                    {
                        image
                        for block in section["blocks"]
                        for image in block.get("image_refs", [])
                        if isinstance(image, str)
                    }
                ),
                "extraction_confidence": "high" if section_text else "low",
                "heading_aliases": headings,
                "locator_aliases": locator_aliases,
                "trust_prior_inputs": build_trust_prior_from_source_entry(source_entry),
            }
        )

    phase_three = compile_docx_visual_compatibility(
        source_dir,
        source_id=str(source_entry["source_id"]),
        units=units,
        document_renders=rendered_assets,
        embedded_media_assets=embedded_media_assets,
    )
    for unit in units:
        unit_update = phase_three["unit_updates"].get(str(unit.get("unit_id")), {})
        if isinstance(unit_update, dict):
            unit.update(unit_update)
    materialize_focus_render_assets(source_dir, evidence_manifest={"units": units})
    evidence_manifest = {
        "source_id": source_entry["source_id"],
        "document_type": "docx",
        "source_fingerprint": source_entry["source_fingerprint"],
        "generated_at": utc_now(),
        "rendering": {
            "renderer": "libreoffice->pdf->png",
            "status": "ready" if rendered_assets else "failed",
        },
        "document_renders": rendered_assets,
        "units": units,
        "failures": normalization_failures + conversion_failures + render_failures,
        "structure_assets": structure_assets,
        "artifact_index_asset": "artifact_index.json",
        "visual_layout_assets": phase_three["visual_layout_assets"],
        "embedded_media": embedded_images,
        "language_candidates": [detect_source_language(all_texts)],
    }
    return (
        source_manifest,
        enrich_evidence_manifest_reference_fields(
            source_manifest,
            evidence_manifest,
            source_dir=source_dir,
        ),
    )


def build_xlsx_source(
    paths: WorkspacePaths,
    source_path: Path,
    source_entry: dict[str, Any],
    source_dir: Path,
    soffice_binary: str,
) -> tuple[dict[str, Any], dict[str, Any]]:
    """Build manifests and extracted artifacts for an XLSX source."""
    source_entry = normalize_source_entry_for_build(paths, source_path, source_entry)
    extracted_dir = source_dir / "extracted"
    renders_dir = source_dir / "renders"
    ensure_directory(extracted_dir)
    ensure_directory(renders_dir)

    source_manifest = build_source_manifest(paths, source_entry, "libreoffice-pdf")
    with tempfile.TemporaryDirectory() as tempdir_name:
        tempdir = Path(tempdir_name)
        render_failures: list[dict[str, str]] = []
        conversion_failures: list[dict[str, str]] = []
        rendered_assets: list[str] = []
        sheet_render_assets: dict[str, list[str]] = {}
        normalized_source_path, normalization_failures = normalize_legacy_office_source(
            source_path,
            source_entry=source_entry,
            output_dir=tempdir,
            soffice_binary=soffice_binary,
        )
        workbook_formula = (
            _load_workbook_quietly(
                normalized_source_path,
                data_only=False,
                failures=render_failures,
                stage="xlsx-formula-load",
            )
            if normalized_source_path is not None
            else None
        )
        workbook_value = (
            _load_workbook_quietly(
                normalized_source_path,
                data_only=True,
                failures=render_failures,
                stage="xlsx-value-load",
            )
            if normalized_source_path is not None
            else None
        )
        if normalized_source_path is not None and workbook_formula is not None:
            (
                sheet_render_assets,
                rendered_assets,
                sheet_render_failures,
            ) = render_xlsx_sheet_documents(
                normalized_source_path,
                sheet_names=[worksheet.title for worksheet in workbook_formula.worksheets],
                renders_dir=renders_dir,
                tempdir=tempdir,
                soffice_binary=soffice_binary,
            )
            render_failures.extend(sheet_render_failures)
        if not rendered_assets:
            converted_pdf, conversion_failures = convert_office_to_pdf(
                source_path,
                tempdir,
                soffice_binary,
            )
            if converted_pdf is not None:
                rendered_assets, fallback_render_failures = render_pdf_document(
                    converted_pdf,
                    renders_dir,
                )
                render_failures.extend(fallback_render_failures)

    units: list[dict[str, Any]] = []
    structure_assets: list[str] = []
    all_texts: list[str] = []
    for index, worksheet in enumerate(
        workbook_value.worksheets if workbook_value is not None else [],
        start=1,
    ):
        non_empty_cells: list[dict[str, str]] = []
        truncated = False
        for row in worksheet.iter_rows():
            for cell in row:
                if cell.value in (None, ""):
                    continue
                if len(non_empty_cells) >= 2000:
                    truncated = True
                    break
                non_empty_cells.append(
                    {"cell": cell.coordinate, "value": sanitize_text(cell.value)}
                )
            if truncated:
                break

        tables = []
        for name in worksheet.tables.keys():
            table = worksheet.tables[name]
            table_ref = table if isinstance(table, str) else getattr(table, "ref", None)
            tables.append({"name": name, "ref": table_ref})
        text_asset = Path("extracted") / f"sheet-{index:03d}.txt"
        structure_asset = Path("extracted") / f"sheet-{index:03d}.json"
        lines = [f"{cell['cell']}: {cell['value']}" for cell in non_empty_cells if cell["value"]]
        sheet_text = "\n".join(lines).strip()
        all_texts.append(sheet_text)
        write_text(source_dir / text_asset, sheet_text + ("\n" if sheet_text else ""))
        write_json(
            source_dir / structure_asset,
            {
                "unit_id": f"sheet-{index:03d}",
                "ordinal": index,
                "sheet_name": worksheet.title,
                "max_row": worksheet.max_row,
                "max_column": worksheet.max_column,
                "non_empty_cells": non_empty_cells,
                "truncated": truncated,
                "tables": tables,
                "image_count": len(getattr(worksheet, "_images", [])),
                "has_drawing": bool(getattr(worksheet, "_drawing", None)),
            },
        )
        structure_assets.append(str(structure_asset))
        unit_render_assets = list(sheet_render_assets.get(worksheet.title, []))
        units.append(
            {
                "unit_id": f"sheet-{index:03d}",
                "unit_type": "sheet",
                "ordinal": index,
                "title": worksheet.title,
                "rendered_asset": unit_render_assets[0] if unit_render_assets else None,
                "render_assets": unit_render_assets,
                "render_reference_ids": [Path(asset).stem for asset in unit_render_assets],
                "render_page_span": None,
                "text_asset": str(text_asset),
                "structure_asset": str(structure_asset),
                "embedded_media": [],
                "extraction_confidence": "high" if non_empty_cells else "low",
                "hidden": worksheet.sheet_state != "visible",
                "trust_prior_inputs": build_trust_prior_from_source_entry(source_entry),
            }
        )

    embedded_media_assets = (
        extract_xlsx_embedded_media(
            workbook_value,
            units=units,
            source_dir=source_dir,
        )
        if workbook_value is not None
        else {}
    )
    for unit in units:
        unit_id = str(unit.get("unit_id") or "")
        if not unit_id:
            continue
        unit["embedded_media"] = sorted(
            embedded_media_assets.get(unit_id, {}).values()
        )

    phase_three = (
        compile_xlsx_artifacts(
            source_dir,
            source_id=str(source_entry["source_id"]),
            workbook_formula=workbook_formula,
            workbook_value=workbook_value,
            units=units,
            sheet_render_assets=sheet_render_assets,
            embedded_media_assets=embedded_media_assets,
        )
        if workbook_formula is not None and workbook_value is not None
        else {
            "artifact_index": write_empty_artifact_index(
                source_dir,
                source_id=str(source_entry["source_id"]),
            ),
            "spreadsheet_workbook_asset": None,
            "spreadsheet_sheet_assets": [],
            "visual_layout_assets": [],
            "unit_updates": {},
        }
    )
    for unit in units:
        unit_update = phase_three["unit_updates"].get(str(unit.get("unit_id")), {})
        if isinstance(unit_update, dict):
            unit.update(unit_update)
    materialize_focus_render_assets(source_dir, evidence_manifest={"units": units})
    evidence_manifest = {
        "source_id": source_entry["source_id"],
        "document_type": "xlsx",
        "source_fingerprint": source_entry["source_fingerprint"],
        "generated_at": utc_now(),
        "rendering": {
            "renderer": "libreoffice->pdf->png",
            "status": "ready" if rendered_assets else "failed",
        },
        "document_renders": rendered_assets,
        "units": units,
        "failures": normalization_failures + conversion_failures + render_failures,
        "structure_assets": structure_assets,
        "artifact_index_asset": "artifact_index.json",
        "spreadsheet_workbook_asset": phase_three["spreadsheet_workbook_asset"],
        "spreadsheet_sheet_assets": phase_three["spreadsheet_sheet_assets"],
        "visual_layout_assets": phase_three["visual_layout_assets"],
        "embedded_media": [],
        "language_candidates": [detect_source_language(all_texts)],
    }
    return (
        source_manifest,
        enrich_evidence_manifest_reference_fields(
            source_manifest,
            evidence_manifest,
            source_dir=source_dir,
        ),
    )


def _safe_media_filename(index: int, original_name: str) -> str:
    sanitized = re.sub(r"[^A-Za-z0-9._-]+", "-", original_name).strip("-")
    if not sanitized:
        sanitized = f"asset-{index:03d}"
    return f"{index:03d}-{sanitized}"


def _resolve_text_reference_target(
    paths: WorkspacePaths,
    source_path: Path,
    target: str,
) -> dict[str, Any]:
    stripped = str(target).strip()
    if not stripped:
        return {"target": stripped, "kind": "missing"}
    if stripped.startswith("#"):
        return {"target": stripped, "kind": "anchor", "fragment": stripped.removeprefix("#")}
    parsed = urlparse(stripped)
    if parsed.scheme and parsed.scheme not in {"file"}:
        return {"target": stripped, "kind": "external", "scheme": parsed.scheme}
    target_path = stripped
    fragment = None
    if "#" in target_path:
        target_path, fragment = target_path.split("#", 1)
    candidate = (source_path.parent / target_path).resolve()
    resolved_root = paths.root.resolve()
    resolved_source_dir = paths.source_dir.resolve()
    try:
        relative_to_root = candidate.relative_to(resolved_root)
    except ValueError:
        return {
            "target": stripped,
            "kind": "outside-workspace",
            "fragment": fragment,
        }
    try:
        candidate.relative_to(resolved_source_dir)
    except ValueError:
        return {
            "target": stripped,
            "kind": "outside-source-dir",
            "relative_path": str(relative_to_root),
            "fragment": fragment,
        }
    return {
        "target": stripped,
        "kind": "local",
        "absolute_path": candidate,
        "relative_path": str(relative_to_root),
        "fragment": fragment,
        "exists": candidate.exists(),
    }


def _copy_local_text_media_assets(
    paths: WorkspacePaths,
    source_path: Path,
    source_dir: Path,
    media_refs: list[dict[str, Any]],
) -> tuple[dict[str, dict[str, Any]], list[str]]:
    warnings: list[str] = []
    copied: dict[str, dict[str, Any]] = {}
    media_dir = source_dir / "media"
    created_media_dir = False
    for index, ref in enumerate(media_refs, start=1):
        raw_target = str(ref.get("target") or "").strip()
        if not raw_target or raw_target in copied:
            continue
        resolution = _resolve_text_reference_target(paths, source_path, raw_target)
        kind = str(resolution.get("kind") or "missing")
        if kind != "local":
            if kind == "external":
                warnings.append(
                    f"Preserved external image reference `{raw_target}` without copying."
                )
            elif kind in {"outside-workspace", "outside-source-dir"}:
                warnings.append(
                    "Skipped image reference "
                    f"`{raw_target}` because it resolves outside `original_doc/`."
                )
            else:
                warnings.append(f"Skipped unresolved image reference `{raw_target}`.")
            copied[raw_target] = {
                "target": raw_target,
                "status": "unresolved",
                "resolution": resolution,
            }
            continue
        resolved_path = resolution.get("absolute_path")
        if not isinstance(resolved_path, Path) or not resolved_path.exists():
            warnings.append(
                f"Image reference `{raw_target}` could not be resolved to a local file."
            )
            copied[raw_target] = {
                "target": raw_target,
                "status": "missing",
                "resolution": resolution,
            }
            continue
        if not created_media_dir:
            ensure_directory(media_dir)
            created_media_dir = True
        published_name = _safe_media_filename(index, resolved_path.name)
        published_path = media_dir / published_name
        shutil.copy2(resolved_path, published_path)
        copied[raw_target] = {
            "target": raw_target,
            "status": "copied",
            "resolution": resolution,
            "published_asset": str(Path("media") / published_name),
            "source_fingerprint": file_sha256(resolved_path),
        }
    return copied, warnings


def _enrich_text_structure_refs(
    paths: WorkspacePaths,
    source_path: Path,
    structure_data: dict[str, Any],
    media_lookup: dict[str, dict[str, Any]],
) -> dict[str, Any]:
    enriched = dict(structure_data)
    enriched_blocks: list[dict[str, Any]] = []
    for raw_block in structure_data.get("blocks", []):
        if not isinstance(raw_block, dict):
            continue
        block = dict(raw_block)
        enriched_links: list[dict[str, Any]] = []
        for raw_link in raw_block.get("links", []):
            if not isinstance(raw_link, dict):
                continue
            link = dict(raw_link)
            resolution = _resolve_text_reference_target(
                paths,
                source_path,
                str(raw_link.get("target") or ""),
            )
            link["resolution"] = {
                key: value for key, value in resolution.items() if key != "absolute_path"
            }
            enriched_links.append(link)
        block["links"] = enriched_links

        enriched_images: list[dict[str, Any]] = []
        for raw_image in raw_block.get("images", []):
            if not isinstance(raw_image, dict):
                continue
            image = dict(raw_image)
            media = media_lookup.get(str(raw_image.get("target") or "").strip())
            if isinstance(media, dict):
                image["resolution"] = {
                    key: value
                    for key, value in media.get("resolution", {}).items()
                    if key != "absolute_path"
                }
                if isinstance(media.get("published_asset"), str):
                    image["published_asset"] = media["published_asset"]
                image["status"] = media.get("status")
            enriched_images.append(image)
        block["images"] = enriched_images
        enriched_blocks.append(block)
    enriched["blocks"] = enriched_blocks
    return enriched


def _text_unit_line_span(unit: ParsedUnit) -> tuple[int | None, int | None]:
    line_start = unit.structure_data.get("line_start")
    line_end = unit.structure_data.get("line_end")
    return (
        line_start if isinstance(line_start, int) else None,
        line_end if isinstance(line_end, int) else None,
    )


def build_text_source(
    paths: WorkspacePaths,
    source_path: Path,
    source_entry: dict[str, Any],
    source_dir: Path,
) -> tuple[dict[str, Any], dict[str, Any]]:
    """Build manifests and extracted artifacts for a conservative text-like source."""
    source_entry = normalize_source_entry_for_build(paths, source_path, source_entry)
    definition = source_definition_for_entry(source_entry, source_path)
    if definition is None:
        source_extension = source_entry.get("source_extension") or source_path.suffix
        raise RuntimeError(
            f"Unsupported text source type for staging: {source_extension}"
        )
    extracted_dir = source_dir / "extracted"
    ensure_directory(extracted_dir)

    parsed = parse_text_source(source_path, document_type=definition.document_type)
    media_lookup, media_warnings = _copy_local_text_media_assets(
        paths,
        source_path,
        source_dir,
        parsed.document_media,
    )
    units: list[dict[str, Any]] = []
    structure_assets: list[str] = []
    embedded_media_assets: list[str] = []
    for unit in parsed.units:
        text_asset = Path("extracted") / f"{unit.unit_id}.txt"
        structure_asset = Path("extracted") / f"{unit.unit_id}.json"
        enriched_structure = _enrich_text_structure_refs(
            paths,
            source_path,
            unit.structure_data,
            media_lookup,
        )
        line_start, line_end = _text_unit_line_span(unit)
        unit_media = sorted(
            {
                str(image["published_asset"])
                for block in enriched_structure.get("blocks", [])
                if isinstance(block, dict)
                for image in block.get("images", [])
                if isinstance(image, dict) and isinstance(image.get("published_asset"), str)
            }
        )
        write_text(source_dir / text_asset, unit.text + ("\n" if unit.text else ""))
        write_json(source_dir / structure_asset, enriched_structure)
        structure_assets.append(str(structure_asset))
        embedded_media_assets.extend(unit_media)
        units.append(
            {
                "unit_id": unit.unit_id,
                "unit_type": unit.unit_type,
                "ordinal": unit.ordinal,
                "title": unit.title,
                "rendered_asset": None,
                "text_asset": str(text_asset),
                "structure_asset": str(structure_asset),
                "embedded_media": unit_media,
                "extraction_confidence": unit.extraction_confidence,
                "line_start": line_start,
                "line_end": line_end,
                "warnings": unit.warnings,
                "trust_prior_inputs": build_trust_prior_from_source_entry(source_entry),
            }
        )

    source_manifest = build_source_manifest(
        paths, source_entry, "text-native", title=parsed.source_title
    )
    write_empty_artifact_index(source_dir, source_id=str(source_entry["source_id"]))
    evidence_manifest = {
        "source_id": source_entry["source_id"],
        "document_type": source_entry["document_type"],
        "source_fingerprint": source_entry["source_fingerprint"],
        "generated_at": utc_now(),
        "rendering": {"renderer": "text-native", "status": "ready"},
        "document_renders": [],
        "units": units,
        "failures": parsed.failures,
        "warnings": list(dict.fromkeys([*parsed.warnings, *media_warnings])),
        "structure_assets": structure_assets,
        "artifact_index_asset": "artifact_index.json",
        "embedded_media": sorted(set(embedded_media_assets)),
        "language_candidates": [parsed.source_language],
    }
    return (
        source_manifest,
        enrich_evidence_manifest_reference_fields(
            source_manifest,
            evidence_manifest,
            source_dir=source_dir,
        ),
    )


def _synthetic_source_entry_from_manifest(
    source_manifest: dict[str, Any],
    *,
    current_path: str | None = None,
    change_classification: str = "unchanged",
    change_traits: list[str] | None = None,
    change_reason: str | None = None,
) -> dict[str, Any]:
    """Rehydrate a reusable source entry from an existing source manifest."""
    path_value = str(current_path or source_manifest.get("current_path") or "")
    source_fingerprint = str(source_manifest.get("source_fingerprint") or "")
    previous_path = source_manifest.get("current_path")
    traits = [
        *(change_traits or []),
    ]
    if (
        isinstance(previous_path, str)
        and previous_path
        and previous_path != path_value
        and "path_changed" not in traits
    ):
        traits.append("path_changed")
    return {
        "source_id": source_manifest["source_id"],
        "current_path": path_value,
        "prior_paths": list(source_manifest.get("prior_paths", [])),
        "path_history": append_unique_strings(
            [
                *[
                    value
                    for value in source_manifest.get("path_history", [])
                    if isinstance(value, str) and value
                ],
                path_value,
            ]
        ),
        "source_fingerprint": source_fingerprint,
        "file_size": source_manifest.get("file_size"),
        "modified_at": source_manifest.get("modified_at"),
        "first_seen_at": source_manifest.get("first_seen_at", utc_now()),
        "last_seen_at": utc_now(),
        "identity_confidence": str(
            source_manifest.get("identity_confidence") or "derived-attachment"
        ),
        "identity_basis": str(source_manifest.get("identity_basis") or "derived-attachment"),
        "change_classification": change_classification,
        "change_traits": traits,
        "change_reason": change_reason
        or (
            "The derived attachment source fingerprint is unchanged, so the previous staged "
            "artifacts were reused."
        ),
        "ambiguous_match": False,
        "matched_source_ids": [],
        "document_type": source_manifest["document_type"],
        "source_extension": source_manifest.get("source_extension"),
        "support_tier": source_manifest.get("support_tier"),
        "source_origin": source_manifest.get("source_origin", DERIVED_SOURCE_ORIGIN),
        "parent_source_id": source_manifest.get("parent_source_id"),
        "root_email_source_id": source_manifest.get("root_email_source_id"),
        "attachment_lineage": source_manifest.get("attachment_lineage", []),
        "attachment_depth": source_manifest.get("attachment_depth"),
        "attachment_filename": source_manifest.get("attachment_filename"),
        "attachment_mime_type": source_manifest.get("attachment_mime_type"),
        "attachment_disposition": source_manifest.get("attachment_disposition"),
        "email_metadata": source_manifest.get("email_metadata", {}),
        "child_source_ids": source_manifest.get("child_source_ids", []),
        "published_attachment_assets": source_manifest.get("published_attachment_assets", []),
        "email_subject": source_manifest.get("email_subject"),
        "message_id": source_manifest.get("message_id"),
    }


def _derived_source_change_state(
    previous_manifest: dict[str, Any] | None,
    *,
    current_path: str,
    fingerprint: str,
) -> dict[str, Any]:
    """Compute additive change metadata for one derived attachment source."""
    now = utc_now()
    previous_path = (
        str(previous_manifest.get("current_path"))
        if isinstance(previous_manifest, dict)
        and isinstance(previous_manifest.get("current_path"), str)
        else None
    )
    previous_fingerprint = (
        str(previous_manifest.get("source_fingerprint"))
        if isinstance(previous_manifest, dict)
        and isinstance(previous_manifest.get("source_fingerprint"), str)
        else None
    )
    prior_paths: list[str] = []
    if isinstance(previous_manifest, dict) and isinstance(
        previous_manifest.get("prior_paths"), list
    ):
        prior_paths.extend(
            value for value in previous_manifest["prior_paths"] if isinstance(value, str) and value
        )
    if previous_path and previous_path != current_path:
        prior_paths.append(previous_path)
    path_history = append_unique_strings(
        [
            *(
                [
                    value
                    for value in (
                        previous_manifest.get("path_history", [])
                        if isinstance(previous_manifest, dict)
                        else []
                    )
                    if isinstance(previous_manifest, dict)
                    and isinstance(previous_manifest.get("path_history"), list)
                    and isinstance(value, str)
                ]
            ),
            *prior_paths,
            current_path,
        ]
    )
    if not isinstance(previous_manifest, dict) or not previous_manifest:
        return {
            "first_seen_at": now,
            "last_seen_at": now,
            "prior_paths": append_unique_strings(prior_paths),
            "path_history": path_history or [current_path],
            "change_classification": "added",
            "change_traits": [],
            "change_reason": "This derived attachment source is new to the published corpus.",
        }
    traits: list[str] = []
    if previous_path and previous_path != current_path:
        traits.append("path_changed")
    if previous_fingerprint and previous_fingerprint != fingerprint:
        traits.append("binary_changed")
    if previous_fingerprint == fingerprint:
        classification = "moved-or-renamed" if previous_path != current_path else "unchanged"
    else:
        classification = "modified"
    reason = "The derived attachment source changed and must be restaged."
    if classification == "unchanged":
        reason = (
            "The derived attachment fingerprint is unchanged, so the staged "
            "artifacts can be reused."
        )
    elif classification == "moved-or-renamed":
        reason = (
            "The derived attachment path changed but the attachment fingerprint is unchanged, "
            "so staged artifacts can be reused after metadata refresh."
        )
    return {
        "first_seen_at": previous_manifest.get("first_seen_at", now),
        "last_seen_at": now,
        "prior_paths": append_unique_strings(prior_paths),
        "path_history": path_history or [current_path],
        "change_classification": classification,
        "change_traits": traits,
        "change_reason": reason,
    }


def _copy_reused_source_tree(
    paths: WorkspacePaths,
    *,
    source_entry: dict[str, Any],
    destination_root: Path,
    seen_source_ids: set[str] | None = None,
) -> None:
    """Copy one reused source directory plus any reusable child attachment sources."""
    seen = seen_source_ids if seen_source_ids is not None else set()
    source_id = str(source_entry["source_id"])
    if source_id in seen:
        return
    previous_source_dir = locate_previous_source_dir(paths, source_id)
    if previous_source_dir is None:
        raise RuntimeError(f"Could not locate previous staged artifacts for `{source_id}`.")
    seen.add(source_id)
    destination_dir = destination_root / source_id
    if destination_dir.exists():
        shutil.rmtree(destination_dir)
    shutil.copytree(previous_source_dir, destination_dir)
    refresh_reused_source_metadata(paths, source_entry, destination_dir)
    previous_manifest = read_json(previous_source_dir / "source_manifest.json")
    for child_source_id in previous_manifest.get("child_source_ids", []):
        if not isinstance(child_source_id, str) or not child_source_id:
            continue
        child_previous_dir = locate_previous_source_dir(paths, child_source_id)
        if child_previous_dir is None:
            continue
        child_manifest = read_json(child_previous_dir / "source_manifest.json")
        if not child_manifest:
            continue
        _copy_reused_source_tree(
            paths,
            source_entry=_synthetic_source_entry_from_manifest(child_manifest),
            destination_root=destination_root,
            seen_source_ids=seen,
        )


def _attachment_segment(ordinal: int, filename: str) -> str:
    """Return a deterministic human-readable segment for one attachment path step."""
    return _safe_media_filename(ordinal, filename)


def _attachment_source_id(root_email_source_id: str, lineage_slot: str) -> str:
    """Return the stable child source ID for one attachment lineage."""
    return str(uuid.uuid5(uuid.NAMESPACE_URL, f"{root_email_source_id}:{lineage_slot}"))


def _attachment_current_path(parent_current_path: str, lineage_segments: list[str]) -> str:
    """Return the synthetic current path for a derived attachment source."""
    current = parent_current_path
    for segment in lineage_segments:
        current = f"{current}#attachment/{segment}"
    return current


def _copy_email_attachment_asset(
    source_dir: Path,
    *,
    attachment_ordinal: int,
    filename: str,
    payload_bytes: bytes,
    subdirectory: str,
) -> str:
    """Copy one preserved attachment or media payload into the published source dir."""
    relative = Path(subdirectory) / _safe_media_filename(attachment_ordinal, filename)
    write_bytes(source_dir / relative, payload_bytes)
    return str(relative)


def _persist_source_artifacts(
    source_dir: Path,
    *,
    source_manifest: dict[str, Any],
    evidence_manifest: dict[str, Any],
) -> None:
    """Write the standard manifest and guidance files for one staged source dir."""
    evidence_manifest = sync_optional_sidecar_assets(
        source_dir,
        evidence_manifest=evidence_manifest,
    )
    write_json(source_dir / "source_manifest.json", source_manifest)
    write_json(source_dir / "evidence_manifest.json", evidence_manifest)
    create_source_authoring_notes(source_dir, source_manifest)
    knowledge = read_json(source_dir / "knowledge.json")
    summary_path = source_dir / "summary.md"
    summary_text = summary_path.read_text(encoding="utf-8") if summary_path.exists() else ""
    _write_source_affordances(
        source_dir,
        source_manifest=source_manifest,
        evidence_manifest=evidence_manifest,
        knowledge=knowledge or None,
        summary_text=summary_text,
    )


def sync_optional_sidecar_assets(
    source_dir: Path,
    *,
    evidence_manifest: dict[str, Any],
) -> dict[str, Any]:
    """Refresh optional sidecar asset references from the current staged source dir."""
    payload = dict(evidence_manifest)
    semantic_overlay_assets = collect_semantic_overlay_assets(source_dir)
    if semantic_overlay_assets:
        payload["semantic_overlay_assets"] = semantic_overlay_assets
    else:
        payload.pop("semantic_overlay_assets", None)
    pdf_document_path = source_dir / "pdf_document.json"
    if str(payload.get("document_type") or "") == "pdf" and pdf_document_path.exists():
        payload["pdf_document_asset"] = "pdf_document.json"
    elif payload.get("pdf_document_asset") and not pdf_document_path.exists():
        payload.pop("pdf_document_asset", None)
    return payload


def build_email_source_tree(
    paths: WorkspacePaths,
    *,
    source_entry: dict[str, Any],
    source_dir: Path,
    staging_sources_dir: Path,
    office_binary: str | None,
    attachment_depth: int = 0,
    lineage_slots: tuple[str, ...] = (),
    lineage_segments: tuple[str, ...] = (),
) -> tuple[dict[str, Any], dict[str, Any]]:
    """Build one email source plus any supported derived attachment sources."""
    previous_source_dir = locate_previous_source_dir(paths, str(source_entry["source_id"]))
    previous_manifest = (
        read_json(previous_source_dir / "source_manifest.json")
        if previous_source_dir is not None
        else {}
    )
    if (
        previous_source_dir is not None
        and source_entry.get("change_classification") == "unchanged"
        and previous_manifest.get("source_fingerprint") == source_entry.get("source_fingerprint")
    ):
        _copy_reused_source_tree(
            paths,
            source_entry=source_entry,
            destination_root=staging_sources_dir,
        )
        copied_manifest = read_json(source_dir / "source_manifest.json")
        copied_evidence = read_json(source_dir / "evidence_manifest.json")
        return copied_manifest, copied_evidence

    if source_dir.exists():
        shutil.rmtree(source_dir)
    ensure_directory(source_dir / "extracted")
    ensure_directory(source_dir / "media")
    ensure_directory(source_dir / "attachments")

    source_path = source_content_path(paths, source_entry)
    parsed = parse_email_source(source_path)
    child_source_ids: list[str] = []
    published_attachment_assets: list[str] = []
    deterministic_links: list[dict[str, Any]] = []
    attachment_records: list[dict[str, Any]] = []
    attachment_resolution: dict[str, dict[str, Any]] = {}
    embedded_media_assets: list[str] = []

    with tempfile.TemporaryDirectory() as tempdir_name:
        tempdir = Path(tempdir_name)
        for attachment in parsed.attachments:
            unit_warnings = list(attachment.warnings)
            published_asset: str | None = None
            child_source_id: str | None = None
            child_document_type = attachment.document_type
            child_depth = attachment_depth + 1

            if attachment.content_type.startswith("image/") and attachment.payload_bytes:
                published_asset = _copy_email_attachment_asset(
                    source_dir,
                    attachment_ordinal=attachment.ordinal,
                    filename=attachment.filename,
                    payload_bytes=attachment.payload_bytes,
                    subdirectory="media",
                )
                embedded_media_assets.append(published_asset)
                published_attachment_assets.append(published_asset)

            supported_child = child_document_type is not None and child_document_type != "email"
            if child_document_type == "email":
                supported_child = child_depth <= EMAIL_MAX_ATTACHMENT_DEPTH
                if not supported_child:
                    unit_warnings.append(
                        "Nested `.eml` depth exceeded the supported maximum of "
                        f"{EMAIL_MAX_ATTACHMENT_DEPTH}."
                    )
            if supported_child:
                combined_slots = (*lineage_slots, attachment.lineage_slot)
                combined_segments = (*lineage_segments, *attachment.lineage_segments)
                lineage_slot = "/".join(combined_slots)
                child_source_id = _attachment_source_id(
                    str(source_entry.get("root_email_source_id") or source_entry["source_id"]),
                    lineage_slot,
                )
                child_current_path = _attachment_current_path(
                    str(source_entry["current_path"]),
                    list(attachment.lineage_segments),
                )
                child_previous_dir = locate_previous_source_dir(paths, child_source_id)
                child_previous_manifest = (
                    read_json(child_previous_dir / "source_manifest.json")
                    if child_previous_dir is not None
                    else {}
                )
                attachment_path = tempdir / _safe_media_filename(
                    attachment.ordinal,
                    attachment.filename,
                )
                write_bytes(attachment_path, attachment.payload_bytes)
                source_state = _derived_source_change_state(
                    child_previous_manifest,
                    current_path=child_current_path,
                    fingerprint=file_sha256(attachment_path),
                )
                definition = source_type_definition(attachment.source_extension)
                child_entry = {
                    "source_id": child_source_id,
                    "current_path": child_current_path,
                    "prior_paths": source_state["prior_paths"],
                    "path_history": source_state["path_history"],
                    "source_fingerprint": file_sha256(attachment_path),
                    "file_size": len(attachment.payload_bytes),
                    "modified_at": source_entry.get("modified_at"),
                    "first_seen_at": source_state["first_seen_at"],
                    "last_seen_at": source_state["last_seen_at"],
                    "identity_confidence": "derived-attachment",
                    "identity_basis": "derived-attachment",
                    "change_classification": source_state["change_classification"],
                    "change_traits": source_state["change_traits"],
                    "change_reason": source_state["change_reason"],
                    "ambiguous_match": False,
                    "matched_source_ids": [],
                    "document_type": child_document_type,
                    "source_extension": attachment.source_extension,
                    "support_tier": (
                        definition.support_tier
                        if definition is not None
                        else attachment.support_tier
                    ),
                    "source_origin": DERIVED_SOURCE_ORIGIN,
                    "parent_source_id": source_entry["source_id"],
                    "root_email_source_id": source_entry.get("root_email_source_id")
                    or source_entry["source_id"],
                    "attachment_lineage": list(combined_segments),
                    "attachment_depth": child_depth,
                    "attachment_filename": attachment.filename,
                    "attachment_mime_type": attachment.content_type,
                    "attachment_disposition": attachment.disposition,
                    "email_subject": source_entry.get("email_subject")
                    or parsed.email_metadata.get("subject"),
                    "message_id": source_entry.get("message_id")
                    or parsed.email_metadata.get("message_id"),
                    "content_path": str(attachment_path),
                }
                child_source_dir = staging_sources_dir / child_source_id
                if child_document_type == "email":
                    child_manifest, child_evidence = build_email_source_tree(
                        paths,
                        source_entry=child_entry,
                        source_dir=child_source_dir,
                        staging_sources_dir=staging_sources_dir,
                        office_binary=office_binary,
                        attachment_depth=child_depth,
                        lineage_slots=combined_slots,
                        lineage_segments=combined_segments,
                    )
                else:
                    if (
                        child_previous_dir is not None
                        and child_previous_manifest.get("source_fingerprint")
                        == child_entry["source_fingerprint"]
                    ):
                        child_entry["change_traits"] = append_unique_strings(
                            [
                                *[
                                    str(trait)
                                    for trait in child_entry.get("change_traits", [])
                                    if isinstance(trait, str) and trait
                                ],
                                "source_reused",
                                "semantic_outputs_reused",
                            ]
                        )
                        child_entry["change_reason"] = (
                            "The derived attachment fingerprint is unchanged, so the previous "
                            "staged evidence and semantic outputs were reused."
                        )
                        if child_source_dir.exists():
                            shutil.rmtree(child_source_dir)
                        shutil.copytree(child_previous_dir, child_source_dir)
                        refresh_reused_source_metadata(paths, child_entry, child_source_dir)
                        child_manifest = read_json(child_source_dir / "source_manifest.json")
                        child_evidence = read_json(child_source_dir / "evidence_manifest.json")
                    else:
                        if child_source_dir.exists():
                            shutil.rmtree(child_source_dir)
                        ensure_directory(child_source_dir)
                        child_manifest, child_evidence = build_single_source_artifacts(
                            paths,
                            child_entry,
                            child_source_dir,
                            office_binary,
                        )
                        _persist_source_artifacts(
                            child_source_dir,
                            source_manifest=child_manifest,
                            evidence_manifest=child_evidence,
                        )
                        if child_previous_dir is not None:
                            previous_signature = semantic_evidence_signature(child_previous_dir)
                            current_signature = semantic_evidence_signature(child_source_dir)
                            if (
                                previous_signature is not None
                                and previous_signature == current_signature
                            ):
                                preserve_semantic_outputs(child_previous_dir, child_source_dir)
                                _persist_source_artifacts(
                                    child_source_dir,
                                    source_manifest=child_manifest,
                                    evidence_manifest=child_evidence,
                                )
                child_source_ids.append(child_source_id)
                deterministic_links.append(
                    {
                        "source_id": source_entry["source_id"],
                        "related_source_id": child_source_id,
                        "relation_type": EMAIL_ATTACHMENT_RELATION_TYPE,
                        "strength": "high",
                        "status": "derived",
                        "citation_unit_ids": [attachment.unit_id],
                    }
                )
            elif attachment.payload_bytes:
                if published_asset is None:
                    published_asset = _copy_email_attachment_asset(
                        source_dir,
                        attachment_ordinal=attachment.ordinal,
                        filename=attachment.filename,
                        payload_bytes=attachment.payload_bytes,
                        subdirectory="attachments",
                    )
                    published_attachment_assets.append(published_asset)
                unit_warnings.append(
                    f"Attachment `{attachment.filename}` is preserved as raw "
                    "evidence without a dedicated parser."
                )

            attachment_resolution[attachment.unit_id] = {
                "published_asset": published_asset,
                "child_source_id": child_source_id,
                "warnings": unit_warnings,
            }
            attachment_records.append(
                {
                    "unit_id": attachment.unit_id,
                    "filename": attachment.filename,
                    "mime_type": attachment.content_type,
                    "disposition": attachment.disposition,
                    "size": attachment.size,
                    "inline": attachment.inline,
                    "content_id": attachment.content_id,
                    "support_tier": attachment.support_tier,
                    "document_type": attachment.document_type,
                    "attachment_depth": child_depth,
                    "child_source_id": child_source_id,
                    "published_asset": published_asset,
                    "warnings": unit_warnings,
                }
            )

    structure_assets: list[str] = []
    source_warnings = list(parsed.warnings)
    source_warnings.extend(
        warning
        for attachment in attachment_records
        for warning in attachment.get("warnings", [])
        if isinstance(warning, str)
    )
    cid_asset_map = {
        str(item["content_id"]): item["published_asset"]
        for item in attachment_records
        if isinstance(item.get("content_id"), str) and isinstance(item.get("published_asset"), str)
    }
    if parsed.html_body:
        write_text(source_dir / Path("extracted") / "body.html", parsed.html_body)
    mime_structure = dict(parsed.mime_structure)
    if cid_asset_map:
        mime_structure["cid_asset_map"] = cid_asset_map
    write_json(source_dir / Path("extracted") / "mime-structure.json", mime_structure)

    units: list[dict[str, Any]] = []
    for unit in parsed.units:
        text_asset = Path("extracted") / f"{unit.unit_id}.txt"
        structure_asset = Path("extracted") / f"{unit.unit_id}.json"
        structure_data = dict(unit.structure_data)
        unit_media: list[str] = []
        unit_warnings = list(unit.warnings)
        if unit.unit_type == "email-section":
            if cid_asset_map:
                structure_data["cid_asset_map"] = cid_asset_map
            if parsed.html_body:
                structure_data["html_body_asset"] = "extracted/body.html"
        if unit.unit_type == "email-attachment":
            resolution = attachment_resolution.get(unit.unit_id, {})
            if isinstance(resolution.get("published_asset"), str):
                unit_media.append(str(resolution["published_asset"]))
            structure_data.update(
                {
                    "child_source_id": resolution.get("child_source_id"),
                    "published_asset": resolution.get("published_asset"),
                }
            )
            unit_warnings.extend(
                warning for warning in resolution.get("warnings", []) if isinstance(warning, str)
            )
        write_text(source_dir / text_asset, unit.text + ("\n" if unit.text else ""))
        write_json(source_dir / structure_asset, structure_data)
        structure_assets.append(str(structure_asset))
        line_start, line_end = _text_unit_line_span(unit)
        units.append(
            {
                "unit_id": unit.unit_id,
                "unit_type": unit.unit_type,
                "ordinal": unit.ordinal,
                "title": unit.title,
                "rendered_asset": None,
                "text_asset": str(text_asset),
                "structure_asset": str(structure_asset),
                "embedded_media": unit_media,
                "extraction_confidence": unit.extraction_confidence,
                "line_start": line_start,
                "line_end": line_end,
                "warnings": list(dict.fromkeys(unit_warnings)),
                "trust_prior_inputs": build_trust_prior_from_source_entry(source_entry),
            }
        )

    enriched_entry = {
        **source_entry,
        "email_metadata": parsed.email_metadata,
        "child_source_ids": child_source_ids,
        "published_attachment_assets": published_attachment_assets,
        "email_subject": parsed.email_metadata.get("subject"),
        "message_id": parsed.email_metadata.get("message_id"),
        "root_email_source_id": source_entry.get("root_email_source_id")
        or source_entry["source_id"],
    }
    source_manifest = build_source_manifest(
        paths,
        enriched_entry,
        "email-native",
        title=parsed.source_title,
    )
    write_empty_artifact_index(source_dir, source_id=str(source_entry["source_id"]))
    evidence_manifest = {
        "source_id": source_entry["source_id"],
        "document_type": "email",
        "source_fingerprint": source_entry["source_fingerprint"],
        "generated_at": utc_now(),
        "rendering": {"renderer": "email-native", "status": "ready"},
        "document_renders": [],
        "units": units,
        "failures": parsed.failures,
        "warnings": list(dict.fromkeys(source_warnings)),
        "structure_assets": structure_assets,
        "artifact_index_asset": "artifact_index.json",
        "embedded_media": sorted(set(embedded_media_assets)),
        "language_candidates": [parsed.source_language],
        "email_metadata": parsed.email_metadata,
        "mime_structure_asset": "extracted/mime-structure.json",
        "html_body_asset": "extracted/body.html" if parsed.html_body else None,
        "attachments": attachment_records,
        "deterministic_linked_sources": deterministic_links,
    }
    evidence_manifest = enrich_evidence_manifest_reference_fields(
        source_manifest,
        evidence_manifest,
        source_dir=source_dir,
    )
    _persist_source_artifacts(
        source_dir,
        source_manifest=source_manifest,
        evidence_manifest=evidence_manifest,
    )
    if previous_source_dir is not None:
        previous_signature = semantic_evidence_signature(previous_source_dir)
        current_signature = semantic_evidence_signature(source_dir)
        if previous_signature is not None and previous_signature == current_signature:
            preserve_semantic_outputs(previous_source_dir, source_dir)
            _persist_source_artifacts(
                source_dir,
                source_manifest=source_manifest,
                evidence_manifest=evidence_manifest,
            )
    return source_manifest, evidence_manifest


def build_source_manifest(
    paths: WorkspacePaths,
    source_entry: dict[str, Any],
    render_strategy: str,
    *,
    title: str | None = None,
) -> dict[str, Any]:
    """Build the source manifest payload for a staged source."""
    current_path = str(source_entry["current_path"])
    source_relative_path = source_relative_path_from_current_path(current_path)
    content_path = source_content_path(paths, source_entry)
    stat = content_path.stat() if content_path.exists() else None
    modified_at = str(source_entry.get("modified_at") or "") or (
        isoformat_timestamp(stat.st_mtime) if stat is not None else None
    )
    file_size = source_entry.get("file_size")
    if file_size is None and stat is not None:
        file_size = stat.st_size
    source_manifest = {
        "source_id": source_entry["source_id"],
        "current_path": current_path,
        "prior_paths": source_entry["prior_paths"],
        "path_history": source_entry.get("path_history", [current_path]),
        "relative_path_lineage": build_relative_path_lineage(source_relative_path),
        "document_type": source_entry["document_type"],
        "support_tier": source_entry.get("support_tier"),
        "source_extension": source_entry.get("source_extension")
        or content_path.suffix.lower().lstrip("."),
        "source_fingerprint": source_entry["source_fingerprint"],
        "file_size": int(file_size or 0),
        "modified_at": modified_at,
        "first_seen_at": source_entry["first_seen_at"],
        "last_seen_at": source_entry["last_seen_at"],
        "identity_confidence": source_entry["identity_confidence"],
        "identity_basis": source_entry.get("identity_basis", source_entry["identity_confidence"]),
        "change_classification": source_entry.get("change_classification"),
        "change_traits": source_entry.get("change_traits", []),
        "change_reason": source_entry.get("change_reason"),
        "trust_prior": build_trust_prior_from_source_entry(
            {**source_entry, "current_path": current_path, "modified_at": modified_at}
        ),
        "render_strategy": render_strategy,
        "source_origin": source_entry.get("source_origin", "original-document"),
        "staging_generated_at": utc_now(),
    }
    if isinstance(title, str) and title.strip():
        source_manifest["title"] = title.strip()
    for key in (
        "parent_source_id",
        "root_email_source_id",
        "attachment_lineage",
        "attachment_depth",
        "attachment_filename",
        "attachment_mime_type",
        "attachment_disposition",
        "email_metadata",
        "child_source_ids",
        "published_attachment_assets",
        "email_subject",
        "message_id",
    ):
        value = source_entry.get(key)
        if value not in (None, [], {}):
            source_manifest[key] = value
    return enrich_source_manifest_reference_fields(source_manifest, title=title)


def create_source_authoring_notes(source_dir: Path, source_manifest: dict[str, Any]) -> None:
    """Write a small authoring contract to guide agent-written knowledge objects."""
    write_json(
        source_dir / "work_item.json",
        {
            "source_id": source_manifest["source_id"],
            "knowledge_path": "knowledge.json",
            "summary_path": "summary.md",
            "affordance_path": DEFAULT_AFFORDANCE_FILENAME,
            "required_knowledge_keys": list(REQUIRED_KNOWLEDGE_KEYS),
            "summary_contract": [
                "Use `# <title>` as the first line.",
                "Include `## English Summary` and `## Source-Language Summary` sections.",
                "Mention the source ID in the body.",
            ],
            "affordance_contract": [
                "Preserve KB-native published evidence channels for odd questions.",
                "Keep channel_descriptors compact, evidence-backed, and grouped by channel.",
                "Do not relabel derived descriptors as source-authored facts.",
            ],
        },
    )


def _write_source_affordances(
    source_dir: Path,
    *,
    source_manifest: dict[str, Any],
    evidence_manifest: dict[str, Any],
    knowledge: dict[str, Any] | None = None,
    summary_text: str = "",
) -> dict[str, Any]:
    """Write or refresh the derived-affordance sidecar for one source directory."""
    affordance_path = source_dir / DEFAULT_AFFORDANCE_FILENAME
    baseline = derive_source_affordances(
        source_manifest=source_manifest,
        evidence_manifest=evidence_manifest,
        source_dir=source_dir,
        knowledge=knowledge,
        summary_text=summary_text,
    )
    merged = merge_derived_affordances(baseline, read_json(affordance_path))
    write_json(affordance_path, merged)
    return merged


def build_single_source_artifacts(
    paths: WorkspacePaths,
    source_entry: dict[str, Any],
    source_dir: Path,
    office_binary: str | None,
) -> tuple[dict[str, Any], dict[str, Any]]:
    """Build the source and evidence manifests for one source entry."""
    source_path = source_content_path(paths, source_entry)
    definition = source_definition_for_entry(source_entry, source_path)
    if definition is None:
        raise RuntimeError(
            "Unsupported source type for staging: "
            f"{source_entry.get('source_extension') or source_path.suffix}"
        )
    document_type = definition.document_type
    if document_type == "pdf":
        return build_pdf_source(paths, source_path, source_entry, source_dir)
    if document_type == "pptx":
        if office_binary is None:
            raise RuntimeError("LibreOffice is required to render PowerPoint sources.")
        return build_pptx_source(paths, source_path, source_entry, source_dir, office_binary)
    if document_type == "docx":
        if office_binary is None:
            raise RuntimeError("LibreOffice is required to render Word sources.")
        return build_docx_source(paths, source_path, source_entry, source_dir, office_binary)
    if document_type == "xlsx":
        if office_binary is None:
            raise RuntimeError("LibreOffice is required to render spreadsheet sources.")
        return build_xlsx_source(paths, source_path, source_entry, source_dir, office_binary)
    if document_type in TEXT_DOCUMENT_TYPES:
        return build_text_source(paths, source_path, source_entry, source_dir)
    if document_type == "email":
        raise RuntimeError("Email sources must be staged through the email-specific builder.")
    raise RuntimeError(f"Unsupported source type for staging: {source_path.suffix}")


def locate_previous_source_dir(paths: WorkspacePaths, source_id: str) -> Path | None:
    """Find the best previous source directory for reuse without dropping richer semantic state."""

    def semantic_rank(source_dir: Path) -> tuple[int, int, int, float]:
        overlay_assets = collect_semantic_overlay_assets(source_dir)
        semantic_paths = [
            source_dir / "knowledge.json",
            source_dir / "summary.md",
            source_dir / "source_manifest.json",
            source_dir / "evidence_manifest.json",
            *(source_dir / asset for asset in overlay_assets),
        ]
        newest_mtime = max(
            (path.stat().st_mtime for path in semantic_paths if path.exists()),
            default=0.0,
        )
        return (
            len(overlay_assets),
            int((source_dir / "knowledge.json").exists()),
            int((source_dir / "summary.md").exists()),
            newest_mtime,
        )

    staging_candidate = paths.knowledge_base_staging_dir / "sources" / source_id
    current_candidate = paths.knowledge_base_current_dir / "sources" / source_id
    if staging_candidate.exists() and current_candidate.exists():
        if semantic_rank(current_candidate) >= semantic_rank(staging_candidate):
            return current_candidate
        return staging_candidate
    if current_candidate.exists():
        return current_candidate
    if staging_candidate.exists():
        return staging_candidate
    return None


def preserve_semantic_outputs(previous_source_dir: Path, source_dir: Path) -> None:
    """Preserve authored semantic files when rebuilding a source directory."""
    for filename in ("knowledge.json", "summary.md"):
        previous_path = previous_source_dir / filename
        if previous_path.exists():
            shutil.copy2(previous_path, source_dir / filename)
    previous_overlay_dir = previous_source_dir / "semantic_overlay"
    target_overlay_dir = source_dir / "semantic_overlay"
    if previous_overlay_dir.exists():
        if target_overlay_dir.exists():
            shutil.rmtree(target_overlay_dir)
        shutil.copytree(previous_overlay_dir, target_overlay_dir)


def default_render_strategy(document_type: str) -> str:
    """Return the default render strategy for a document type."""
    if document_type == "pdf":
        return "python-pdf"
    if document_type in {"pptx", "docx", "xlsx"}:
        return "libreoffice-pdf"
    if document_type == "email":
        return "email-native"
    return "text-native"


def refresh_reused_source_metadata(
    paths: WorkspacePaths,
    source_entry: dict[str, Any],
    source_dir: Path,
) -> None:
    """Refresh reused source metadata so copied directories reflect the current sync state."""
    source_manifest_path = source_dir / "source_manifest.json"
    existing_manifest = read_json(source_manifest_path)
    merged_entry = dict(source_entry)
    for key in (
        "source_origin",
        "parent_source_id",
        "root_email_source_id",
        "attachment_lineage",
        "attachment_depth",
        "attachment_filename",
        "attachment_mime_type",
        "attachment_disposition",
        "email_metadata",
        "child_source_ids",
        "published_attachment_assets",
        "email_subject",
        "message_id",
    ):
        if key not in merged_entry and key in existing_manifest:
            merged_entry[key] = existing_manifest.get(key)
    render_strategy = str(
        existing_manifest.get(
            "render_strategy",
            default_render_strategy(str(source_entry["document_type"])),
        )
    )
    refreshed_manifest = build_source_manifest(
        paths,
        merged_entry,
        render_strategy,
        title=str(existing_manifest.get("title") or "").strip() or None,
    )
    write_json(source_manifest_path, refreshed_manifest)
    create_source_authoring_notes(source_dir, refreshed_manifest)
    evidence_manifest = read_json(source_dir / "evidence_manifest.json")
    knowledge = read_json(source_dir / "knowledge.json")
    summary_path = source_dir / "summary.md"
    summary_text = summary_path.read_text(encoding="utf-8") if summary_path.exists() else ""
    if evidence_manifest:
        _write_source_affordances(
            source_dir,
            source_manifest=refreshed_manifest,
            evidence_manifest=evidence_manifest,
            knowledge=knowledge or None,
            summary_text=summary_text,
        )


def refresh_staging_source_metadata(
    paths: WorkspacePaths,
    active_sources: list[dict[str, Any]],
) -> None:
    """Refresh metadata for existing staging source directories without rebuilding evidence."""
    for source_entry in active_sources:
        source_dir = paths.knowledge_base_staging_dir / "sources" / str(source_entry["source_id"])
        if source_dir.exists():
            refresh_reused_source_metadata(paths, source_entry, source_dir)


def source_artifact_contract_complete(
    source_dir: Path,
    *,
    document_type: str,
) -> bool:
    """Return whether a source directory satisfies the current Phase 3 artifact contract."""
    evidence_manifest = read_json(source_dir / "evidence_manifest.json")
    if not evidence_manifest:
        return False
    artifact_index_value = evidence_manifest.get("artifact_index_asset")
    artifact_index_path = (
        source_dir / artifact_index_value
        if isinstance(artifact_index_value, str) and artifact_index_value
        else source_dir / "artifact_index.json"
    )
    if not artifact_index_path.exists():
        return False

    def asset_list_complete(values: Any) -> bool:
        return (
            isinstance(values, list)
            and bool(values)
            and all(
                isinstance(value, str) and value and (source_dir / value).exists()
                for value in values
            )
        )

    if document_type in {"pdf", "pptx", "docx", "xlsx"}:
        if not asset_list_complete(evidence_manifest.get("visual_layout_assets")):
            return False
    if document_type == "pdf":
        pdf_document_asset = evidence_manifest.get("pdf_document_asset")
        if not isinstance(pdf_document_asset, str) or not pdf_document_asset:
            return False
        if not (source_dir / pdf_document_asset).exists():
            return False
    if document_type == "xlsx":
        workbook_asset = evidence_manifest.get("spreadsheet_workbook_asset")
        if not isinstance(workbook_asset, str) or not workbook_asset:
            return False
        if not (source_dir / workbook_asset).exists():
            return False
        if not asset_list_complete(evidence_manifest.get("spreadsheet_sheet_assets")):
            return False
    if document_type in {"pdf", "pptx", "docx", "xlsx"} and not focus_render_contract_complete(
        source_dir
    ):
        return False
    return True


def staging_source_artifacts_complete(
    paths: WorkspacePaths,
    active_sources: list[dict[str, Any]],
) -> bool:
    """Return whether staging still contains the required per-source artifacts."""
    staged_sources_dir = paths.knowledge_base_staging_dir / "sources"
    if not staged_sources_dir.exists():
        return False
    seen_source_ids: set[str] = set()

    def source_dir_complete(source_id: str) -> bool:
        if source_id in seen_source_ids:
            return True
        seen_source_ids.add(source_id)
        source_dir = staged_sources_dir / source_id
        if not source_dir.exists():
            return False
        if not (source_dir / "source_manifest.json").exists():
            return False
        if not (source_dir / "evidence_manifest.json").exists():
            return False
        source_manifest = read_json(source_dir / "source_manifest.json")
        if not source_artifact_contract_complete(
            source_dir,
            document_type=str(source_manifest.get("document_type") or "unknown"),
        ):
            return False
        for child_source_id in source_manifest.get("child_source_ids", []):
            if not isinstance(child_source_id, str) or not child_source_id:
                return False
            if not source_dir_complete(child_source_id):
                return False
        return True

    for source_entry in active_sources:
        source_id = str(source_entry.get("source_id") or "")
        if not source_id or not source_dir_complete(source_id):
            return False
    return True


def write_staging_root_artifacts(
    paths: WorkspacePaths,
    active_sources: list[dict[str, Any]],
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], bool]:
    """Regenerate the root staging artifacts from the current staged source directories."""
    catalog_sources: list[dict[str, Any]] = []
    source_summaries: list[dict[str, Any]] = []
    active_lookup = {str(source["source_id"]): source for source in active_sources}
    ambiguous_match = any(bool(source.get("ambiguous_match")) for source in active_sources)
    staged_sources_dir = paths.knowledge_base_staging_dir / "sources"
    for source_dir in sorted(path for path in staged_sources_dir.iterdir() if path.is_dir()):
        source_id = source_dir.name
        source_manifest = read_json(source_dir / "source_manifest.json")
        evidence_manifest = read_json(source_dir / "evidence_manifest.json")
        source_entry = active_lookup.get(source_id, {})
        classification = str(
            source_entry.get(
                "change_classification",
                source_manifest.get("change_classification", "unknown"),
            )
        )
        catalog_sources.append(
            {
                "source_id": source_manifest["source_id"],
                "current_path": source_manifest["current_path"],
                "document_type": source_manifest["document_type"],
                "support_tier": source_manifest.get("support_tier"),
                "source_origin": source_manifest.get("source_origin", "original-document"),
                "parent_source_id": source_manifest.get("parent_source_id"),
                "root_email_source_id": source_manifest.get("root_email_source_id"),
                "source_dir": str(Path("sources") / source_id),
                "change_classification": classification,
            }
        )
        source_summaries.append(
            {
                "source_id": source_manifest["source_id"],
                "current_path": source_manifest["current_path"],
                "document_type": source_manifest["document_type"],
                "support_tier": source_manifest.get("support_tier"),
                "source_origin": source_manifest.get("source_origin", "original-document"),
                "parent_source_id": source_manifest.get("parent_source_id"),
                "root_email_source_id": source_manifest.get("root_email_source_id"),
                "unit_count": len(evidence_manifest["units"]),
                "render_count": len(evidence_manifest["document_renders"]),
                "failure_count": len(evidence_manifest["failures"]),
                "ambiguous_identity": bool(source_entry.get("ambiguous_match", False)),
                "change_classification": classification,
                "reused": bool(
                    source_entry.get(
                        "source_reused",
                        classification == "unchanged"
                        or "source_reused"
                        in [
                            str(trait)
                            for trait in source_manifest.get("change_traits", [])
                            if isinstance(trait, str)
                        ],
                    )
                ),
            }
        )

    catalog = {
        "generated_at": utc_now(),
        "source_signature": source_inventory_signature(paths),
        "source_count": len(catalog_sources),
        "sources": catalog_sources,
    }
    coverage = {
        "generated_at": utc_now(),
        "source_count": len(source_summaries),
        "unit_count": sum(source["unit_count"] for source in source_summaries),
        "render_count": sum(source["render_count"] for source in source_summaries),
        "failure_count": sum(source["failure_count"] for source in source_summaries),
        "ambiguous_identity_detected": ambiguous_match,
        "reused_sources": sum(1 for source in source_summaries if source["reused"]),
        "rebuilt_sources": sum(1 for source in source_summaries if not source["reused"]),
        "sources": source_summaries,
    }
    write_json(paths.staging_catalog_path, catalog)
    write_json(paths.staging_coverage_manifest_path, coverage)
    write_json(paths.staging_graph_edges_path, {"generated_at": utc_now(), "edges": []})
    write_json(
        paths.staging_pending_work_path,
        {"generated_at": utc_now(), "pending_sources": []},
    )
    write_json(
        paths.staging_hybrid_work_path,
        {"generated_at": utc_now(), "target": "staging", "sources": []},
    )
    write_json(
        paths.staging_publish_manifest_path,
        {
            "staged_at": utc_now(),
            "published_at": None,
            "source_signature": source_inventory_signature(paths),
            "validation_status": "not-run",
        },
    )
    return catalog_sources, source_summaries, ambiguous_match


def _sorted_string_list(value: Any) -> list[str]:
    return (
        sorted(str(item) for item in value if isinstance(item, str))
        if isinstance(value, list)
        else []
    )


def _read_unit_text(source_dir: Path, unit: dict[str, Any]) -> str:
    text_asset = unit.get("text_asset")
    if not isinstance(text_asset, str) or not text_asset:
        return ""
    asset_path = source_dir / text_asset
    if not asset_path.exists():
        return ""
    return sanitize_text(asset_path.read_text(encoding="utf-8"))


def semantic_evidence_signature(
    source_dir: Path,
    *,
    include_artifacts: bool = True,
) -> str | None:
    """Return a stable semantic signature for staged evidence, independent of file path."""
    evidence_manifest = read_json(source_dir / "evidence_manifest.json")
    if not evidence_manifest:
        return None
    units: list[dict[str, Any]] = []
    for raw_unit in evidence_manifest.get("units", []):
        if not isinstance(raw_unit, dict):
            continue
        unit = dict(raw_unit)
        units.append(
            {
                "unit_id": unit.get("unit_id"),
                "unit_type": unit.get("unit_type"),
                "ordinal": unit.get("ordinal"),
                "title": unit.get("title"),
                "line_start": unit.get("line_start"),
                "line_end": unit.get("line_end"),
                "slug_anchor": unit.get("slug_anchor"),
                "row_count": unit.get("row_count"),
                "header_names": _sorted_string_list(unit.get("header_names")),
                "text": _read_unit_text(source_dir, unit),
            }
        )
    payload = {
        "document_type": evidence_manifest.get("document_type"),
        "units": units,
        "warnings": evidence_warning_messages(evidence_manifest),
    }
    if include_artifacts:
        artifact_index = read_json(source_dir / "artifact_index.json")
        payload["artifacts"] = [
            {
                "artifact_id": artifact.get("artifact_id"),
                "artifact_type": artifact.get("artifact_type"),
                "unit_id": artifact.get("unit_id"),
                "title": artifact.get("title"),
            }
            for artifact in artifact_index.get("artifacts", [])
            if isinstance(artifact, dict)
        ]
    return hashlib.sha256(str(payload).encode("utf-8")).hexdigest()


def _prune_source_related_sources(
    related_sources: Any,
    *,
    active_source_ids: set[str],
) -> tuple[list[dict[str, Any]], int]:
    filtered: list[dict[str, Any]] = []
    pruned = 0
    if not isinstance(related_sources, list):
        return filtered, pruned
    for related in related_sources:
        if not isinstance(related, dict):
            pruned += 1
            continue
        related_source_id = related.get("source_id")
        if not isinstance(related_source_id, str) or related_source_id not in active_source_ids:
            pruned += 1
            continue
        filtered.append(related)
    return filtered, pruned


def render_summary_markdown(
    *,
    title: str,
    source_id: str,
    summary_en: str,
    summary_source: str,
) -> str:
    """Render the canonical summary markdown contract for one semantic object."""
    return "\n".join(
        [
            f"# {title.strip()}",
            "",
            f"Source ID: {source_id}",
            "",
            "## English Summary",
            summary_en.strip(),
            "",
            "## Source-Language Summary",
            summary_source.strip(),
            "",
        ]
    )


def refresh_source_semantic_outputs(
    source_dir: Path,
    *,
    active_source_ids: set[str],
) -> dict[str, int]:
    """Refresh reused source knowledge metadata and silently prune deleted relations."""
    source_manifest = read_json(source_dir / "source_manifest.json")
    evidence_manifest_path = source_dir / "evidence_manifest.json"
    evidence_manifest = read_json(evidence_manifest_path)
    if evidence_manifest:
        write_json(
            evidence_manifest_path,
            sync_optional_sidecar_assets(
                source_dir,
                evidence_manifest=evidence_manifest,
            ),
        )
    knowledge_path = source_dir / "knowledge.json"
    knowledge = read_json(knowledge_path)
    if not source_manifest or not knowledge:
        return {"knowledge_refreshed": 0, "summary_rebuilt": 0, "relations_pruned": 0}

    changed = False
    if knowledge.get("source_id") != source_manifest.get("source_id"):
        knowledge["source_id"] = source_manifest.get("source_id")
        changed = True
    if knowledge.get("source_fingerprint") != source_manifest.get("source_fingerprint"):
        knowledge["source_fingerprint"] = source_manifest.get("source_fingerprint")
        changed = True
    if knowledge.get("document_type") != source_manifest.get("document_type"):
        knowledge["document_type"] = source_manifest.get("document_type")
        changed = True

    filtered_related_sources, relations_pruned = _prune_source_related_sources(
        knowledge.get("related_sources", []),
        active_source_ids=active_source_ids,
    )
    if knowledge.get("related_sources") != filtered_related_sources:
        knowledge["related_sources"] = filtered_related_sources
        changed = True

    if changed:
        write_json(knowledge_path, knowledge)

    summary_path = source_dir / "summary.md"
    summary_rebuilt = 0
    if not summary_path.exists() or not summary_path.read_text(encoding="utf-8").strip():
        title = str(
            knowledge.get("title") or source_manifest.get("title") or source_manifest["source_id"]
        )
        summary_en = (
            sanitize_text(knowledge.get("summary_en"))
            or "No concise English summary is currently available."
        )
        summary_source = sanitize_text(knowledge.get("summary_source")) or summary_en
        summary_path.write_text(
            render_summary_markdown(
                title=title,
                source_id=str(source_manifest["source_id"]),
                summary_en=summary_en,
                summary_source=summary_source,
            ),
            encoding="utf-8",
        )
        summary_rebuilt = 1

    return {
        "knowledge_refreshed": 1 if changed else 0,
        "summary_rebuilt": summary_rebuilt,
        "relations_pruned": relations_pruned,
    }


def _compact_excerpt(text: str, *, limit: int = 220) -> str:
    compact = re.sub(r"\s+", " ", sanitize_text(text)).strip()
    if len(compact) <= limit:
        return compact
    truncated = compact[:limit].rsplit(" ", 1)[0].strip()
    return (truncated or compact[:limit]).rstrip(" ,;:.") + "..."


def _fallback_title(source_manifest: dict[str, Any]) -> str:
    title = sanitize_text(source_manifest.get("title"))
    if title:
        return title
    current_path = str(source_manifest.get("current_path") or "")
    stem = Path(current_path).stem.replace("_", " ").replace("-", " ").strip()
    return stem or str(source_manifest.get("source_id") or "Untitled Source")


def _collect_unit_snippets(
    source_dir: Path,
    evidence_manifest: dict[str, Any],
    *,
    limit: int = 3,
) -> list[dict[str, str]]:
    snippets: list[dict[str, str]] = []
    for raw_unit in evidence_manifest.get("units", []):
        if not isinstance(raw_unit, dict):
            continue
        unit_id = raw_unit.get("unit_id")
        if not isinstance(unit_id, str) or not unit_id:
            continue
        excerpt = _compact_excerpt(_read_unit_text(source_dir, raw_unit))
        if not excerpt:
            title = sanitize_text(raw_unit.get("title"))
            if title:
                excerpt = title
        if not excerpt:
            continue
        snippets.append({"unit_id": unit_id, "excerpt": excerpt})
        if len(snippets) >= limit:
            break
    return snippets


def _collect_artifact_snippets(
    source_dir: Path,
    *,
    limit: int = 4,
) -> list[dict[str, str]]:
    artifact_index = read_json(source_dir / "artifact_index.json")
    snippets: list[dict[str, str]] = []
    for artifact in artifact_index.get("artifacts", []):
        if not isinstance(artifact, dict):
            continue
        artifact_id = artifact.get("artifact_id")
        unit_id = artifact.get("unit_id")
        if not isinstance(artifact_id, str) or not artifact_id:
            continue
        if not isinstance(unit_id, str) or not unit_id:
            continue
        title = _compact_excerpt(
            str(artifact.get("title") or artifact.get("linked_text") or ""),
            limit=140,
        )
        if not title:
            continue
        if not (
            bool(artifact.get("graph_promoted"))
            or artifact.get("artifact_type") in {"chart", "table", "major-region"}
        ):
            continue
        snippets.append(
            {
                "unit_id": unit_id,
                "artifact_id": artifact_id,
                "excerpt": title,
            }
        )
        if len(snippets) >= limit:
            break
    return snippets


def _collect_overlay_snippets(
    source_dir: Path,
    *,
    limit: int = 4,
) -> list[dict[str, str]]:
    overlays = load_semantic_overlays(source_dir)
    snippets: list[dict[str, str]] = []
    for unit_id, payload in overlays.items():
        confidence = overlay_confidence(payload)
        if confidence == "low":
            continue
        for excerpt in overlay_search_strings(payload)[:3]:
            compact = _compact_excerpt(excerpt, limit=160)
            if not compact:
                continue
            snippets.append(
                {
                    "unit_id": unit_id,
                    "excerpt": compact,
                }
            )
            if len(snippets) >= limit:
                return snippets
    return snippets


def _visual_role_hints(source_dir: Path, evidence_manifest: dict[str, Any]) -> list[str]:
    hints: list[str] = []
    for asset in evidence_manifest.get("visual_layout_assets", []):
        if not isinstance(asset, str) or not asset:
            continue
        payload = read_json(source_dir / asset)
        if not isinstance(payload, dict):
            continue
        hints.extend(
            hint for hint in payload.get("role_hints", []) if isinstance(hint, str) and hint
        )
    return list(dict.fromkeys(hints))


def write_conservative_semantic_outputs(
    source_dir: Path,
    *,
    active_source_ids: set[str],
) -> dict[str, Any]:
    """Write a minimal but publishable knowledge object directly from staged evidence."""
    source_manifest = read_json(source_dir / "source_manifest.json")
    evidence_manifest = read_json(source_dir / "evidence_manifest.json")
    if not source_manifest or not evidence_manifest:
        raise RuntimeError(f"Missing staged manifests for autonomous authoring in {source_dir}.")

    source_id = str(source_manifest["source_id"])
    document_type = str(source_manifest.get("document_type") or "unknown")
    title = _fallback_title(source_manifest)
    snippets = _collect_unit_snippets(source_dir, evidence_manifest, limit=3)
    artifact_snippets = _collect_artifact_snippets(source_dir, limit=4)
    overlay_snippets = _collect_overlay_snippets(source_dir, limit=4)
    visual_role_hints = _visual_role_hints(source_dir, evidence_manifest)
    warnings = evidence_warning_messages(evidence_manifest)
    failures = evidence_manifest.get("failures", [])
    if not isinstance(failures, list):
        failures = []
    source_language = detect_source_language(
        [snippet["excerpt"] for snippet in [*overlay_snippets, *artifact_snippets, *snippets]]
        or [title]
    )
    primary_citations = [
        {"unit_id": snippet["unit_id"], "support": "Autonomous in-repo semantic summary"}
        for snippet in snippets[:3]
    ]
    artifact_citations = [
        {
            "unit_id": snippet["unit_id"],
            "artifact_id": snippet["artifact_id"],
            "support": "Autonomous artifact-backed semantic summary",
        }
        for snippet in artifact_snippets[:3]
    ]
    overlay_citations = [
        {
            "unit_id": snippet["unit_id"],
            "support": "Autonomous semantic-overlay-backed summary",
        }
        for snippet in overlay_snippets[:2]
    ]
    if artifact_citations:
        primary_citations = artifact_citations + primary_citations
    elif overlay_citations:
        primary_citations = overlay_citations + primary_citations
    if not primary_citations:
        first_unit_id = next(
            (
                str(unit["unit_id"])
                for unit in evidence_manifest.get("units", [])
                if isinstance(unit, dict) and isinstance(unit.get("unit_id"), str)
            ),
            None,
        )
        if first_unit_id is not None:
            primary_citations.append(
                {
                    "unit_id": first_unit_id,
                    "support": "Rendered or structured evidence preserved for manual inspection",
                }
            )

    if document_type == "interaction":
        interaction_context = read_json(source_dir / "interaction_context.json")
        filtered_related_sources, _relations_pruned = _prune_source_related_sources(
            interaction_context.get("related_sources", []),
            active_source_ids=active_source_ids,
        )
        primary_snippet_excerpt = (
            snippets[0]["excerpt"]
            if snippets
            else "The promoted turns were preserved, but extracted text remains limited."
        )
        summary_en = (
            f"This interaction-derived memory preserves conversation evidence for `{title}`. "
            f"{primary_snippet_excerpt}"
        )
        summary_source = (
            snippets[0]["excerpt"] if source_language != "en" and snippets else summary_en
        )
        key_points = [
            {
                "text_en": (
                    snippets[0]["excerpt"]
                    if snippets
                    else "The promoted interaction memory preserves bounded turn-level evidence."
                ),
                "text_source": (
                    snippets[0]["excerpt"]
                    if snippets
                    else "The promoted interaction memory preserves bounded turn-level evidence."
                ),
                "citations": primary_citations[:1],
            }
        ]
        claims = [
            {
                "statement_en": (
                    "This memory is interaction-derived context and should remain distinct from "
                    "source-authored corpus documents."
                ),
                "statement_source": (
                    "This memory is interaction-derived context and should remain distinct from "
                    "source-authored corpus documents."
                ),
                "citations": primary_citations[:1],
            }
        ]
        entities = [{"name": title, "type": "interaction-memory"}]
        related_sources = filtered_related_sources
    elif document_type == "xlsx":
        workbook_payload = read_json(source_dir / "spreadsheet_workbook.json")
        sheet_inventory = workbook_payload.get("sheet_inventory", [])
        visible_sheets = [
            item
            for item in sheet_inventory
            if isinstance(item, dict) and item.get("visibility") == "visible"
        ]
        hidden_sheets = [
            item
            for item in sheet_inventory
            if isinstance(item, dict) and item.get("visibility") != "visible"
        ]
        chart_registry = workbook_payload.get("chart_registry", [])
        cross_sheet_summary = workbook_payload.get("cross_sheet_reference_summary", [])
        workbook_summary_bits = [
            f"{len(sheet_inventory)} sheets",
            f"{len(visible_sheets)} visible",
        ]
        if hidden_sheets:
            workbook_summary_bits.append(f"{len(hidden_sheets)} hidden")
        if isinstance(chart_registry, list) and chart_registry:
            workbook_summary_bits.append(f"{len(chart_registry)} charts")
        if isinstance(cross_sheet_summary, list) and cross_sheet_summary:
            workbook_summary_bits.append(f"{len(cross_sheet_summary)} cross-sheet reference groups")
        summary_en = "This workbook preserves " + ", ".join(workbook_summary_bits) + "."
        if overlay_snippets:
            summary_en += " Semantic overlay signals include: " + "; ".join(
                snippet["excerpt"] for snippet in overlay_snippets[:2]
            )
        elif artifact_snippets:
            summary_en += " Key artifact signals include: " + "; ".join(
                snippet["excerpt"] for snippet in artifact_snippets[:2]
            )
        summary_source = summary_en
        key_points = []
        for sheet in visible_sheets[:2]:
            if not isinstance(sheet, dict):
                continue
            point_parts = [str(sheet.get("sheet_name") or "Sheet")]
            metric_candidates = sheet.get("metric_candidates", [])
            if isinstance(metric_candidates, list) and metric_candidates:
                point_parts.append(
                    "metrics: " + ", ".join(str(value) for value in metric_candidates[:3])
                )
            time_axis_candidates = sheet.get("time_axis_candidates", [])
            if isinstance(time_axis_candidates, list) and time_axis_candidates:
                point_parts.append(
                    "time axes: " + ", ".join(str(value) for value in time_axis_candidates[:2])
                )
            key_points.append(
                {
                    "text_en": "; ".join(point_parts),
                    "text_source": "; ".join(point_parts),
                    "citations": primary_citations[:1],
                }
            )
        if not key_points:
            key_points = [
                {
                    "text_en": summary_en,
                    "text_source": summary_en,
                    "citations": primary_citations[:1],
                }
            ]
        claims = [
            {
                "statement_en": (
                    f"`{title}` is preserved as a structured workbook with "
                    "sheet-level, chart, and tabular evidence."
                ),
                "statement_source": (
                    f"`{title}` is preserved as a structured workbook with "
                    "sheet-level, chart, and tabular evidence."
                ),
                "citations": primary_citations[:1],
            }
        ]
        if hidden_sheets:
            claims.append(
                {
                    "statement_en": (
                        f"The workbook contains {len(hidden_sheets)} hidden "
                        "supporting sheets."
                    ),
                    "statement_source": (
                        f"The workbook contains {len(hidden_sheets)} hidden "
                        "supporting sheets."
                    ),
                    "citations": primary_citations[:1],
                }
            )
        entities = [{"name": title, "type": "xlsx-workbook"}]
        related_sources = []
    elif document_type in {"pdf", "pptx", "docx"} and (
        overlay_snippets or artifact_snippets or visual_role_hints
    ):
        summary_seed = (
            "; ".join(snippet["excerpt"] for snippet in overlay_snippets[:2])
            if overlay_snippets
            else "; ".join(snippet["excerpt"] for snippet in artifact_snippets[:2])
            if artifact_snippets
            else ""
        )
        role_seed = ", ".join(visual_role_hints[:3]) if visual_role_hints else ""
        summary_en = (
            f"This source preserves published visual evidence with role hints such as {role_seed}. "
            f"{summary_seed}".strip()
        ).strip()
        if not summary_en:
            summary_en = (
                "This source preserves published visual evidence for manual "
                "and downstream reasoning."
            )
        summary_source = summary_en
        key_points = []
        for snippet in overlay_snippets[:2]:
            key_points.append(
                {
                    "text_en": snippet["excerpt"],
                    "text_source": snippet["excerpt"],
                    "citations": [
                        {
                            "unit_id": snippet["unit_id"],
                            "support": "Autonomous semantic-overlay-backed summary",
                        }
                    ],
                }
            )
        if not key_points:
            key_points = [
                {
                    "text_en": snippet["excerpt"],
                    "text_source": snippet["excerpt"],
                    "citations": [
                        {
                            "unit_id": snippet["unit_id"],
                            "artifact_id": snippet["artifact_id"],
                            "support": "Autonomous artifact-backed semantic summary",
                        }
                    ],
                }
                for snippet in artifact_snippets[:2]
            ]
        if not key_points:
            key_points = [
                {
                    "text_en": summary_en,
                    "text_source": summary_en,
                    "citations": primary_citations[:1],
                }
            ]
        claims = [
            {
                "statement_en": (
                    f"`{title}` includes visual-layout evidence that remains "
                    "usable without immediate source rerendering."
                ),
                "statement_source": (
                    f"`{title}` includes visual-layout evidence that remains "
                    "usable without immediate source rerendering."
                ),
                "citations": primary_citations[:1],
            }
        ]
        if visual_role_hints:
            claims.append(
                {
                    "statement_en": "Detected visual role hints include: "
                    + ", ".join(visual_role_hints[:4]),
                    "statement_source": "Detected visual role hints include: "
                    + ", ".join(visual_role_hints[:4]),
                    "citations": primary_citations[:1],
                }
            )
        entities = [{"name": title, "type": document_type}]
        related_sources = []
    else:
        if overlay_snippets:
            snippet_seed = list(overlay_snippets)
        elif artifact_snippets:
            snippet_seed = list(artifact_snippets)
        else:
            snippet_seed = list(snippets)
        topic_summary = (
            "; ".join(snippet["excerpt"] for snippet in snippet_seed[:2])
            if snippet_seed
            else (
                "The source preserved little extracted text, so interpretation "
                "remains conservative."
            )
        )
        summary_en = f"This source covers: {topic_summary}"
        summary_source = (
            snippet_seed[0]["excerpt"]
            if source_language != "en" and snippet_seed
            else summary_en
        )
        key_points = [
            {
                "text_en": (
                    snippet["excerpt"]
                    if snippets
                    else "The source preserved limited evidence for manual follow-up."
                ),
                "text_source": (
                    snippet["excerpt"]
                    if snippets
                    else "The source preserved limited evidence for manual follow-up."
                ),
                "citations": [citation],
            }
            for snippet, citation in zip(
                snippets[:2],
                primary_citations[:2],
                strict=False,
            )
        ]
        if not key_points:
            key_points = [
                {
                    "text_en": "The source preserved limited evidence for manual follow-up.",
                    "text_source": "The source preserved limited evidence for manual follow-up.",
                    "citations": primary_citations[:1],
                }
            ]
        evidence_claim = (
            f"The staged evidence for `{title}` explicitly includes: "
            f"{snippet_seed[0]['excerpt']}"
            if snippet_seed
            else (
                f"`{title}` is preserved with limited extracted text and may "
                "require render inspection."
            )
        )
        claims = [
            {
                "statement_en": evidence_claim,
                "statement_source": evidence_claim,
                "citations": primary_citations[:1],
            }
        ]
        entities = [{"name": title, "type": document_type}]
        related_sources = []

    known_gaps = [
        {
            "text_en": "Evidence warnings were preserved during staging: "
            + "; ".join(warnings[:3]),
            "text_source": "Evidence warnings were preserved during staging: "
            + "; ".join(warnings[:3]),
        }
        for _warning in ([warnings[0]] if warnings else [])
    ]
    if not snippets:
        limited_text_gap = (
            "Extracted text is limited, so manual render or structure "
            "inspection may still be required."
        )
        known_gaps.append(
            {
                "text_en": limited_text_gap,
                "text_source": limited_text_gap,
            }
        )

    ambiguities = [
        {
            "text_en": "Staging recorded partial extraction failures that may hide source details.",
            "text_source": (
                "Staging recorded partial extraction failures that may hide "
                "source details."
            ),
            "citations": primary_citations[:1],
        }
        for _failure in ([failures[0]] if failures else [])
    ]
    confidence_level = (
        "high" if snippets and not warnings and not failures else "medium" if snippets else "low"
    )
    confidence_note = (
        "Autonomous in-repo authoring synthesized this knowledge directly from staged evidence."
    )
    knowledge = {
        "source_id": source_id,
        "source_fingerprint": source_manifest.get("source_fingerprint"),
        "title": title,
        "source_language": source_language,
        "summary_en": summary_en,
        "summary_source": summary_source,
        "document_type": document_type,
        "key_points": key_points,
        "entities": entities,
        "claims": claims,
        "known_gaps": known_gaps,
        "ambiguities": ambiguities,
        "confidence": {
            "level": confidence_level,
            "notes_en": confidence_note,
            "notes_source": confidence_note,
        },
        "citations": primary_citations,
        "related_sources": related_sources,
    }
    write_json(source_dir / "knowledge.json", knowledge)
    (source_dir / "summary.md").write_text(
        render_summary_markdown(
            title=title,
            source_id=source_id,
            summary_en=summary_en,
            summary_source=summary_source,
        ),
        encoding="utf-8",
    )
    _write_source_affordances(
        source_dir,
        source_manifest=source_manifest,
        evidence_manifest=evidence_manifest,
        knowledge=knowledge,
        summary_text=render_summary_markdown(
            title=title,
            source_id=source_id,
            summary_en=summary_en,
            summary_source=summary_source,
        ),
    )
    return {
        "source_id": source_id,
        "document_type": document_type,
        "title": title,
        "confidence_level": confidence_level,
        "warning_count": len(warnings),
        "failure_count": len(failures),
    }


def repair_staging_semantic_artifacts(
    paths: WorkspacePaths,
    *,
    active_source_ids: set[str],
    active_sources: list[dict[str, Any]],
) -> dict[str, Any]:
    """Refresh reused staged semantic outputs before authoring or validation."""
    staged_source_ids = {
        *active_source_ids,
        *[
            source_dir.name
            for source_dir in sorted((paths.knowledge_base_staging_dir / "sources").glob("*"))
            if source_dir.is_dir()
        ],
    }
    source_lookup = {
        str(source.get("source_id")): source
        for source in active_sources
        if isinstance(source, dict) and isinstance(source.get("source_id"), str)
    }
    source_repairs = {
        "knowledge_refreshed": 0,
        "summary_rebuilt": 0,
        "relations_pruned": 0,
    }
    for source_dir in sorted((paths.knowledge_base_staging_dir / "sources").glob("*")):
        if not source_dir.is_dir():
            continue
        source_id = source_dir.name
        refreshed = refresh_source_semantic_outputs(
            source_dir,
            active_source_ids=staged_source_ids,
        )
        for key, value in refreshed.items():
            source_repairs[key] += value
        if int(refreshed.get("relations_pruned", 0)) > 0 and source_id in source_lookup:
            source_entry = source_lookup[source_id]
            source_entry["change_traits"] = append_unique_strings(
                [
                    *[
                        str(trait)
                        for trait in source_entry.get("change_traits", [])
                        if isinstance(trait, str) and trait
                    ],
                    "relation_shrink_only",
                ]
            )
            if source_entry.get("change_classification") == "unchanged":
                source_entry["change_reason"] = (
                    "Published related-source references were silently pruned because one or more "
                    "linked sources were deleted from the active corpus."
                )

    interaction_repairs = {
        "knowledge_relations_pruned": 0,
        "context_relations_pruned": 0,
    }
    for memory_dir in sorted(paths.interaction_memories_dir("staging").glob("*")):
        if not memory_dir.is_dir():
            continue
        refreshed = repair_interaction_memory_related_sources(
            memory_dir,
            active_source_ids=staged_source_ids,
        )
        interaction_repairs["knowledge_relations_pruned"] += int(
            refreshed.get("knowledge_pruned", 0)
        )
        interaction_repairs["context_relations_pruned"] += int(refreshed.get("context_pruned", 0))

    return {
        "source_repairs": source_repairs,
        "interaction_repairs": interaction_repairs,
        "repair_count": sum(source_repairs.values()) + sum(interaction_repairs.values()),
    }


def auto_author_pending_semantics(
    paths: WorkspacePaths,
    *,
    pending_sources: list[dict[str, Any]],
    active_source_ids: set[str],
) -> dict[str, Any]:
    """Write conservative semantic outputs for remaining pending staged items."""
    authored_items: list[dict[str, Any]] = []
    for item in pending_sources:
        if not isinstance(item, dict):
            continue
        source_id = item.get("source_id")
        if not isinstance(source_id, str) or not source_id:
            continue
        kind = str(item.get("kind") or "corpus")
        source_dir = (
            paths.interaction_memories_dir("staging") / source_id
            if kind == "interaction-memory"
            else paths.knowledge_base_staging_dir / "sources" / source_id
        )
        if not source_dir.exists():
            raise RuntimeError(f"Pending staged directory is missing for `{source_id}`.")
        authored_items.append(
            {
                **write_conservative_semantic_outputs(
                    source_dir,
                    active_source_ids=active_source_ids,
                ),
                "kind": kind,
                "reason": item.get("reason"),
            }
        )
    return {
        "attempted": len([item for item in pending_sources if isinstance(item, dict)]),
        "authored": authored_items,
        "authored_count": len(authored_items),
        "mode": "conservative-in-repo",
    }


def build_hybrid_work_queue(
    paths: WorkspacePaths,
    *,
    target: str = "staging",
) -> dict[str, Any]:
    """Build the staged hard-artifact hybrid work queue."""
    target_root = paths.knowledge_target_dir(target)
    source_root = target_root / "sources"
    payload: dict[str, Any] = {"generated_at": utc_now(), "target": target, "sources": []}
    if not source_root.exists():
        write_json(paths.hybrid_work_path(target), payload)
        return payload
    for source_dir in sorted(source_root.glob("*")):
        if not source_dir.is_dir():
            continue
        evidence_manifest = read_json(source_dir / "evidence_manifest.json")
        source_manifest = read_json(source_dir / "source_manifest.json")
        if not evidence_manifest or not source_manifest:
            continue
        source_packet = build_source_hybrid_packet(
            source_dir,
            evidence_manifest=evidence_manifest,
            source_manifest=source_manifest,
        )
        if not source_packet:
            continue
        payload["sources"].append(source_packet)
    payload["sources"].sort(
        key=lambda item: (
            -int(item.get("highest_remaining_priority", 0)),
            -int(item.get("remaining_candidate_count", 0)),
            str(item.get("source_path", "")),
            str(item.get("source_id", "")),
        )
    )
    write_json(paths.hybrid_work_path(target), payload)
    return payload


def hybrid_enrichment_status(paths: WorkspacePaths, *, target: str = "staging") -> dict[str, Any]:
    """Describe the current semantic-overlay coverage without blocking publication."""
    hybrid_work = build_hybrid_work_queue(paths, target=target)
    summary: dict[str, Any] = {
        "target": target,
        "capability_detected": False,
        "workflow_auto_supported": True,
        "hybrid_work_path": str(paths.hybrid_work_path(target).relative_to(paths.root)),
        "capability_gap_reason": "",
        **summarize_hybrid_work(hybrid_work),
    }
    if summary["mode"] in {"candidate-prepared", "partially-covered"}:
        summary["capability_gap_reason"] = (
            "Deterministic sync cannot complete hard-artifact multimodal enrichment by itself. "
            "A capable host agent workflow must consume hybrid_work.json and write additive "
            "semantic_overlay sidecars."
        )
    return summary


def refresh_change_set_details(
    change_set: dict[str, Any],
    active_sources: list[dict[str, Any]],
) -> dict[str, Any]:
    """Merge build-time change traits and reasons back into the public change set."""
    changes = change_set.get("changes", [])
    if not isinstance(changes, list):
        return change_set
    source_lookup = {
        str(source.get("source_id")): source
        for source in active_sources
        if isinstance(source, dict) and isinstance(source.get("source_id"), str)
    }
    enriched_changes: list[dict[str, Any]] = []
    for change in changes:
        if not isinstance(change, dict):
            continue
        source_id = change.get("source_id")
        if isinstance(source_id, str) and source_id in source_lookup:
            source_entry = source_lookup[source_id]
            updated_change = dict(change)
            updated_change["change_traits"] = source_entry.get("change_traits", [])
            updated_change["change_reason"] = source_entry.get("change_reason")
            enriched_changes.append(updated_change)
        else:
            enriched_changes.append(change)
    return {**change_set, "changes": enriched_changes}


def build_staging_artifacts(
    paths: WorkspacePaths,
    active_sources: list[dict[str, Any]],
    office_binary: str | None,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], bool, dict[str, int]]:
    """Assemble a new staging tree with selective reuse plus targeted rebuilds."""
    staging_temp_dir = paths.knowledge_base_dir / ".staging-build"
    if staging_temp_dir.exists():
        shutil.rmtree(staging_temp_dir)
    ensure_directory(staging_temp_dir / "sources")

    reused_sources = 0
    rebuilt_sources = 0
    for source_entry in active_sources:
        source_id = str(source_entry["source_id"])
        source_dir = staging_temp_dir / "sources" / source_id
        previous_source_dir = locate_previous_source_dir(paths, source_id)
        classification = str(source_entry.get("change_classification", "added"))
        previous_manifest = (
            read_json(previous_source_dir / "source_manifest.json")
            if previous_source_dir is not None
            else {}
        )
        previous_fingerprint = previous_manifest.get("source_fingerprint")
        source_entry["source_reused"] = False
        source_entry["semantic_outputs_reused"] = False
        source_entry["semantic_signature_stable"] = False
        change_traits = [
            str(trait)
            for trait in source_entry.get("change_traits", [])
            if isinstance(trait, str) and trait
        ]
        is_email_source = str(source_entry.get("document_type")) == "email"

        reuse_previous_directory = bool(
            previous_source_dir is not None
            and (
                classification == "unchanged"
                or (
                    classification == "moved-or-renamed"
                    and previous_fingerprint == source_entry.get("source_fingerprint")
                )
            )
        )
        if is_email_source and classification == "moved-or-renamed":
            reuse_previous_directory = False
        previous_contract_complete = bool(
            previous_source_dir is not None
            and source_artifact_contract_complete(
                previous_source_dir,
                document_type=str(source_entry.get("document_type") or "unknown"),
            )
        )
        if reuse_previous_directory and not previous_contract_complete:
            reuse_previous_directory = False

        if reuse_previous_directory and previous_source_dir is not None:
            if is_email_source:
                _copy_reused_source_tree(
                    paths,
                    source_entry=source_entry,
                    destination_root=staging_temp_dir / "sources",
                )
            else:
                if source_dir.exists():
                    shutil.rmtree(source_dir)
                shutil.copytree(previous_source_dir, source_dir)
                refresh_reused_source_metadata(paths, source_entry, source_dir)
            source_entry["source_reused"] = True
            source_entry["semantic_outputs_reused"] = True
            change_traits = append_unique_strings(
                [*change_traits, "source_reused", "semantic_outputs_reused"]
            )
            if "path_changed" in change_traits:
                source_entry["change_reason"] = (
                    "The source path changed but the binary fingerprint stayed stable, so the "
                    "previous staged evidence and semantic outputs were reused."
                )
            else:
                source_entry["change_reason"] = (
                    "The source fingerprint is unchanged, so the previous staged evidence and "
                    "semantic outputs were reused."
                )
            reused_sources += 1
        else:
            if source_dir.exists():
                shutil.rmtree(source_dir)
            ensure_directory(source_dir)
            if is_email_source:
                build_email_source_tree(
                    paths,
                    source_entry=source_entry,
                    source_dir=source_dir,
                    staging_sources_dir=staging_temp_dir / "sources",
                    office_binary=office_binary,
                )
                if previous_source_dir is not None:
                    previous_signature = semantic_evidence_signature(previous_source_dir)
                    current_signature = semantic_evidence_signature(source_dir)
                    previous_base_signature = semantic_evidence_signature(
                        previous_source_dir,
                        include_artifacts=False,
                    )
                    current_base_signature = semantic_evidence_signature(
                        source_dir,
                        include_artifacts=False,
                    )
                    if previous_signature is not None and previous_signature == current_signature:
                        preserve_semantic_outputs(previous_source_dir, source_dir)
                        _persist_source_artifacts(
                            source_dir,
                            source_manifest=read_json(source_dir / "source_manifest.json"),
                            evidence_manifest=read_json(source_dir / "evidence_manifest.json"),
                        )
                        source_entry["semantic_outputs_reused"] = True
                        source_entry["semantic_signature_stable"] = True
                        change_traits = append_unique_strings(
                            [*change_traits, "semantic_outputs_reused"]
                        )
                        if "binary_changed" in change_traits and "path_changed" in change_traits:
                            source_entry["change_reason"] = (
                                "The email path and binary fingerprint changed, but the rebuilt "
                                "email evidence stayed semantically stable, so prior "
                                "semantic outputs were refreshed and reused."
                            )
                        elif "binary_changed" in change_traits:
                            source_entry["change_reason"] = (
                                "The email fingerprint changed, but the rebuilt "
                                "email evidence stayed semantically stable enough "
                                "to reuse prior semantic outputs."
                            )
                        elif "path_changed" in change_traits:
                            source_entry["change_reason"] = (
                                "The email path changed and staged evidence was rebuilt, but the "
                                "published semantic content stayed stable enough to reuse."
                            )
                    elif (
                        not previous_contract_complete
                        and previous_base_signature is not None
                        and previous_base_signature == current_base_signature
                    ):
                        preserve_semantic_outputs(previous_source_dir, source_dir)
                        _persist_source_artifacts(
                            source_dir,
                            source_manifest=read_json(source_dir / "source_manifest.json"),
                            evidence_manifest=read_json(source_dir / "evidence_manifest.json"),
                        )
                        source_entry["semantic_outputs_reused"] = True
                        source_entry["semantic_signature_stable"] = True
                        change_traits = append_unique_strings(
                            [
                                *change_traits,
                                "semantic_outputs_reused",
                                "artifact_contract_backfilled",
                            ]
                        )
                        source_entry["change_reason"] = (
                            "The source needed a Phase 3 artifact-contract backfill, "
                            "but the rebuilt evidence stayed semantically stable "
                            "enough to reuse prior semantic outputs."
                        )
                    else:
                        change_traits = append_unique_strings(
                            [*change_traits, "extracted_content_changed"]
                        )
                        if "path_changed" in change_traits and "binary_changed" in change_traits:
                            source_entry["change_reason"] = (
                                "The email path changed and the rebuilt email "
                                "evidence also changed semantically, so semantic "
                                "outputs must be refreshed to current truth."
                            )
                        elif "binary_changed" in change_traits:
                            source_entry["change_reason"] = (
                                "The rebuilt email evidence changed semantically, "
                                "so semantic outputs must be refreshed to current "
                                "truth."
                            )
                        elif "path_changed" in change_traits:
                            source_entry["change_reason"] = (
                                "The email path changed and the rebuilt evidence "
                                "could not be proven semantically stable, so "
                                "semantic outputs must be refreshed."
                            )
            else:
                source_manifest, evidence_manifest = build_single_source_artifacts(
                    paths,
                    source_entry,
                    source_dir,
                    office_binary,
                )
                _persist_source_artifacts(
                    source_dir,
                    source_manifest=source_manifest,
                    evidence_manifest=evidence_manifest,
                )
                if previous_source_dir is not None:
                    previous_signature = semantic_evidence_signature(previous_source_dir)
                    current_signature = semantic_evidence_signature(source_dir)
                    previous_base_signature = semantic_evidence_signature(
                        previous_source_dir,
                        include_artifacts=False,
                    )
                    current_base_signature = semantic_evidence_signature(
                        source_dir,
                        include_artifacts=False,
                    )
                    if previous_signature is not None and previous_signature == current_signature:
                        preserve_semantic_outputs(previous_source_dir, source_dir)
                        source_entry["semantic_outputs_reused"] = True
                        source_entry["semantic_signature_stable"] = True
                        change_traits = append_unique_strings(
                            [*change_traits, "semantic_outputs_reused"]
                        )
                        if "binary_changed" in change_traits and "path_changed" in change_traits:
                            source_entry["change_reason"] = (
                                "The source path and binary fingerprint changed, but the rebuilt "
                                "staged evidence stayed semantically stable, so prior semantic "
                                "outputs were refreshed and reused."
                            )
                        elif "binary_changed" in change_traits:
                            source_entry["change_reason"] = (
                                "The binary fingerprint changed, but the rebuilt staged evidence "
                                "stayed semantically stable, so prior semantic outputs were reused."
                            )
                        elif "path_changed" in change_traits:
                            source_entry["change_reason"] = (
                                "The source path changed and staged evidence was rebuilt, but the "
                                "semantic content stayed stable enough to reuse "
                                "prior semantic outputs."
                            )
                    elif (
                        not previous_contract_complete
                        and previous_base_signature is not None
                        and previous_base_signature == current_base_signature
                    ):
                        preserve_semantic_outputs(previous_source_dir, source_dir)
                        source_entry["semantic_outputs_reused"] = True
                        source_entry["semantic_signature_stable"] = True
                        change_traits = append_unique_strings(
                            [
                                *change_traits,
                                "semantic_outputs_reused",
                                "artifact_contract_backfilled",
                            ]
                        )
                        source_entry["change_reason"] = (
                            "The source needed a Phase 3 artifact-contract backfill, "
                            "but the rebuilt evidence stayed semantically stable "
                            "enough to reuse prior semantic outputs."
                        )
                    else:
                        change_traits = append_unique_strings(
                            [*change_traits, "extracted_content_changed"]
                        )
                        if "path_changed" in change_traits and "binary_changed" in change_traits:
                            source_entry["change_reason"] = (
                                "The source path changed and the rebuilt staged "
                                "evidence also changed semantically, so semantic "
                                "outputs must be refreshed to current truth."
                            )
                        elif "binary_changed" in change_traits:
                            source_entry["change_reason"] = (
                                "The rebuilt staged evidence changed semantically, "
                                "so semantic outputs must be refreshed to current "
                                "truth."
                            )
                        elif "path_changed" in change_traits:
                            source_entry["change_reason"] = (
                                "The source path changed and staged evidence could not be proven "
                                "semantically stable, so semantic outputs must be refreshed."
                            )
            rebuilt_sources += 1

        source_entry["change_traits"] = change_traits
        source_manifest = read_json(source_dir / "source_manifest.json")
        evidence_manifest = read_json(source_dir / "evidence_manifest.json")
    if paths.knowledge_base_staging_dir.exists():
        shutil.rmtree(paths.knowledge_base_staging_dir)
    staging_temp_dir.rename(paths.knowledge_base_staging_dir)
    catalog_sources, source_summaries, ambiguous_match = write_staging_root_artifacts(
        paths,
        active_sources,
    )
    return (
        catalog_sources,
        source_summaries,
        ambiguous_match,
        {
            "reused_sources": reused_sources,
            "rebuilt_sources": rebuilt_sources,
        },
    )


def collect_pending_synthesis(paths: WorkspacePaths) -> list[dict[str, Any]]:
    """Return the list of staged sources that still need agent-authored knowledge."""
    catalog = read_json(paths.staging_catalog_path)
    pending: list[dict[str, Any]] = []
    for source in catalog.get("sources", []):
        if not isinstance(source, dict):
            continue
        source_id = source.get("source_id")
        if not isinstance(source_id, str):
            continue
        source_dir = paths.knowledge_base_staging_dir / "sources" / source_id
        knowledge_path = source_dir / "knowledge.json"
        summary_path = source_dir / "summary.md"
        source_manifest = read_json(source_dir / "source_manifest.json")
        expected_fingerprint = source_manifest.get("source_fingerprint")

        stale = False
        if knowledge_path.exists():
            knowledge = read_json(knowledge_path)
            stale = knowledge.get("source_fingerprint") != expected_fingerprint
        needs_summary = (
            not summary_path.exists() or not summary_path.read_text(encoding="utf-8").strip()
        )
        if not knowledge_path.exists() or needs_summary or stale:
            pending.append(
                {
                    "source_id": source_id,
                    "current_path": source.get("current_path"),
                    "knowledge_path": str(knowledge_path.relative_to(paths.root)),
                    "summary_path": str(summary_path.relative_to(paths.root)),
                    "evidence_manifest_path": str(
                        (source_dir / "evidence_manifest.json").relative_to(paths.root)
                    ),
                    "reason": "stale" if stale else "missing",
                }
            )

    interaction_manifest = read_json(paths.interaction_manifest_path("staging"))
    for memory in interaction_manifest.get("memories", []):
        if not isinstance(memory, dict):
            continue
        memory_id = memory.get("memory_id")
        if not isinstance(memory_id, str):
            continue
        memory_dir = paths.interaction_memories_dir("staging") / memory_id
        knowledge_path = memory_dir / "knowledge.json"
        summary_path = memory_dir / "summary.md"
        source_manifest = read_json(memory_dir / "source_manifest.json")
        expected_fingerprint = source_manifest.get("source_fingerprint")

        stale = False
        if knowledge_path.exists():
            knowledge = read_json(knowledge_path)
            stale = knowledge.get("source_fingerprint") != expected_fingerprint
        needs_summary = (
            not summary_path.exists() or not summary_path.read_text(encoding="utf-8").strip()
        )
        if not knowledge_path.exists() or needs_summary or stale:
            pending.append(
                {
                    "source_id": memory_id,
                    "current_path": source_manifest.get("current_path"),
                    "knowledge_path": str(knowledge_path.relative_to(paths.root)),
                    "summary_path": str(summary_path.relative_to(paths.root)),
                    "evidence_manifest_path": str(
                        (memory_dir / "evidence_manifest.json").relative_to(paths.root)
                    ),
                    "work_item_path": str((memory_dir / "work_item.json").relative_to(paths.root)),
                    "kind": "interaction-memory",
                    "reason": "stale" if stale else "missing",
                }
            )

    pending.sort(
        key=lambda item: (
            str(item.get("current_path", "")),
            str(item.get("kind", "corpus")),
            str(item.get("source_id", "")),
        )
    )
    write_json(
        paths.staging_pending_work_path, {"generated_at": utc_now(), "pending_sources": pending}
    )
    return pending


def artifact_signature(path: Path) -> str | None:
    """Return a content signature for a file when it exists."""
    if not path.exists() or not path.is_file():
        return None
    return file_sha256(path)


def update_dependency_state(
    paths: WorkspacePaths,
    *,
    target: str,
    pending_sources: list[dict[str, Any]],
    retrieval_manifest: dict[str, Any] | None,
    trace_manifest: dict[str, Any] | None,
) -> dict[str, Any]:
    """Persist a machine-readable dependency summary for the latest target build."""
    target_root = paths.knowledge_target_dir(target)
    catalog = read_json(target_root / "catalog.json")
    coverage = read_json(target_root / "coverage_manifest.json")
    graph_edges = read_json(target_root / "graph_edges.json")
    pending_lookup = {
        str(item.get("source_id")): item for item in pending_sources if isinstance(item, dict)
    }
    coverage_lookup = {
        str(item.get("source_id")): item
        for item in coverage.get("sources", [])
        if isinstance(item, dict) and isinstance(item.get("source_id"), str)
    }
    source_dependencies: list[dict[str, Any]] = []
    for source in catalog.get("sources", []):
        if not isinstance(source, dict) or not isinstance(source.get("source_id"), str):
            continue
        source_id = source["source_id"]
        source_dir = target_root / "sources" / source_id
        source_manifest = read_json(source_dir / "source_manifest.json")
        source_dependencies.append(
            {
                "source_id": source_id,
                "current_path": source.get("current_path"),
                "document_type": source.get("document_type"),
                "support_tier": source_manifest.get("support_tier"),
                "change_classification": coverage_lookup.get(source_id, {}).get(
                    "change_classification"
                ),
                "source_fingerprint": source_manifest.get("source_fingerprint"),
                "artifact_fingerprints": {
                    "source_manifest": artifact_signature(source_dir / "source_manifest.json"),
                    "evidence_manifest": artifact_signature(source_dir / "evidence_manifest.json"),
                    "artifact_index": artifact_signature(source_dir / "artifact_index.json"),
                    "knowledge": artifact_signature(source_dir / "knowledge.json"),
                    "summary": artifact_signature(source_dir / "summary.md"),
                    "derived_affordances": artifact_signature(
                        source_dir / DEFAULT_AFFORDANCE_FILENAME
                    ),
                    "spreadsheet_workbook": artifact_signature(
                        source_dir / "spreadsheet_workbook.json"
                    ),
                },
                "pending_reason": pending_lookup.get(source_id, {}).get("reason"),
            }
        )

    payload = {
        "generated_at": utc_now(),
        "target": target,
        "source_signature": catalog.get("source_signature"),
        "pending_sources": len(pending_sources),
        "sources": source_dependencies,
        "root_artifacts": {
            "catalog": {
                "path": str((target_root / "catalog.json").relative_to(paths.root)),
                "fingerprint": artifact_signature(target_root / "catalog.json"),
                "dependent_source_ids": [source["source_id"] for source in source_dependencies],
            },
            "coverage_manifest": {
                "path": str((target_root / "coverage_manifest.json").relative_to(paths.root)),
                "fingerprint": artifact_signature(target_root / "coverage_manifest.json"),
                "dependent_source_ids": [source["source_id"] for source in source_dependencies],
            },
            "graph_edges": {
                "path": str((target_root / "graph_edges.json").relative_to(paths.root)),
                "fingerprint": artifact_signature(target_root / "graph_edges.json"),
                "dependent_source_ids": [
                    source["source_id"] for source in source_dependencies if source.get("source_id")
                ],
                "edge_count": len(graph_edges.get("edges", [])),
            },
            "pending_work": {
                "path": str((target_root / "pending_work.json").relative_to(paths.root)),
                "fingerprint": artifact_signature(target_root / "pending_work.json"),
                "dependent_source_ids": [source["source_id"] for source in source_dependencies],
            },
            "hybrid_work": {
                "path": str(paths.hybrid_work_path(target).relative_to(paths.root)),
                "fingerprint": artifact_signature(paths.hybrid_work_path(target)),
                "dependent_source_ids": [source["source_id"] for source in source_dependencies],
            },
            "retrieval_manifest": {
                "path": str(paths.retrieval_manifest_path(target).relative_to(paths.root)),
                "fingerprint": artifact_signature(paths.retrieval_manifest_path(target)),
                "present": retrieval_manifest is not None,
            },
            "retrieval_artifact_records": {
                "path": str(paths.retrieval_artifact_records_path(target).relative_to(paths.root)),
                "fingerprint": artifact_signature(paths.retrieval_artifact_records_path(target)),
                "present": paths.retrieval_artifact_records_path(target).exists(),
            },
            "trace_manifest": {
                "path": str(paths.trace_manifest_path(target).relative_to(paths.root)),
                "fingerprint": artifact_signature(paths.trace_manifest_path(target)),
                "present": trace_manifest is not None,
            },
            "interaction_manifest": {
                "path": str(paths.interaction_manifest_path(target).relative_to(paths.root)),
                "fingerprint": artifact_signature(paths.interaction_manifest_path(target)),
                "present": paths.interaction_manifest_path(target).exists(),
            },
        },
    }
    write_json(paths.dependency_state_path, payload)
    return payload


def iter_strings(value: Any) -> Iterable[str]:
    """Yield nested string values from a JSON-like structure."""
    if isinstance(value, str):
        yield value
    elif isinstance(value, list):
        for item in value:
            yield from iter_strings(item)
    elif isinstance(value, dict):
        for item in value.values():
            yield from iter_strings(item)


def contains_placeholder(value: Any) -> bool:
    """Return whether a JSON-like structure contains placeholder content."""
    for text in iter_strings(value):
        normalized = text.lower()
        if any(term in normalized for term in PLACEHOLDER_TERMS):
            return True
    return False


def citations_from_knowledge(knowledge: dict[str, Any]) -> list[dict[str, Any]]:
    """Collect all citation objects from a knowledge payload."""
    citations: list[dict[str, Any]] = []
    if isinstance(knowledge.get("citations"), list):
        citations.extend(item for item in knowledge["citations"] if isinstance(item, dict))
    for key in ("key_points", "claims", "ambiguities"):
        items = knowledge.get(key)
        if not isinstance(items, list):
            continue
        for item in items:
            if isinstance(item, dict) and isinstance(item.get("citations"), list):
                citations.extend(
                    citation for citation in item["citations"] if isinstance(citation, dict)
                )
    return citations


def validate_bilingual_field(value: dict[str, Any], key_en: str, key_source: str) -> bool:
    """Validate a required English/source-language pair."""
    return bool(sanitize_text(value.get(key_en)) and sanitize_text(value.get(key_source)))


def evidence_warning_messages(evidence_manifest: dict[str, Any]) -> list[str]:
    """Collect non-blocking evidence warnings for source-report visibility."""
    messages: list[str] = []
    manifest_warnings = evidence_manifest.get("warnings", [])
    if isinstance(manifest_warnings, list):
        messages.extend(
            str(message).strip()
            for message in manifest_warnings
            if isinstance(message, str) and str(message).strip()
        )
    for unit in evidence_manifest.get("units", []):
        if not isinstance(unit, dict):
            continue
        unit_id = unit.get("unit_id")
        unit_warnings = unit.get("warnings", [])
        if not isinstance(unit_id, str) or not isinstance(unit_warnings, list):
            continue
        messages.extend(
            f"{unit_id}: {message.strip()}"
            for message in unit_warnings
            if isinstance(message, str) and message.strip()
        )
    return list(dict.fromkeys(messages))


def validate_target(paths: WorkspacePaths, target: str) -> dict[str, Any]:
    """Validate a staged or published knowledge-base target and write the report."""
    target_root = paths.knowledge_target_dir(target)
    if not target_root.exists():
        return {
            "target": target,
            "status": "not-run",
            "blocking_errors": [f"{target_root} does not exist."],
            "warnings": [],
            "source_reports": [],
        }

    blocking_errors: list[str] = []
    warnings: list[str] = []
    source_reports: list[dict[str, Any]] = []
    required_root_files = [
        target_root / "catalog.json",
        target_root / "coverage_manifest.json",
        target_root / "graph_edges.json",
        target_root / "pending_work.json",
        target_root / "hybrid_work.json",
        target_root / "publish_manifest.json",
    ]
    for path in required_root_files:
        if not path.exists():
            blocking_errors.append(f"Missing required root artifact: {path.name}")

    catalog = read_json(target_root / "catalog.json")
    graph_edges: list[dict[str, Any]] = []
    source_contexts: list[dict[str, Any]] = []
    catalog_source_ids = {
        str(source["source_id"])
        for source in catalog.get("sources", [])
        if isinstance(source, dict) and isinstance(source.get("source_id"), str)
    }
    hybrid_work = read_json(paths.hybrid_work_path(target))
    if hybrid_work:
        blocking_errors.extend(
            validate_hybrid_work(
                hybrid_work,
                target=target,
                known_source_ids=catalog_source_ids,
            )
        )
    else:
        blocking_errors.append("Missing or invalid hybrid_work.json")
    pending_synthesis = False

    for source in catalog.get("sources", []):
        if not isinstance(source, dict):
            continue
        source_id = source.get("source_id")
        if not isinstance(source_id, str):
            continue
        source_dir = target_root / "sources" / source_id
        source_errors: list[str] = []
        source_warnings: list[str] = []
        for required in (
            "source_manifest.json",
            "evidence_manifest.json",
            "knowledge.json",
            "summary.md",
        ):
            if not (source_dir / required).exists():
                source_errors.append(f"Missing {required}")

        source_manifest = read_json(source_dir / "source_manifest.json")
        evidence_manifest = read_json(source_dir / "evidence_manifest.json")
        knowledge = read_json(source_dir / "knowledge.json")
        source_manifest = enrich_source_manifest_reference_fields(
            source_manifest,
            title=knowledge.get("title") if isinstance(knowledge, dict) else None,
        )
        evidence_manifest = sync_optional_sidecar_assets(
            source_dir,
            evidence_manifest=evidence_manifest,
        )
        evidence_manifest = enrich_evidence_manifest_reference_fields(
            source_manifest,
            evidence_manifest,
            source_dir=source_dir,
        )
        if target == "staging":
            write_json(source_dir / "source_manifest.json", source_manifest)
            write_json(source_dir / "evidence_manifest.json", evidence_manifest)
        affordance_path = source_dir / DEFAULT_AFFORDANCE_FILENAME
        affordances = read_json(affordance_path)
        summary_text = (
            (source_dir / "summary.md").read_text(encoding="utf-8")
            if (source_dir / "summary.md").exists()
            else ""
        )
        artifact_index_path_value = evidence_manifest.get("artifact_index_asset")
        artifact_index_path = (
            source_dir / artifact_index_path_value
            if isinstance(artifact_index_path_value, str) and artifact_index_path_value
            else source_dir / "artifact_index.json"
        )
        artifact_index = read_json(artifact_index_path)
        artifact_ids = {
            str(artifact["artifact_id"])
            for artifact in artifact_index.get("artifacts", [])
            if isinstance(artifact, dict) and isinstance(artifact.get("artifact_id"), str)
        }
        if artifact_index:
            source_errors.extend(validate_artifact_index(artifact_index, source_id=source_id))
        else:
            source_errors.append("Missing or invalid artifact_index.json")
        unit_ids = {
            str(unit["unit_id"])
            for unit in evidence_manifest.get("units", [])
            if isinstance(unit, dict) and isinstance(unit.get("unit_id"), str)
        }
        if str(source_manifest.get("document_type") or "") == "pdf":
            pdf_document_asset = evidence_manifest.get("pdf_document_asset")
            pdf_document_path = (
                source_dir / pdf_document_asset
                if isinstance(pdf_document_asset, str) and pdf_document_asset
                else source_dir / "pdf_document.json"
            )
            pdf_document = read_json(pdf_document_path)
            if pdf_document:
                source_errors.extend(
                    validate_pdf_document(
                        pdf_document,
                        source_id=source_id,
                        unit_ids=unit_ids,
                        artifact_ids=artifact_ids,
                    )
                )
            else:
                source_errors.append("Missing or invalid pdf_document.json")

        semantic_overlay_assets = evidence_manifest.get("semantic_overlay_assets", [])
        if semantic_overlay_assets is not None and not isinstance(semantic_overlay_assets, list):
            source_errors.append("semantic_overlay_assets must be a list when present")
        elif isinstance(semantic_overlay_assets, list):
            for asset in semantic_overlay_assets:
                if not isinstance(asset, str) or not asset:
                    source_errors.append(
                        "semantic_overlay_assets entries must be non-empty strings"
                    )
                    continue
                overlay_payload = read_json(source_dir / asset)
                if not overlay_payload:
                    source_errors.append(f"Missing or invalid semantic overlay `{asset}`")
                    continue
                source_errors.extend(
                    validate_semantic_overlay(
                        overlay_payload,
                        source_id=source_id,
                        unit_ids=unit_ids,
                        artifact_ids=artifact_ids,
                    )
                )

        if not knowledge:
            pending_synthesis = True
            source_errors.append("Missing or invalid knowledge.json")
        else:
            missing_keys = [key for key in REQUIRED_KNOWLEDGE_KEYS if key not in knowledge]
            if missing_keys:
                source_errors.append("Missing required knowledge keys: " + ", ".join(missing_keys))
            if knowledge.get("source_id") != source_id:
                source_errors.append("knowledge.json source_id does not match the staged source ID")
            if knowledge.get("source_fingerprint") != source_manifest.get("source_fingerprint"):
                source_errors.append(
                    "knowledge.json source_fingerprint does not match source_manifest.json"
                )
            if not validate_bilingual_field(knowledge, "summary_en", "summary_source"):
                source_errors.append(
                    "knowledge.json must include non-empty summary_en and summary_source"
                )
            if contains_placeholder(knowledge):
                source_errors.append("knowledge.json contains placeholder content")

            citations = citations_from_knowledge(knowledge)
            for citation in citations:
                if citation.get("unit_id") not in unit_ids:
                    source_errors.append(
                        f"Unresolved citation unit_id `{citation.get('unit_id')}` in knowledge.json"
                    )
                artifact_id = citation.get("artifact_id")
                if artifact_id is not None and artifact_id not in artifact_ids:
                    source_errors.append(
                        f"Unresolved citation artifact_id `{artifact_id}` in knowledge.json"
                    )
                if not sanitize_text(citation.get("support")):
                    source_errors.append("Each citation must include a non-empty support field")

            for related in knowledge.get("related_sources", []):
                if not isinstance(related, dict):
                    source_errors.append("related_sources entries must be objects")
                    continue
                related_source_id = related.get("source_id")
                if related_source_id not in catalog_source_ids:
                    source_errors.append(
                        f"related_sources entry references unknown source_id `{related_source_id}`"
                    )
                    continue
                citation_unit_ids = [
                    unit_id
                    for unit_id in related.get("citation_unit_ids", [])
                    if isinstance(unit_id, str)
                ]
                missing_citations = [
                    unit_id for unit_id in citation_unit_ids if unit_id not in unit_ids
                ]
                if missing_citations:
                    source_errors.append(
                        "related_sources contains unresolved citation_unit_ids: "
                        + ", ".join(missing_citations)
                    )
                citation_artifact_ids = [
                    artifact_id
                    for artifact_id in related.get("citation_artifact_ids", [])
                    if isinstance(artifact_id, str)
                ]
                missing_artifact_citations = [
                    artifact_id
                    for artifact_id in citation_artifact_ids
                    if artifact_id not in artifact_ids
                ]
                if missing_artifact_citations:
                    source_errors.append(
                        "related_sources contains unresolved citation_artifact_ids: "
                        + ", ".join(missing_artifact_citations)
                    )
                graph_edges.append(
                    {
                        "source_id": source_id,
                        "related_source_id": related_source_id,
                        "relation_type": related.get("relation_type"),
                        "strength": related.get("strength"),
                        "status": related.get("status"),
                        "citation_unit_ids": citation_unit_ids,
                        "source_unit_id": related.get("source_unit_id"),
                        "related_unit_id": related.get("related_unit_id"),
                        "source_artifact_id": related.get("source_artifact_id"),
                        "related_artifact_id": related.get("related_artifact_id"),
                        "citation_artifact_ids": citation_artifact_ids,
                    }
                )
            for related in evidence_manifest.get("deterministic_linked_sources", []):
                if not isinstance(related, dict):
                    source_errors.append("deterministic_linked_sources entries must be objects")
                    continue
                related_source_id = related.get("related_source_id")
                if related_source_id not in catalog_source_ids:
                    source_errors.append(
                        "deterministic_linked_sources entry references unknown source_id "
                        f"`{related_source_id}`"
                    )
                    continue
                citation_unit_ids = [
                    unit_id
                    for unit_id in related.get("citation_unit_ids", [])
                    if isinstance(unit_id, str)
                ]
                missing_citations = [
                    unit_id for unit_id in citation_unit_ids if unit_id not in unit_ids
                ]
                if missing_citations:
                    source_errors.append(
                        "deterministic_linked_sources contains unresolved citation_unit_ids: "
                        + ", ".join(missing_citations)
                    )
                citation_artifact_ids = [
                    artifact_id
                    for artifact_id in related.get("citation_artifact_ids", [])
                    if isinstance(artifact_id, str)
                ]
                missing_artifact_citations = [
                    artifact_id
                    for artifact_id in citation_artifact_ids
                    if artifact_id not in artifact_ids
                ]
                if missing_artifact_citations:
                    source_errors.append(
                        "deterministic_linked_sources contains unresolved citation_artifact_ids: "
                        + ", ".join(missing_artifact_citations)
                    )
                graph_edges.append(
                    {
                        "source_id": source_id,
                        "related_source_id": related_source_id,
                        "relation_type": related.get("relation_type"),
                        "strength": related.get("strength"),
                        "status": related.get("status"),
                        "citation_unit_ids": citation_unit_ids,
                        "source_unit_id": related.get("source_unit_id"),
                        "related_unit_id": related.get("related_unit_id"),
                        "source_artifact_id": related.get("source_artifact_id"),
                        "related_artifact_id": related.get("related_artifact_id"),
                        "citation_artifact_ids": citation_artifact_ids,
                    }
                )

        if not summary_text.strip():
            pending_synthesis = True
            source_errors.append("summary.md is missing or empty")
        else:
            if (
                "## English Summary" not in summary_text
                or "## Source-Language Summary" not in summary_text
            ):
                source_errors.append(
                    "summary.md must include `## English Summary` and `## Source-Language Summary`"
                )
            if source_id not in summary_text:
                source_errors.append("summary.md must mention the source ID")
            if contains_placeholder(summary_text):
                source_errors.append("summary.md contains placeholder content")

        units = [
            unit
            for unit in evidence_manifest.get("units", [])
            if isinstance(unit, dict) and isinstance(unit.get("unit_id"), str)
        ]
        source_warnings.extend(evidence_warning_messages(evidence_manifest))
        renders = [
            render
            for render in evidence_manifest.get("document_renders", [])
            if isinstance(render, str)
        ]
        for render in renders:
            if not (source_dir / render).exists():
                source_errors.append(f"Missing rendered asset `{render}`")
        if not artifact_index_path.exists():
            source_errors.append("Missing artifact_index.json")
        for artifact in artifact_index.get("artifacts", []):
            if not isinstance(artifact, dict):
                continue
            artifact_path_value = artifact.get("artifact_path")
            if isinstance(artifact_path_value, str) and artifact_path_value:
                if not (source_dir / artifact_path_value).exists():
                    source_errors.append(f"Missing artifact asset `{artifact_path_value}`")
            for render_asset in artifact.get("render_assets", []):
                if (
                    isinstance(render_asset, str)
                    and render_asset
                    and not (source_dir / render_asset).exists()
                ):
                    source_errors.append(f"Missing artifact render asset `{render_asset}`")
            for focus_asset in artifact.get("focus_render_assets", []):
                if (
                    isinstance(focus_asset, str)
                    and focus_asset
                    and not (source_dir / focus_asset).exists()
                ):
                    source_errors.append(f"Missing artifact focus render asset `{focus_asset}`")

        if source_manifest.get("document_type") in {"pdf", "pptx"}:
            for unit in units:
                if source_manifest.get("document_type") == "pptx" and bool(unit.get("hidden")):
                    continue
                asset = unit.get("rendered_asset")
                if not isinstance(asset, str) or not asset:
                    source_errors.append(f"{unit['unit_id']} is missing a rendered asset")
                    continue
                if not (source_dir / asset).exists():
                    source_errors.append(f"Missing rendered asset `{asset}`")
        elif source_manifest.get("document_type") in {"docx", "xlsx"} and not renders:
            source_errors.append(
                f"{source_manifest.get('document_type')} evidence must include "
                "at least one document-level render"
            )
        elif source_manifest.get("document_type") in NON_RENDERED_DOCUMENT_TYPES and not units:
            source_errors.append(
                f"{source_manifest.get('document_type')} evidence must preserve at least one unit"
            )

        for unit in units:
            text_asset = unit.get("text_asset")
            structure_asset = unit.get("structure_asset")
            if (
                isinstance(text_asset, str)
                and text_asset
                and not (source_dir / text_asset).exists()
            ):
                source_errors.append(f"Missing text asset `{text_asset}`")
            if (
                isinstance(structure_asset, str)
                and structure_asset
                and not (source_dir / structure_asset).exists()
            ):
                source_errors.append(f"Missing structure asset `{structure_asset}`")
            for render_asset in unit.get("render_assets", []):
                if (
                    isinstance(render_asset, str)
                    and render_asset
                    and not (source_dir / render_asset).exists()
                ):
                    source_errors.append(f"Missing unit render asset `{render_asset}`")
            render_page_span = unit.get("render_page_span")
            if render_page_span is not None and not (
                isinstance(render_page_span, dict)
                and isinstance(render_page_span.get("start"), int)
                and isinstance(render_page_span.get("end"), int)
            ):
                source_errors.append(f"{unit['unit_id']} has an invalid render_page_span")

        document_type = source_manifest.get("document_type")
        if document_type == "xlsx":
            workbook_asset = evidence_manifest.get("spreadsheet_workbook_asset")
            if not isinstance(workbook_asset, str) or not workbook_asset:
                source_errors.append("xlsx evidence must include spreadsheet_workbook_asset")
            elif not (source_dir / workbook_asset).exists():
                source_errors.append(f"Missing spreadsheet workbook asset `{workbook_asset}`")
            sheet_assets = evidence_manifest.get("spreadsheet_sheet_assets", [])
            if not isinstance(sheet_assets, list) or not sheet_assets:
                source_errors.append("xlsx evidence must include spreadsheet_sheet_assets")
            else:
                for asset in sheet_assets:
                    if isinstance(asset, str) and asset and not (source_dir / asset).exists():
                        source_errors.append(f"Missing spreadsheet sheet asset `{asset}`")
        if document_type == "pdf":
            pdf_document_asset = evidence_manifest.get("pdf_document_asset")
            if not isinstance(pdf_document_asset, str) or not pdf_document_asset:
                source_errors.append("pdf evidence must include pdf_document_asset")
            elif not (source_dir / pdf_document_asset).exists():
                source_errors.append(f"Missing pdf document asset `{pdf_document_asset}`")
        if document_type in {"pdf", "pptx", "docx", "xlsx"}:
            visual_assets = evidence_manifest.get("visual_layout_assets", [])
            if not isinstance(visual_assets, list) or not visual_assets:
                source_errors.append(f"{document_type} evidence must include visual_layout_assets")
            else:
                for asset in visual_assets:
                    if isinstance(asset, str) and asset and not (source_dir / asset).exists():
                        source_errors.append(f"Missing visual layout asset `{asset}`")

        failures = evidence_manifest.get("failures", [])
        if isinstance(failures, list) and failures:
            if units or renders:
                source_warnings.append(
                    "Extraction recorded partial failures but preserved surviving evidence"
                )
            else:
                source_errors.append("Extraction failed without preserving surviving evidence")

        derived_affordances = derive_source_affordances(
            source_manifest=source_manifest,
            evidence_manifest=evidence_manifest,
            source_dir=source_dir,
            knowledge=knowledge or None,
            summary_text=summary_text,
        )
        if affordance_path.exists():
            affordance_errors = validate_derived_affordances(
                affordances,
                source_manifest=source_manifest,
                evidence_manifest=evidence_manifest,
            )
            if affordance_errors:
                source_errors.extend(affordance_errors)
            else:
                derived_affordances = merge_derived_affordances(derived_affordances, affordances)
        else:
            source_warnings.append(
                "derived_affordances.json is missing; using deterministic fallback affordances"
            )

        if knowledge and summary_text.strip():
            source_contexts.append(
                {
                    "source_manifest": source_manifest,
                    "evidence_manifest": evidence_manifest,
                    "knowledge": knowledge,
                    "affordances": derived_affordances,
                    "summary_text": summary_text,
                    "artifact_dir": source_dir,
                    "source_family": "corpus",
                    "trust_tier": "source",
                }
            )

        source_reports.append(
            {
                "source_id": source_id,
                "current_path": source_manifest.get("current_path"),
                "errors": list(dict.fromkeys(source_errors)),
                "warnings": list(dict.fromkeys(source_warnings)),
            }
        )
        blocking_errors.extend(f"{source_id}: {message}" for message in source_errors)

    interaction_contexts = load_promoted_interaction_contexts(paths, target=target)
    interaction_source_ids = {
        str(context["source_manifest"]["source_id"])
        for context in interaction_contexts
        if isinstance(context.get("source_manifest"), dict)
        and isinstance(context["source_manifest"].get("source_id"), str)
    }
    all_known_source_ids = catalog_source_ids | interaction_source_ids
    for context in interaction_contexts:
        source_manifest = context["source_manifest"]
        evidence_manifest = context["evidence_manifest"]
        knowledge = context["knowledge"]
        source_dir = Path(context["artifact_dir"])
        source_manifest = enrich_source_manifest_reference_fields(
            source_manifest,
            title=knowledge.get("title") if isinstance(knowledge, dict) else None,
        )
        evidence_manifest = enrich_evidence_manifest_reference_fields(
            source_manifest,
            evidence_manifest,
            source_dir=source_dir,
        )
        if target == "staging":
            write_json(source_dir / "source_manifest.json", source_manifest)
            write_json(source_dir / "evidence_manifest.json", evidence_manifest)
        affordance_path = source_dir / DEFAULT_AFFORDANCE_FILENAME
        affordances = read_json(affordance_path)
        summary_text = context["summary_text"]
        source_id = str(source_manifest.get("source_id"))
        interaction_errors: list[str] = []
        interaction_warnings: list[str] = []

        for required in (
            "source_manifest.json",
            "evidence_manifest.json",
            "knowledge.json",
            "summary.md",
        ):
            if not (source_dir / required).exists():
                interaction_errors.append(f"Missing {required}")

        missing_keys = [key for key in REQUIRED_KNOWLEDGE_KEYS if key not in knowledge]
        if missing_keys:
            interaction_errors.append("Missing required knowledge keys: " + ", ".join(missing_keys))
        if knowledge.get("source_id") != source_id:
            interaction_errors.append(
                "knowledge.json source_id does not match the staged source ID"
            )
        if knowledge.get("source_fingerprint") != source_manifest.get("source_fingerprint"):
            interaction_errors.append(
                "knowledge.json source_fingerprint does not match source_manifest.json"
            )
        if not validate_bilingual_field(knowledge, "summary_en", "summary_source"):
            interaction_errors.append(
                "knowledge.json must include non-empty summary_en and summary_source"
            )
        if contains_placeholder(knowledge):
            interaction_errors.append("knowledge.json contains placeholder content")

        citations = citations_from_knowledge(knowledge)
        unit_ids = {
            str(unit["unit_id"])
            for unit in evidence_manifest.get("units", [])
            if isinstance(unit, dict) and isinstance(unit.get("unit_id"), str)
        }
        for citation in citations:
            if citation.get("unit_id") not in unit_ids:
                interaction_errors.append(
                    f"Unresolved citation unit_id `{citation.get('unit_id')}` in knowledge.json"
                )
            if not sanitize_text(citation.get("support")):
                interaction_errors.append("Each citation must include a non-empty support field")

        for related in knowledge.get("related_sources", []):
            if not isinstance(related, dict):
                interaction_errors.append("related_sources entries must be objects")
                continue
            related_source_id = related.get("source_id")
            if related_source_id not in all_known_source_ids:
                interaction_errors.append(
                    f"related_sources entry references unknown source_id `{related_source_id}`"
                )
                continue
            citation_unit_ids = [
                unit_id
                for unit_id in related.get("citation_unit_ids", [])
                if isinstance(unit_id, str)
            ]
            missing_citations = [
                unit_id for unit_id in citation_unit_ids if unit_id not in unit_ids
            ]
            if missing_citations:
                interaction_errors.append(
                    "related_sources contains unresolved citation_unit_ids: "
                    + ", ".join(missing_citations)
                )
            graph_edges.append(
                {
                    "source_id": source_id,
                    "related_source_id": related_source_id,
                    "relation_type": related.get("relation_type"),
                    "strength": related.get("strength"),
                    "status": related.get("status"),
                    "citation_unit_ids": citation_unit_ids,
                }
            )

        if not summary_text.strip():
            interaction_errors.append("summary.md is missing or empty")
        else:
            if (
                "## English Summary" not in summary_text
                or "## Source-Language Summary" not in summary_text
            ):
                interaction_errors.append(
                    "summary.md must include `## English Summary` and `## Source-Language Summary`"
                )
            if source_id not in summary_text:
                interaction_errors.append("summary.md must mention the source ID")
            if contains_placeholder(summary_text):
                interaction_errors.append("summary.md contains placeholder content")

        renders = [
            render
            for render in evidence_manifest.get("document_renders", [])
            if isinstance(render, str)
        ]
        for render in renders:
            if not (source_dir / render).exists():
                interaction_errors.append(f"Missing rendered asset `{render}`")

        for unit in evidence_manifest.get("units", []):
            if not isinstance(unit, dict) or not isinstance(unit.get("unit_id"), str):
                continue
            text_asset = unit.get("text_asset")
            structure_asset = unit.get("structure_asset")
            if (
                isinstance(text_asset, str)
                and text_asset
                and not (source_dir / text_asset).exists()
            ):
                interaction_errors.append(f"Missing text asset `{text_asset}`")
            if (
                isinstance(structure_asset, str)
                and structure_asset
                and not (source_dir / structure_asset).exists()
            ):
                interaction_errors.append(f"Missing structure asset `{structure_asset}`")

        derived_affordances = derive_source_affordances(
            source_manifest=source_manifest,
            evidence_manifest=evidence_manifest,
            source_dir=source_dir,
            knowledge=knowledge,
            summary_text=summary_text,
        )
        if affordance_path.exists():
            affordance_errors = validate_derived_affordances(
                affordances,
                source_manifest=source_manifest,
                evidence_manifest=evidence_manifest,
            )
            if affordance_errors:
                interaction_errors.extend(affordance_errors)
            else:
                derived_affordances = merge_derived_affordances(derived_affordances, affordances)
        else:
            interaction_warnings.append(
                "derived_affordances.json is missing; using deterministic fallback affordances"
            )

        source_reports.append(
            {
                "source_id": source_id,
                "current_path": source_manifest.get("current_path"),
                "errors": list(dict.fromkeys(interaction_errors)),
                "warnings": list(dict.fromkeys(interaction_warnings)),
            }
        )
        blocking_errors.extend(f"{source_id}: {message}" for message in interaction_errors)
        if not interaction_errors:
            context["source_manifest"] = source_manifest
            context["evidence_manifest"] = evidence_manifest
            context["affordances"] = derived_affordances
            source_contexts.append(context)

    write_json(target_root / "graph_edges.json", {"generated_at": utc_now(), "edges": graph_edges})
    retrieval_manifest: dict[str, Any] | None = None
    trace_manifest: dict[str, Any] | None = None
    if not pending_synthesis and not blocking_errors:
        retrieval_manifest = build_retrieval_artifacts(
            paths,
            target=target,
            source_contexts=source_contexts,
            graph_edges=graph_edges,
            source_signature=catalog.get("source_signature"),
        )
        trace_manifest = build_trace_artifacts(
            paths,
            target=target,
            source_contexts=source_contexts,
            graph_edges=graph_edges,
            source_signature=catalog.get("source_signature"),
        )

    status = "valid"
    if pending_synthesis:
        status = "pending-synthesis"
    elif blocking_errors:
        status = "blocking-errors"
    elif warnings:
        status = "warnings"

    report = {
        "generated_at": utc_now(),
        "target": target,
        "status": status,
        "source_signature": catalog.get("source_signature"),
        "blocking_errors": blocking_errors,
        "warnings": warnings,
        "source_reports": source_reports,
        "edge_count": len(graph_edges),
        "interaction_memory_count": len(interaction_contexts),
        "retrieval_artifacts_built": retrieval_manifest is not None,
        "trace_artifacts_built": trace_manifest is not None,
    }
    write_json(target_root / "validation_report.json", report)
    update_dependency_state(
        paths,
        target=target,
        pending_sources=read_json(target_root / "pending_work.json").get("pending_sources", []),
        retrieval_manifest=retrieval_manifest,
        trace_manifest=trace_manifest,
    )
    return report


def publish_staging(paths: WorkspacePaths, validation_report: dict[str, Any]) -> dict[str, Any]:
    """Publish the staged knowledge base as an immutable snapshot and activate `current`."""
    return publish_staging_snapshot(
        paths,
        validation_report=validation_report,
        published_at=utc_now(),
    )


def update_sync_state(
    paths: WorkspacePaths,
    *,
    staging_signature: str,
    published_signature: str | None,
    validation_status: str,
    validation_target: str,
    pending_sources: list[dict[str, Any]],
    change_set: dict[str, Any] | None = None,
    reused_sources: int = 0,
    rebuilt_sources: int = 0,
) -> dict[str, Any]:
    """Persist the current Phase 4 sync state."""
    state = sync_state(paths)
    payload = {
        **state,
        "last_sync_at": utc_now(),
        "last_validation_at": utc_now(),
        "last_validation_status": validation_status,
        "last_validation_target": validation_target,
        "staging_present": paths.knowledge_base_staging_dir.exists(),
        "staging_source_signature": staging_signature,
        "published_source_signature": published_signature
        or state.get("published_source_signature"),
        "last_publish_at": state.get("last_publish_at"),
        "pending_sources": len(pending_sources),
        "last_change_set": change_set or state.get("last_change_set"),
        "reused_sources": reused_sources,
        "rebuilt_sources": rebuilt_sources,
        "retrieval_artifact_signature": artifact_signature(
            paths.retrieval_manifest_path(validation_target)
        ),
        "trace_artifact_signature": artifact_signature(
            paths.trace_manifest_path(validation_target)
        ),
    }
    write_json(paths.sync_state_path, payload)
    return payload


def sync_workspace(paths: WorkspacePaths, *, autonomous: bool = True) -> dict[str, Any]:
    """Run the staged sync workflow and, by default, close the loop to final publish."""
    autonomous_steps: list[dict[str, Any]] = []
    phase_costs = {
        "detect": 0.0,
        "stage": 0.0,
        "repair": 0.0,
        "author": 0.0,
        "validate": 0.0,
        "publish": 0.0,
        "projection_enqueue": 0.0,
    }
    repair_actions: list[dict[str, Any]] = []

    pdf_snapshot = pdf_renderer_snapshot()
    if (
        any(path.suffix.lower() == ".pdf" for path in supported_source_documents(paths))
        and not pdf_snapshot["ready"]
    ):
        return {
            "status": "action-required",
            "detail": pdf_snapshot["detail"],
            "pending_sources": [],
            "validation": None,
            "published": False,
            "auto_repairs": {"repair_count": 0},
            "auto_authoring": {"attempted": 0, "authored": [], "authored_count": 0},
            "autonomous_steps": [
                {
                    "step": "detect",
                    "status": "blocked",
                    "detail": pdf_snapshot["detail"],
                }
            ],
            "required_capabilities": ["pdf-rendering"],
            "phase_costs": phase_costs,
            "publish_skipped": False,
            "publish_skip_reason": None,
            "repair_actions": repair_actions,
            "projection_state": projection_state_summary(paths),
        }

    office_snapshot = office_renderer_snapshot(paths)
    if office_snapshot["required"] and not office_snapshot["ready"]:
        return {
            "status": "action-required",
            "detail": office_snapshot["detail"],
            "pending_sources": [],
            "validation": None,
            "published": False,
            "auto_repairs": {"repair_count": 0},
            "auto_authoring": {"attempted": 0, "authored": [], "authored_count": 0},
            "autonomous_steps": [
                {
                    "step": "detect",
                    "status": "blocked",
                    "detail": office_snapshot["detail"],
                }
            ],
            "required_capabilities": ["office-rendering"],
            "phase_costs": phase_costs,
            "publish_skipped": False,
            "publish_skip_reason": None,
            "repair_actions": repair_actions,
            "projection_state": projection_state_summary(paths),
        }
    with workspace_lease(paths, "sync", timeout_seconds=600.0):
        if paths.knowledge_base_dir.joinpath(".staging-build").exists():
            shutil.rmtree(paths.knowledge_base_dir / ".staging-build")
            repair_actions.append(
                {
                    "kind": "removed-stale-staging-build",
                    "path": str(
                        (paths.knowledge_base_dir / ".staging-build").relative_to(paths.root)
                    ),
                }
            )
        repair_actions.extend(repair_stale_shared_jobs(paths))
        detect_started = perf_counter()
        current_signature = source_inventory_signature(paths)
        state = sync_state(paths)
        index_payload, active_sources, ambiguous_match, change_set = update_source_index(paths)
        active_source_ids = {
            str(source["source_id"])
            for source in active_sources
            if isinstance(source.get("source_id"), str)
        }
        change_stats = change_set.get("stats", {})
        autonomous_steps.append(
            {
                "step": "detect",
                "status": "completed",
                "detail": (
                    "Detected source changes: "
                    f"unchanged={change_stats.get('unchanged', 0)}, "
                    f"added={change_stats.get('added', 0)}, "
                    f"modified={change_stats.get('modified', 0)}, "
                    f"moved_or_renamed={change_stats.get('moved_or_renamed', 0)}, "
                    f"deleted={change_stats.get('deleted', 0)}, "
                    f"ambiguous={change_stats.get('ambiguous', 0)}"
                ),
            }
        )
        phase_costs["detect"] = perf_counter() - detect_started
        rebuild_required = (
            not paths.knowledge_base_staging_dir.exists()
            or state.get("staging_source_signature") != current_signature
            or not staging_source_artifacts_complete(paths, active_sources)
            or ambiguous_match
            or any(
                int(change_stats.get(key, 0)) > 0
                for key in ("added", "modified", "moved_or_renamed", "deleted", "ambiguous")
            )
        )
        build_stats = {"reused_sources": 0, "rebuilt_sources": 0}

        stage_started = perf_counter()
        if rebuild_required:
            _catalog_sources, _source_summaries, _ambiguous, build_stats = build_staging_artifacts(
                paths,
                active_sources,
                office_snapshot["binary"],
            )
        else:
            refresh_staging_source_metadata(paths, active_sources)
            _catalog_sources, _source_summaries, _ambiguous = write_staging_root_artifacts(
                paths,
                active_sources,
            )
        change_set = refresh_change_set_details(change_set, active_sources)
        autonomous_steps.append(
            {
                "step": "stage",
                "status": "completed",
                "detail": (
                    f"Staging {'rebuilt' if rebuild_required else 'refreshed'} with "
                    f"reused_sources={build_stats['reused_sources']} and "
                    f"rebuilt_sources={build_stats['rebuilt_sources']}."
                ),
            }
        )
        phase_costs["stage"] = perf_counter() - stage_started

        repair_started = perf_counter()
        interaction_manifest = build_promoted_interaction_memories(
            paths,
            target="staging",
            active_source_ids=active_source_ids,
        )
        interaction_snapshot = interaction_ingest_snapshot(paths)
        auto_repairs = repair_staging_semantic_artifacts(
            paths,
            active_source_ids=active_source_ids,
            active_sources=active_sources,
        )
        change_set = refresh_change_set_details(change_set, active_sources)
        autonomous_steps.append(
            {
                "step": "repair",
                "status": "completed",
                "detail": (
                    "Applied silent staging repairs: "
                    f"repair_count={auto_repairs.get('repair_count', 0)}."
                ),
            }
        )
        phase_costs["repair"] = perf_counter() - repair_started

        pending_sources = collect_pending_synthesis(paths)
        auto_authoring = {
            "attempted": 0,
            "authored": [],
            "authored_count": 0,
            "mode": "disabled" if not autonomous else "conservative-in-repo",
        }
        author_started = perf_counter()
        if pending_sources and autonomous:
            auto_authoring = auto_author_pending_semantics(
                paths,
                pending_sources=pending_sources,
                active_source_ids=active_source_ids,
            )
            pending_sources = collect_pending_synthesis(paths)
            autonomous_steps.append(
                {
                    "step": "author",
                    "status": "completed" if not pending_sources else "blocked",
                    "detail": (
                        "Autonomous authoring refreshed staged semantic outputs for "
                        f"{auto_authoring.get('authored_count', 0)} item(s)."
                    ),
                }
            )
        elif pending_sources:
            autonomous_steps.append(
                {
                    "step": "author",
                    "status": "skipped",
                    "detail": "Autonomous authoring is disabled for this sync call.",
                }
            )
        else:
            autonomous_steps.append(
                {
                    "step": "author",
                    "status": "completed",
                    "detail": "No staged semantic authoring was required.",
                }
            )
        phase_costs["author"] = perf_counter() - author_started

        hybrid_enrichment = hybrid_enrichment_status(paths, target="staging")
        autonomous_steps.append(
            {
                "step": "hybrid-enrichment",
                "status": (
                    "completed"
                    if hybrid_enrichment.get("mode") in {"not-needed", "covered"}
                    else "degraded"
                ),
                "detail": hybrid_enrichment.get("detail"),
            }
        )

        if pending_sources and not autonomous:
            update_sync_state(
                paths,
                staging_signature=current_signature,
                published_signature=state.get("published_source_signature"),
                validation_status="pending-synthesis",
                validation_target="staging",
                pending_sources=pending_sources,
                change_set=change_set,
                reused_sources=build_stats["reused_sources"],
                rebuilt_sources=build_stats["rebuilt_sources"],
            )
            update_dependency_state(
                paths,
                target="staging",
                pending_sources=pending_sources,
                retrieval_manifest=None,
                trace_manifest=None,
            )
            return {
                "status": "pending-synthesis",
                "detail": (
                    "Staged evidence is ready, but autonomous authoring is "
                    "disabled for this run."
                ),
                "pending_sources": pending_sources,
                "validation": None,
                "published": False,
                "rebuilt": rebuild_required,
                "build_stats": build_stats,
                "auto_repairs": auto_repairs,
                "auto_authoring": auto_authoring,
                "hybrid_enrichment": hybrid_enrichment,
                "autonomous_steps": autonomous_steps,
                "required_capabilities": ["autonomous-authoring"],
                "interaction_ingest": {
                    **interaction_snapshot,
                    "promoted_memory_count": interaction_manifest.get("memory_count", 0),
                    "promotion_ready": True,
                },
                "change_set": change_set,
                "source_index": index_payload,
                "phase_costs": phase_costs,
                "publish_skipped": False,
                "publish_skip_reason": None,
                "repair_actions": repair_actions,
                "projection_state": projection_state_summary(paths),
            }
        if pending_sources:
            update_sync_state(
                paths,
                staging_signature=current_signature,
                published_signature=state.get("published_source_signature"),
                validation_status="blocking-errors",
                validation_target="staging",
                pending_sources=pending_sources,
                change_set=change_set,
                reused_sources=build_stats["reused_sources"],
                rebuilt_sources=build_stats["rebuilt_sources"],
            )
            update_dependency_state(
                paths,
                target="staging",
                pending_sources=pending_sources,
                retrieval_manifest=None,
                trace_manifest=None,
            )
            return {
                "status": "action-required",
                "detail": (
                    "Autonomous authoring could not fully resolve the staged "
                    "semantic outputs."
                ),
                "pending_sources": pending_sources,
                "validation": None,
                "published": False,
                "rebuilt": rebuild_required,
                "build_stats": build_stats,
                "auto_repairs": auto_repairs,
                "auto_authoring": auto_authoring,
                "hybrid_enrichment": hybrid_enrichment,
                "autonomous_steps": autonomous_steps,
                "required_capabilities": [],
                "interaction_ingest": {
                    **interaction_snapshot,
                    "promoted_memory_count": interaction_manifest.get("memory_count", 0),
                    "promotion_ready": True,
                },
                "change_set": change_set,
                "source_index": index_payload,
                "phase_costs": phase_costs,
                "publish_skipped": False,
                "publish_skip_reason": None,
                "repair_actions": repair_actions,
                "projection_state": projection_state_summary(paths),
            }

        validate_started = perf_counter()
        validation_report = validate_target(paths, target="staging")
        phase_costs["validate"] = perf_counter() - validate_started
        autonomous_steps.append(
            {
                "step": "validate",
                "status": (
                    "completed"
                    if validation_report["status"] in {"valid", "warnings"}
                    else "blocked"
                ),
                "detail": (
                    f"Validation finished with status={validation_report['status']}, "
                    f"blocking_errors={len(validation_report['blocking_errors'])}, "
                    f"warnings={len(validation_report['warnings'])}."
                ),
            }
        )
        published = False
        publish_skipped = False
        publish_skip_reason: str | None = None
        published_manifest: dict[str, Any] | None = None
        published_signature: str | None = state.get("published_source_signature")
        if validation_report["status"] in {"valid", "warnings"}:
            if (
                not rebuild_required
                and int(auto_repairs.get("repair_count", 0) or 0) == 0
                and int(interaction_snapshot.get("pending_promotion_count", 0) or 0) == 0
                and state.get("published_source_signature") == current_signature
            ):
                publish_skipped = True
                publish_skip_reason = (
                    "Published truth already matches the current source signature and no "
                    "interaction promotion or repair changed publishable state."
                )
                autonomous_steps.append(
                    {
                        "step": "publish",
                        "status": "completed",
                        "detail": (
                            "Skipped publication because the current snapshot "
                            "is already current."
                        ),
                    }
                )
            else:
                publish_started = perf_counter()
                published_manifest = publish_staging(paths, validation_report)
                phase_costs["publish"] = perf_counter() - publish_started
                published = True
                published_signature = current_signature
                promoted_result = mark_promoted_interaction_entries(paths, target="current")
                interaction_snapshot = {
                    **interaction_snapshot,
                    "last_overlay_at": promoted_result["pending_overlay"].get("generated_at"),
                    "pending_capture_count": promoted_result["pending_overlay"].get(
                        "pending_entry_count", 0
                    ),
                    "pending_promotion_count": promoted_result["pending_overlay"].get(
                        "pending_entry_count",
                        0,
                    ),
                    "sync_recommended": bool(
                        promoted_result["pending_overlay"].get("pending_entry_count", 0)
                    ),
                }
                autonomous_steps.append(
                    {
                        "step": "publish",
                        "status": "completed",
                        "detail": (
                            "Published an immutable snapshot and activated "
                            "`knowledge_base/current`."
                        ),
                    }
                )
                projection_started = perf_counter()
                queue_projection_refresh(
                    paths,
                    reason="A knowledge-base publish activated new canonical runtime truth.",
                )
                phase_costs["projection_enqueue"] = perf_counter() - projection_started
        else:
            autonomous_steps.append(
                {
                    "step": "publish",
                    "status": "blocked",
                    "detail": "Validation blocked final publication.",
                }
            )

        sync_payload = update_sync_state(
            paths,
            staging_signature=current_signature,
            published_signature=published_signature,
            validation_status=validation_report["status"],
            validation_target="staging",
            pending_sources=[],
            change_set=change_set,
            reused_sources=build_stats["reused_sources"],
            rebuilt_sources=build_stats["rebuilt_sources"],
        )
        if published:
            sync_payload["last_publish_at"] = (
                published_manifest.get("published_at") if published_manifest else None
            )
            update_dependency_state(
                paths,
                target="current",
                pending_sources=[],
                retrieval_manifest=read_json(paths.retrieval_manifest_path("current")),
                trace_manifest=read_json(paths.trace_manifest_path("current")),
            )
            write_json(paths.sync_state_path, sync_payload)

        return {
            "status": validation_report["status"],
            "detail": (
                "Published the staged knowledge base."
                if published
                else (
                    "Published truth was already current, so final publication was skipped."
                    if publish_skipped
                    else "Validation blocked publication."
                )
            ),
            "pending_sources": [],
            "validation": validation_report,
            "published": published,
            "rebuilt": rebuild_required,
            "build_stats": build_stats,
            "auto_repairs": auto_repairs,
            "auto_authoring": auto_authoring,
            "hybrid_enrichment": hybrid_enrichment,
            "autonomous_steps": autonomous_steps,
            "required_capabilities": [],
            "interaction_ingest": {
                **interaction_snapshot,
                "promoted_memory_count": interaction_manifest.get("memory_count", 0),
                "promotion_ready": True,
            },
            "change_set": change_set,
            "source_index": index_payload,
            "phase_costs": phase_costs,
            "publish_skipped": publish_skipped,
            "publish_skip_reason": publish_skip_reason,
            "repair_actions": repair_actions,
            "projection_state": projection_state_summary(paths),
        }


def validate_workspace(paths: WorkspacePaths, *, target: str) -> dict[str, Any]:
    """Validate a selected knowledge-base target and update runtime state."""
    report = validate_target(paths, target=target)
    published_signature = sync_state(paths).get("published_source_signature")
    if target == "current":
        published_signature = source_inventory_signature(paths)
    update_sync_state(
        paths,
        staging_signature=sync_state(paths).get(
            "staging_source_signature", source_inventory_signature(paths)
        ),
        published_signature=published_signature,
        validation_status=report["status"],
        validation_target=target,
        pending_sources=read_json(paths.staging_pending_work_path).get("pending_sources", []),
        change_set=sync_state(paths).get("last_change_set"),
        reused_sources=int(sync_state(paths).get("reused_sources", 0) or 0),
        rebuilt_sources=int(sync_state(paths).get("rebuilt_sources", 0) or 0),
    )
    return report
